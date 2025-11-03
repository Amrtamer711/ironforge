from pathlib import Path
from typing import List, Tuple
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import config


def format_date_for_display(date_str: str) -> str:
    """Convert various date formats to 'Xth Month YYYY' format (e.g., '5th November 2025')."""
    if not date_str:
        return date_str

    # Try to parse the date in various formats
    date_obj = None
    formats_to_try = [
        '%Y-%m-%d',  # 2025-11-05
        '%d/%m/%Y',  # 05/11/2025
        '%m/%d/%Y',  # 11/05/2025
        '%d-%m-%Y',  # 05-11-2025
        '%Y/%m/%d',  # 2025/11/05
        '%d %B %Y',  # 5 November 2025
        '%d %b %Y',  # 5 Nov 2025
        '%B %d, %Y',  # November 5, 2025
        '%b %d, %Y',  # Nov 5, 2025
    ]

    for fmt in formats_to_try:
        try:
            date_obj = datetime.strptime(date_str.strip(), fmt)
            break
        except ValueError:
            continue

    # If we couldn't parse it, check if it already looks like the target format
    if date_obj is None:
        # Check if it already has the ordinal suffix (1st, 2nd, 3rd, etc.)
        if any(suffix in date_str.lower() for suffix in ['st ', 'nd ', 'rd ', 'th ']):
            return date_str  # Already in good format
        return date_str  # Return as-is if we can't parse it

    # Format with ordinal suffix
    day = date_obj.day
    if day in [1, 21, 31]:
        suffix = "st"
    elif day in [2, 22]:
        suffix = "nd"
    elif day in [3, 23]:
        suffix = "rd"
    else:
        suffix = "th"

    return f"{day}{suffix} {date_obj.strftime('%B %Y')}"


def add_location_text_with_colored_sov(paragraph, location_text: str, scale: float) -> None:
    """Render location text in a single consistent style (all black)."""
    run = paragraph.add_run()
    run.text = location_text
    run.font.size = Pt(int(20 * scale))
    run.font.color.rgb = RGBColor(0, 0, 0)


def set_cell_border(cell, edges=("L", "R", "T", "B")) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(side))
        if existing is not None:
            tcPr.remove(existing)

    for edge in edges:
        ln = OxmlElement(f"a:ln{edge}")
        ln.set("w", "25400")
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", "000000")
        solidFill.append(srgbClr)
        ln.append(solidFill)

        prstDash = OxmlElement("a:prstDash")
        prstDash.set("val", "solid")
        ln.append(prstDash)

        headEnd = OxmlElement("a:headEnd")
        headEnd.set("type", "none")
        ln.append(headEnd)

        tailEnd = OxmlElement("a:tailEnd")
        tailEnd.set("type", "none")
        ln.append(tailEnd)

        round_join = OxmlElement("a:round")
        ln.append(round_join)

        tcPr.append(ln)


def _calc_vat_and_total_for_rates(net_rates: List[str], upload_fee: int, municipality_fee: int = 520) -> Tuple[List[str], List[str]]:
    vat_amounts = []
    total_amounts = []
    for net_rate_str in net_rates:
        net_rate = float(net_rate_str.replace("AED", "").replace(",", "").strip())
        subtotal = net_rate + upload_fee + municipality_fee
        vat = subtotal * 0.05
        total = subtotal + vat
        vat_amounts.append(f"AED {vat:,.0f}")
        total_amounts.append(f"AED {total:,.0f}")
    return vat_amounts, total_amounts


def _spots_text(spots: int) -> str:
    return f"{spots} Spot" + ("s" if spots != 1 else "")


def build_location_text(location_key: str, spots: int) -> str:
    """Build location description: Series: Location - Size (H x W) - Faces - Spots - Duration - SOV - Loop
    Format: Series: Location - Size (Height x Width) - Number of faces - Number of spots - Spot Duration x spots - SOV x spots - Loop duration
    """
    logger = config.logger
    logger.info(f"[BUILD_LOC_TEXT] Building text for location '{location_key}' with {spots} spots")
    
    # Get metadata from config (loaded from metadata.txt files)
    meta = config.LOCATION_METADATA.get(location_key.lower(), {})
    logger.info(f"[BUILD_LOC_TEXT] Metadata for '{location_key}': {meta}")
    
    # Extract values from metadata
    series = meta.get("series", "")
    location_name = meta.get("display_name", location_key.title())
    height = meta.get("height", "")
    width = meta.get("width", "")
    num_faces = meta.get("number_of_faces", 1)
    display_type = meta.get("display_type", "Digital").lower()
    spot_duration = meta.get("spot_duration", 16)
    loop_duration = meta.get("loop_duration", 96)
    base_sov = float(meta.get("sov", "16.6").replace("%", ""))
    
    # Build description parts
    parts = []
    
    # Series: Location
    if series:
        parts.append(f"{series}: {location_name}")
    else:
        parts.append(location_name)
    
    # Size (Height x Width)
    if height and width:
        # Check for "Multiple Sizes" special case
        if "multiple sizes" in str(height).lower() or "multiple sizes" in str(width).lower():
            parts.append("Multiple Sizes")
        else:
            # Remove 'm' suffix if present and re-add it
            h = str(height).replace('m', '').strip()
            w = str(width).replace('m', '').strip()
            parts.append(f"Size ({h}m x {w}m)")
    
    # Number of faces
    parts.append(f"{num_faces} faces")
    
    # For digital displays, add spot-related info
    if display_type == "digital":
        # Number of spots
        parts.append(f"{spots} {'spot' if spots == 1 else 'spots'}")

        # Spot Duration x Number of spots
        total_spot_duration = int(spot_duration) * spots
        parts.append(f"{total_spot_duration} Seconds")

        # SOV x Number of spots
        effective_sov = base_sov * spots
        parts.append(f"{effective_sov:.1f}% SOV")

        # Loop duration
        parts.append(f"{loop_duration} seconds loop")
    # Note: Static displays don't show spot numbers (makes no sense for static billboards)
    
    # Join all parts with " - "
    description = " - ".join(parts)
    
    logger.info(f"[BUILD_LOC_TEXT] Final description: '{description}'")
    return description


def create_financial_proposal_slide(slide, financial_data: dict, slide_width, slide_height) -> Tuple[List[str], List[str]]:
    logger = config.logger
    logger.info(f"[CREATE_FINANCIAL] Creating financial slide with data: {financial_data}")
    
    scale_x = slide_width / Inches(20)
    scale_y = slide_height / Inches(12)
    scale = min(scale_x, scale_y)

    rows = 9
    left = int(Inches(0.75) * scale_x)
    top = int(Inches(0.5) * scale_y)
    table_width = int(Inches(18.5) * scale_x)
    col1_width = int(Inches(4.0) * scale_x)
    col2_width = table_width - col1_width

    client_name = financial_data.get("client_name", "").strip()
    header_text = f"{client_name} Investment Sheet" if client_name else "Investment Sheet"

    location_name = financial_data["location"]
    start_date = format_date_for_display(financial_data["start_date"])
    end_date = format_date_for_display(financial_data.get("end_date", ""))
    date_range = f"{start_date} - {end_date}" if end_date else start_date
    durations = financial_data["durations"]
    net_rates = financial_data["net_rates"]
    spots = int(financial_data.get("spots", 1))
    production_fee_str = financial_data.get("production_fee")
    
    logger.info(f"[CREATE_FINANCIAL] Location: '{location_name}', Spots: {spots}")
    logger.info(f"[CREATE_FINANCIAL] Durations: {durations}, Net rates: {net_rates}")
    logger.info(f"[CREATE_FINANCIAL] Production fee: {production_fee_str}")

    location_text = build_location_text(location_name, spots)

    # Check if location is static
    location_meta = config.LOCATION_METADATA.get(location_name.lower(), {})
    is_static = location_meta.get('display_type', '').lower() == 'static'
    
    if is_static and production_fee_str:
        # Use production fee for static locations
        fee_str = production_fee_str
        fee_label = "Production Fee:"
        # Parse production fee to get numeric value
        production_fee = float(production_fee_str.replace("AED", "").replace(",", "").strip())
        upload_fee = production_fee
    else:
        # Use upload fee for digital locations
        upload_fee = config.UPLOAD_FEES_MAPPING.get(location_name.lower(), 3000)
        fee_str = f"AED {upload_fee:,}"
        fee_label = "Upload Fee:"
    
    municipality_fee = 520
    logger.info(f"[CREATE_FINANCIAL] Fee for '{location_name}': {fee_str} (static: {is_static})")

    vat_amounts, total_amounts = _calc_vat_and_total_for_rates(net_rates, upload_fee, municipality_fee)

    data = [
        (header_text, None),
        ("Location:", location_text),
        ("Start/End Date:", date_range),
        ("Duration:", durations if len(durations) > 1 else durations[0]),
        ("Net Rate:", net_rates if len(net_rates) > 1 else net_rates[0]),
        (fee_label, fee_str),
        ("Municipality Fee:", "AED 520 Per Image/Message"),
        ("VAT 5% :", vat_amounts if len(vat_amounts) > 1 else vat_amounts[0]),
        ("Total:", total_amounts if len(total_amounts) > 1 else total_amounts[0]),
    ]

    split_start_index = 3
    max_splits = max(len(v) if isinstance(v, list) else 1 for _, v in data[split_start_index:])
    cols = 1 + max_splits

    image_path = config.BASE_DIR / "image.png"
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width=table_width)

    row_height = int(Inches(0.9) * scale_y)
    table_height = int(row_height * rows)

    table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
    table = table_shape.table

    table.columns[0].width = col1_width
    split_col_width = int(col2_width / (cols - 1))
    for j in range(1, cols):
        table.columns[j].width = split_col_width

    for row in table.rows:
        row.height = int(table_height / rows)

    for i, (label, value) in enumerate(data):
        label_cell = table.cell(i, 0)

        if i == 0:
            label_cell.merge(table.cell(i, cols - 1))
            label_cell.fill.background()
            tf = label_cell.text_frame
            tf.clear()
            p_empty = tf.paragraphs[0]
            p_empty.text = " "
            p_empty.font.size = Pt(8)
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(int(36 * scale))
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            continue

        label_cell.text = label
        label_cell.fill.solid()
        if label == "Total:":
            label_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
        else:
            label_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        tf = label_cell.text_frame
        tf.clear()
        p_empty = tf.paragraphs[0]
        p_empty.text = " "
        p_empty.font.size = Pt(8)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(int(20 * scale))

        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.bold = False

        if label == "Total:":
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.bold = True
            run.font.size = Pt(int(28 * scale))
        elif label == "Net Rate:":
            run.font.bold = True

        if isinstance(value, list):
            for j, val in enumerate(value):
                val_cell = table.cell(i, j + 1)
                val_cell.text = val
                val_cell.fill.solid()
                if label == "Total:":
                    val_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
                else:
                    val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

                tf = val_cell.text_frame
                tf.clear()
                p_empty = tf.paragraphs[0]
                p_empty.text = " "
                p_empty.font.size = Pt(8)
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = val
                run.font.size = Pt(int(20 * scale))
                run.font.color.rgb = RGBColor(0, 0, 0)
                run.font.bold = label == "Net Rate:"

                if label == "Total:":
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                    run.font.size = Pt(int(28 * scale))
        else:
            val_cell = table.cell(i, 1)
            val_cell.merge(table.cell(i, cols - 1))
            val_cell.text = value
            val_cell.fill.solid()
            if label == "Total:":
                val_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
            else:
                val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            tf = val_cell.text_frame
            tf.clear()
            p_empty = tf.paragraphs[0]
            p_empty.text = " "
            p_empty.font.size = Pt(8)
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER

            if label == "Location":
                add_location_text_with_colored_sov(p, value, scale)
                run = None
            else:
                run = p.add_run()
                run.text = value
                run.font.size = Pt(int(20 * scale))
                run.font.color.rgb = RGBColor(0, 0, 0)
                run.font.bold = label == "Net Rate:"

            if label == "Total:" and run is not None:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(int(28 * scale))

    for row in table.rows:
        for cell in row.cells:
            set_cell_border(cell)

    table_element = table._tbl
    tblPr = table_element.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = OxmlElement('a:tblPr')
        table_element.insert(0, tblPr)
    for style in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(style)

    from datetime import datetime, timedelta
    validity_date = datetime.now() + timedelta(days=30)
    validity_date_str = validity_date.strftime("%d{} of %B, %Y").format(
        "st" if validity_date.day in [1, 21, 31] else
        "nd" if validity_date.day in [2, 22] else
        "rd" if validity_date.day in [3, 23] else
        "th"
    )

    # Get payment terms (default to 100% upfront) from financial_data
    payment_terms = financial_data.get("payment_terms", "100% upfront")

    bullet_text = f"""Payment Terms: {payment_terms}

• A DM fee of AED 520 per image/message applies. The final fee will be confirmed after the final artwork is received.
• An official booking order is required to secure the location/spot.
• Once a booking is confirmed, cancellations are not allowed even in case an artwork is rejected by the authorities, the client will be required to submit a revised artwork.
• All artworks are subject to approval by BackLite Media and DM.
• Location availability is subject to change.
• The artwork must comply with DM's guidelines.
• This proposal is valid until the {validity_date_str}."""

    # Smart T&C positioning: calculate available space and auto-scale font
    min_spacing = int(Inches(0.8) * scale_y)  # Minimum gap between table and T&C (increased to prevent overlap)
    bottom_margin = int(Inches(0.2) * scale_y)  # Space at bottom of slide

    bullet_top = table_shape.top + table_shape.height + min_spacing
    available_height = slide_height - bullet_top - bottom_margin

    # Ensure we have reasonable space
    if available_height < int(Inches(1.0) * scale_y):
        available_height = int(Inches(1.0) * scale_y)
        bullet_top = slide_height - available_height - bottom_margin

    bullet_box = slide.shapes.add_textbox(
        left=int(Inches(0.75) * scale_x),
        top=bullet_top,
        width=int(Inches(18.5) * scale_x),
        height=available_height,
    )

    tf = bullet_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    # Auto-scale font to fit available space (7 bullet points with line spacing)
    # Estimate: each line needs ~1.2x font size in height
    num_lines = 7  # Number of bullet points
    line_spacing = 1.15  # Tighter line spacing
    target_font_size = int((available_height / (num_lines * line_spacing * scale)) * 0.85)  # 0.85 for safety margin
    font_size = max(8, min(target_font_size, int(11 * scale)))  # Clamp between 8pt and 11pt*scale

    p = tf.paragraphs[0]
    p.text = bullet_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.line_spacing = line_spacing

    # Add proposal creation date to bottom right
    date_box = slide.shapes.add_textbox(
        left=int(Inches(15) * scale_x),
        top=int((slide_height - Inches(0.5)) * scale_y),
        width=int(Inches(4) * scale_x),
        height=int(Inches(0.4) * scale_y)
    )
    date_tf = date_box.text_frame
    date_p = date_tf.paragraphs[0]
    date_p.text = f"Proposal Date: {datetime.now().strftime('%d/%m/%Y')}"
    date_p.alignment = PP_ALIGN.RIGHT
    date_p.font.size = Pt(int(9 * scale))
    date_p.font.color.rgb = RGBColor(100, 100, 100)

    return vat_amounts, total_amounts


def create_combined_financial_proposal_slide(
    slide,
    proposals_data: list,
    combined_net_rate: str,
    slide_width,
    slide_height,
    client_name: str = "",
) -> str:
    logger = config.logger
    logger.info(f"[CREATE_COMBINED] Creating combined slide for {len(proposals_data)} locations")
    logger.info(f"[CREATE_COMBINED] Proposals data: {proposals_data}")
    logger.info(f"[CREATE_COMBINED] Combined net rate: {combined_net_rate}")

    scale_x = slide_width / Inches(20)
    scale_y = slide_height / Inches(12)
    scale = min(scale_x, scale_y)

    num_locations = len(proposals_data)  # Will be used for dynamic T&C spacing
    cols = num_locations + 1
    rows = 9

    left = int(Inches(0.75) * scale_x)
    top = int(Inches(0.5) * scale_y)
    table_width = int(Inches(18.5) * scale_x)
    col1_width = int(Inches(4.0) * scale_x)
    location_col_width = int((table_width - col1_width) / num_locations)

    image_path = config.BASE_DIR / "image.png"
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width=table_width)

    row_height = int(Inches(0.9) * scale_y)
    table_height = int(row_height * rows)

    table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
    table = table_shape.table

    table.columns[0].width = col1_width
    for j in range(1, cols):
        table.columns[j].width = location_col_width

    for row in table.rows:
        row.height = row_height

    locations = []
    start_dates = []
    durations = []
    upload_fees = []
    fee_label = "Upload Fee:"  # Default label
    has_static = False
    has_digital = False
    total_fees = 0

    for idx, proposal in enumerate(proposals_data):
        loc_name = proposal["location"]
        spots = int(proposal.get("spots", 1))
        production_fee_str = proposal.get("production_fee")
        logger.info(f"[CREATE_COMBINED] Processing location {idx + 1}: '{loc_name}' with {spots} spots")

        location_text = build_location_text(loc_name, spots)
        locations.append(location_text)

        # Format start/end date range
        start_date_fmt = format_date_for_display(proposal["start_date"])
        end_date_fmt = format_date_for_display(proposal.get("end_date", ""))
        date_range = f"{start_date_fmt} - {end_date_fmt}" if end_date_fmt else start_date_fmt
        start_dates.append(date_range)

        durations.append(proposal["durations"][0] if proposal["durations"] else "2 Weeks")
        
        # Check if location is static
        location_meta = config.LOCATION_METADATA.get(loc_name.lower(), {})
        is_static = location_meta.get('display_type', '').lower() == 'static'
        
        if is_static:
            has_static = True
            if production_fee_str:
                # Use production fee for static locations
                upload_fees.append(production_fee_str)
                # Parse production fee to get numeric value
                fee_numeric = float(production_fee_str.replace("AED", "").replace(",", "").strip())
                total_fees += fee_numeric
            else:
                # Fallback to stored fee
                fee = config.UPLOAD_FEES_MAPPING.get(loc_name.lower(), 3000)
                upload_fees.append(f"AED {fee:,}")
                total_fees += fee
        else:
            has_digital = True
            upload_fee = config.UPLOAD_FEES_MAPPING.get(loc_name.lower(), 3000)
            upload_fees.append(f"AED {upload_fee:,}")
            total_fees += upload_fee
        
        logger.info(f"[CREATE_COMBINED] Location {idx + 1} text: '{location_text}'")
        logger.info(f"[CREATE_COMBINED] Location {idx + 1} fee: {upload_fees[-1]} (static: {is_static})")

    # Determine fee label based on location types
    if has_static and has_digital:
        fee_label = "Upload/Production Fee:"
    elif has_static:
        fee_label = "Production Fee:"
    else:
        fee_label = "Upload Fee:"

    municipality_fee = 520
    total_upload_fees = total_fees  # Use calculated total fees

    net_rate_numeric = float(combined_net_rate.replace("AED", "").replace(",", "").strip())
    subtotal = net_rate_numeric + total_upload_fees + municipality_fee
    vat = subtotal * 0.05
    total = subtotal + vat

    header_text = f"{client_name.strip()} Investment Sheet" if client_name.strip() else "Investment Sheet"

    data = [
        (header_text, None),
        ("Location:", locations),
        ("Start/End Date:", start_dates),
        ("Duration:", durations),
        ("Net Rate:", combined_net_rate),
        (fee_label, upload_fees),
        ("Municipality Fee:", "AED 520 Per Image/Message"),
        ("VAT 5% :", f"AED {vat:,.0f}"),
        ("Total:", f"AED {total:,.0f}"),
    ]

    for i, (label, value) in enumerate(data):
        label_cell = table.cell(i, 0)
        if i == 0:
            label_cell.merge(table.cell(i, cols - 1))
            label_cell.fill.background()
            tf = label_cell.text_frame
            tf.clear()
            p_empty = tf.paragraphs[0]
            p_empty.text = " "
            p_empty.font.size = Pt(8)
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(int(36 * scale))
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            continue

        label_cell.text = label
        label_cell.fill.solid()
        if label == "Total:":
            label_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
        else:
            label_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        tf = label_cell.text_frame
        tf.clear()
        p_empty = tf.paragraphs[0]
        p_empty.text = " "
        p_empty.font.size = Pt(8)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(int(20 * scale))

        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.bold = False

        if label == "Total:":
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.bold = True
            run.font.size = Pt(int(28 * scale))
        elif label == "Net Rate:":
            run.font.bold = True

        if isinstance(value, list):
            for j, val in enumerate(value[:num_locations]):
                val_cell = table.cell(i, j + 1)
                val_cell.text = val
                val_cell.fill.solid()
                val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                tf = val_cell.text_frame
                tf.clear()
                p_empty = tf.paragraphs[0]
                p_empty.text = " "
                p_empty.font.size = Pt(8)
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                if label == "Location:":
                    add_location_text_with_colored_sov(p, val, scale)
                else:
                    run = p.add_run()
                    run.text = val
                    run.font.size = Pt(int(20 * scale))
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.bold = label == "Net Rate:"
        else:
            val_cell = table.cell(i, 1)
            val_cell.merge(table.cell(i, cols - 1))
            val_cell.text = value
            val_cell.fill.solid()
            if label == "Total:":
                val_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
            else:
                val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            tf = val_cell.text_frame
            tf.clear()
            p_empty = tf.paragraphs[0]
            p_empty.text = " "
            p_empty.font.size = Pt(8)
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = value
            run.font.size = Pt(int(20 * scale))
            run.font.color.rgb = RGBColor(0, 0, 0)
            run.font.bold = label == "Net Rate:"
            if label == "Total":
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(int(28 * scale))

    for row in table.rows:
        for cell in row.cells:
            set_cell_border(cell)

    from datetime import datetime, timedelta
    validity_date = datetime.now() + timedelta(days=30)
    validity_date_str = validity_date.strftime("%d{} of %B, %Y").format(
        "st" if validity_date.day in [1, 21, 31] else
        "nd" if validity_date.day in [2, 22] else
        "rd" if validity_date.day in [3, 23] else
        "th"
    )

    # Get payment terms (default to 100% upfront) from proposals_data
    payment_terms = "100% upfront"
    if isinstance(proposals_data, list) and len(proposals_data) > 0:
        payment_terms = proposals_data[0].get("payment_terms", "100% upfront")

    bullet_text = f"""Payment Terms: {payment_terms}

• A DM fee of AED 520 per image/message applies. The final fee will be confirmed after the final artwork is received.
• An official booking order is required to secure the location/spot.
• Once a booking is confirmed, cancellations are not allowed even in case an artwork is rejected by the authorities, the client will be required to submit a revised artwork.
• All artworks are subject to approval by BackLite Media and DM.
• Location availability is subject to change.
• The artwork must comply with DM's guidelines.
• This proposal is valid until the {validity_date_str}."""

    # Smart T&C positioning: calculate available space and auto-scale font
    # Scale spacing based on number of locations (more locations = taller location row = need more spacing)
    # Formula calibrated so 3 locations = 0.8" (known good value)
    base_spacing = 0.5  # Base spacing in inches
    location_factor = num_locations * 0.1  # Add 0.1" per location
    min_spacing = int(Inches(base_spacing + location_factor) * scale_y)
    bottom_margin = int(Inches(0.2) * scale_y)  # Space at bottom of slide

    bullet_top = table_shape.top + table_shape.height + min_spacing
    available_height = slide_height - bullet_top - bottom_margin

    # Ensure we have reasonable space
    if available_height < int(Inches(1.0) * scale_y):
        available_height = int(Inches(1.0) * scale_y)
        bullet_top = slide_height - available_height - bottom_margin

    bullet_box = slide.shapes.add_textbox(
        left=int(Inches(0.75) * scale_x),
        top=bullet_top,
        width=int(Inches(18.5) * scale_x),
        height=available_height,
    )

    tf = bullet_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    # Auto-scale font to fit available space (7 bullet points with line spacing)
    # Estimate: each line needs ~1.2x font size in height
    num_lines = 7  # Number of bullet points
    line_spacing = 1.15  # Tighter line spacing
    target_font_size = int((available_height / (num_lines * line_spacing * scale)) * 0.85)  # 0.85 for safety margin
    font_size = max(8, min(target_font_size, int(11 * scale)))  # Clamp between 8pt and 11pt*scale

    p = tf.paragraphs[0]
    p.text = bullet_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.line_spacing = line_spacing

    # Add proposal creation date to bottom right
    date_box = slide.shapes.add_textbox(
        left=int(Inches(15) * scale_x),
        top=int((slide_height - Inches(0.5)) * scale_y),
        width=int(Inches(4) * scale_x),
        height=int(Inches(0.4) * scale_y)
    )
    date_tf = date_box.text_frame
    date_p = date_tf.paragraphs[0]
    date_p.text = f"Proposal Date: {datetime.now().strftime('%d/%m/%Y')}"
    date_p.alignment = PP_ALIGN.RIGHT
    date_p.font.size = Pt(int(9 * scale))
    date_p.font.color.rgb = RGBColor(100, 100, 100)

    return f"AED {total:,.0f}" 