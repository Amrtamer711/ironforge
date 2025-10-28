import json
import asyncio
from typing import Dict, Any, Optional
import os
from pathlib import Path
import aiohttp
from datetime import datetime, timedelta
from pptx import Presentation
import shutil

import config
import db
from proposals import process_proposals
from slack_formatting import SlackResponses
from booking_parser import BookingOrderParser, COMBINED_BOS_DIR
from task_queue import mockup_queue
import bo_slack_messaging

user_history: Dict[str, list] = {}

# Global for pending location additions (waiting for PPT upload)
pending_location_additions: Dict[str, Dict[str, Any]] = {}

# Global for mockup history (30-minute memory per user)
# Structure: {user_id: {"creative_paths": List[Path], "metadata": dict, "timestamp": datetime}}
# Stores individual creative files (1-N files) so they can be reused on different locations with matching frame count
mockup_history: Dict[str, Dict[str, Any]] = {}

# Global for booking order draft review sessions (active until approved/cancelled)
# Structure: {user_id: {"data": dict, "warnings": List[str], "missing_required": List[str],
#                        "original_file_path": Path, "company": str, "file_type": str}}
pending_booking_orders: Dict[str, Dict[str, Any]] = {}

def cleanup_expired_mockups():
    """Remove creative files that have expired (older than 30 minutes)"""
    import os
    import gc

    now = datetime.now()
    expired_users = []

    for user_id, data in mockup_history.items():
        timestamp = data.get("timestamp")
        if timestamp and (now - timestamp) > timedelta(minutes=30):
            # Delete all creative files for this user
            creative_paths = data.get("creative_paths", [])
            deleted_count = 0
            for creative_path in creative_paths:
                if creative_path and creative_path.exists():
                    try:
                        os.unlink(creative_path)
                        deleted_count += 1
                    except Exception as e:
                        config.logger.error(f"[MOCKUP HISTORY] Failed to delete {creative_path}: {e}")

            if deleted_count > 0:
                config.logger.info(f"[MOCKUP HISTORY] Cleaned up {deleted_count} expired creative file(s) for user {user_id}")
            expired_users.append(user_id)

    # Remove from memory
    for user_id in expired_users:
        del mockup_history[user_id]
        config.logger.info(f"[MOCKUP HISTORY] Removed user {user_id} from mockup history")

    # Force garbage collection if we cleaned up any files
    if expired_users:
        gc.collect()
        config.logger.info(f"[MOCKUP HISTORY] Forced garbage collection after cleanup")

def store_mockup_history(user_id: str, creative_paths: list, metadata: dict):
    """Store creative files in user's history with 30-minute expiry

    Args:
        user_id: Slack user ID
        creative_paths: List of Path objects to creative files (1-N files)
        metadata: Dict with location_key, location_name, num_frames, etc.
    """
    import gc

    # Clean up old creative files for this user if exists
    if user_id in mockup_history:
        old_data = mockup_history[user_id]
        old_creative_paths = old_data.get("creative_paths", [])
        deleted_count = 0
        for old_path in old_creative_paths:
            if old_path and old_path.exists():
                try:
                    os.unlink(old_path)
                    deleted_count += 1
                except Exception as e:
                    config.logger.error(f"[MOCKUP HISTORY] Failed to delete old creative: {e}")
        if deleted_count > 0:
            config.logger.info(f"[MOCKUP HISTORY] Replaced {deleted_count} old creative file(s) for user {user_id}")
            # Force garbage collection when replacing files
            gc.collect()

    # Store new creative files
    mockup_history[user_id] = {
        "creative_paths": creative_paths,
        "metadata": metadata,
        "timestamp": datetime.now()
    }
    config.logger.info(f"[MOCKUP HISTORY] Stored {len(creative_paths)} creative file(s) for user {user_id}")

    # Run cleanup to remove expired creatives from other users
    cleanup_expired_mockups()

def get_mockup_history(user_id: str) -> Optional[Dict[str, Any]]:
    """Get user's creative files from history if still valid (within 30 minutes)

    Returns:
        Dict with creative_paths (List[Path]), metadata, timestamp, or None if expired/not found
    """
    import gc

    if user_id not in mockup_history:
        return None

    data = mockup_history[user_id]
    timestamp = data.get("timestamp")

    # Check if expired
    if timestamp and (datetime.now() - timestamp) > timedelta(minutes=30):
        # Expired - clean up all creative files
        creative_paths = data.get("creative_paths", [])
        deleted_count = 0
        for creative_path in creative_paths:
            if creative_path and creative_path.exists():
                try:
                    os.unlink(creative_path)
                    deleted_count += 1
                except:
                    pass
        del mockup_history[user_id]

        # Force garbage collection if we deleted files
        if deleted_count > 0:
            gc.collect()
            config.logger.info(f"[MOCKUP HISTORY] Auto-cleaned {deleted_count} expired file(s) for user {user_id}")

        return None

    return data

def get_location_frame_count(location_key: str, time_of_day: str = "all", finish: str = "all") -> Optional[int]:
    """Get the number of frames for a specific location configuration.

    Returns:
        Number of frames, or None if location not found or no mockups configured
    """
    import db

    # Get available variations for the location
    variations = db.list_mockup_variations(location_key)
    if not variations:
        return None

    # Get the first available variation that matches time_of_day/finish
    # variations structure: {'day': ['gold', 'silver'], 'night': ['gold']}
    for tod, finish_list in variations.items():
        if time_of_day != "all" and tod != time_of_day:
            continue

        for fin in finish_list:
            if finish != "all" and fin != finish:
                continue

            # Get all photos for this time_of_day/finish combination
            photos = db.list_mockup_photos(location_key, tod, fin)
            if photos:
                # Get frames data for the first photo
                frames_data = db.get_mockup_frames(location_key, photos[0], tod, fin)
                if frames_data:
                    return len(frames_data)

    return None

def _validate_pdf_file(file_path: Path) -> bool:
    """Validate that uploaded file is actually a PDF."""
    try:
        # Check if file exists first
        if not file_path.exists():
            config.logger.error(f"[VALIDATION] File does not exist: {file_path}")
            return False

        file_size = file_path.stat().st_size
        config.logger.info(f"[VALIDATION] Validating PDF file: {file_path} (size: {file_size} bytes)")

        # Quick sanity checks before invoking PyMuPDF
        if file_size <= 0:
            config.logger.warning(f"[VALIDATION] PDF file is empty: {file_path}")
            return False

        # Try to open with PyMuPDF
        import fitz
        doc = fitz.open(file_path)
        page_count = len(doc)
        doc.close()

        if page_count <= 0:
            config.logger.warning(f"[VALIDATION] PDF has no pages: {file_path}")
            return False

        config.logger.info(f"[VALIDATION] ✓ Valid PDF with {page_count} pages")
        return True
    except Exception as e:
        config.logger.error(f"[VALIDATION] Invalid PDF file: {e}")
        return False


async def _convert_pdf_to_pptx(pdf_path: Path) -> Optional[Path]:
    """Convert PDF to PowerPoint with maximum quality (300 DPI).

    Args:
        pdf_path: Path to input PDF file

    Returns:
        Path to converted PPTX file, or None if conversion failed
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import tempfile
        import os

        config.logger.info(f"[PDF_CONVERT] Starting PDF to PPTX conversion: {pdf_path}")

        # Open PDF
        doc = fitz.open(pdf_path)
        page_count = len(doc)
        config.logger.info(f"[PDF_CONVERT] PDF has {page_count} pages")

        # Create temporary directory for images
        temp_dir = tempfile.mkdtemp(prefix="pdf_convert_")

        # Convert pages to high-resolution images (4x zoom = ~300 DPI)
        zoom = 4.0
        matrix = fitz.Matrix(zoom, zoom)
        image_paths = []

        for page_num in range(page_count):
            page = doc[page_num]
            pix = page.get_pixmap(matrix=matrix, alpha=False)

            # Save as PNG
            img_path = os.path.join(temp_dir, f"page_{page_num + 1}.png")
            pix.save(img_path)
            image_paths.append(img_path)
            config.logger.info(f"[PDF_CONVERT] Extracted page {page_num + 1}/{page_count} ({pix.width}x{pix.height}px)")

        doc.close()

        # Create PowerPoint presentation
        config.logger.info("[PDF_CONVERT] Creating PowerPoint presentation...")
        prs = Presentation()

        # Set slide size to 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Add pages as slides
        for i, img_path in enumerate(image_paths, 1):
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Get image dimensions
            img = Image.open(img_path)
            img_width, img_height = img.size
            img_aspect = img_width / img_height
            slide_aspect = float(prs.slide_width) / float(prs.slide_height)

            # Calculate dimensions to fill slide while maintaining aspect ratio
            if img_aspect > slide_aspect:
                # Image is wider - fit to width
                width = prs.slide_width
                height = int(float(prs.slide_width) / img_aspect)
                left = 0
                top = int((prs.slide_height - height) / 2)
            else:
                # Image is taller - fit to height
                height = prs.slide_height
                width = int(float(prs.slide_height) * img_aspect)
                left = int((prs.slide_width - width) / 2)
                top = 0

            # Add image to slide
            slide.shapes.add_picture(img_path, left, top, width, height)
            config.logger.info(f"[PDF_CONVERT] Added slide {i}/{page_count}")

        # Save PPTX to temporary file
        pptx_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(pptx_temp.name)
        pptx_temp.close()

        # Cleanup temporary images
        for img_path in image_paths:
            try:
                os.remove(img_path)
            except:
                pass
        try:
            os.rmdir(temp_dir)
        except:
            pass

        file_size_mb = os.path.getsize(pptx_temp.name) / (1024 * 1024)
        config.logger.info(f"[PDF_CONVERT] ✓ Conversion complete: {pptx_temp.name} ({file_size_mb:.1f} MB, {page_count} slides, 300 DPI)")

        return Path(pptx_temp.name)

    except Exception as e:
        config.logger.error(f"[PDF_CONVERT] Conversion failed: {e}", exc_info=True)
        return None


def _validate_powerpoint_file(file_path: Path) -> bool:
    """Validate that uploaded file is actually a PowerPoint presentation."""
    try:
        # Check if file exists first
        if not file_path.exists():
            config.logger.error(f"[VALIDATION] File does not exist: {file_path}")
            return False

        file_size = file_path.stat().st_size
        config.logger.info(f"[VALIDATION] Validating PowerPoint file: {file_path} (size: {file_size} bytes)")

        # Quick sanity checks before invoking python-pptx
        if file_size <= 0:
            config.logger.warning(f"[VALIDATION] PowerPoint file is empty: {file_path}")
            return False

        # PPTX files are ZIP packages starting with 'PK\x03\x04'
        try:
            with open(file_path, 'rb') as fp:
                magic = fp.read(4)
            if magic != b'PK\x03\x04':
                config.logger.warning(
                    f"[VALIDATION] File does not look like a PPTX (ZIP signature missing). "
                    f"Likely an HTML error/permission issue from Slack. Path: {file_path}"
                )
                return False
        except Exception as e:
            config.logger.warning(f"[VALIDATION] Failed to read file header for {file_path}: {e}")
            return False

        # Try to open as PowerPoint - this will fail if not a valid PPTX
        pres = Presentation(str(file_path))
        # Basic validation: must have at least 1 slide
        slide_count = len(pres.slides)

        # CRITICAL: Delete presentation object to free memory
        # python-pptx loads entire PPT into RAM (50-100MB+)
        del pres
        import gc
        gc.collect()

        if slide_count < 1:
            config.logger.warning(f"[VALIDATION] PowerPoint file has no slides: {file_path}")
            return False

        config.logger.info(f"[VALIDATION] PowerPoint validation successful: {slide_count} slides")
        return True
    except Exception as e:
        config.logger.warning(f"[VALIDATION] PowerPoint validation failed: {e}")
        # Cleanup on error path too
        try:
            del pres
        except:
            pass
        import gc
        gc.collect()
        return False


async def _download_slack_file(file_info: Dict[str, Any]) -> Path:
    url = file_info.get("url_private_download") or file_info.get("url_private")
    if not url:
        raise ValueError("Missing file download URL")
    headers = {"Authorization": f"Bearer {config.SLACK_BOT_TOKEN}"}
    suffix = Path(file_info.get("name", "upload.bin")).suffix or ".bin"
    import tempfile
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    tmp.close()

    config.logger.info(f"[DOWNLOAD] Downloading file to: {tmp.name}")

    async with aiohttp.ClientSession() as session:
        async with session.get(url, headers=headers) as resp:
            resp.raise_for_status()
            content_type = resp.headers.get('Content-Type', '')
            config.logger.info(f"[DOWNLOAD] HTTP {resp.status}, Content-Type: {content_type}")
            content = await resp.read()
            config.logger.info(f"[DOWNLOAD] Downloaded {len(content)} bytes")
            with open(tmp.name, "wb") as f:
                f.write(content)

            # Immediately delete content bytes to free memory
            del content

    # Verify file was written
    file_path = Path(tmp.name)
    if file_path.exists():
        config.logger.info(f"[DOWNLOAD] File successfully written: {file_path} (size: {file_path.stat().st_size} bytes)")
    else:
        config.logger.error(f"[DOWNLOAD] File not found after write: {file_path}")

    return file_path


async def _generate_mockup_queued(
    location_key: str,
    creative_paths: list,
    time_of_day: str,
    finish: str,
    specific_photo: str = None,
    config_override: dict = None
):
    """
    Wrapper function for mockup generation that runs through the task queue.
    This limits concurrent mockup generation to prevent memory exhaustion.

    Args:
        location_key: Location identifier
        creative_paths: List of creative file paths
        time_of_day: Time of day variation
        finish: Finish type
        specific_photo: Optional specific photo to use
        config_override: Optional config override

    Returns:
        Tuple of (result_path, metadata)
    """
    import mockup_generator
    import gc

    logger = config.logger
    logger.info(f"[QUEUE] Mockup generation requested for {location_key}")

    # This function will be queued and executed when a slot is available
    async def _generate():
        try:
            logger.info(f"[QUEUE] Starting mockup generation for {location_key}")
            result_path, metadata = mockup_generator.generate_mockup(
                location_key,
                creative_paths,
                time_of_day=time_of_day,
                finish=finish,
                specific_photo=specific_photo,
                config_override=config_override
            )
            logger.info(f"[QUEUE] Mockup generation completed for {location_key}")
            gc.collect()
            return result_path, metadata
        except Exception as e:
            logger.error(f"[QUEUE] Mockup generation failed for {location_key}: {e}")
            raise

    # Submit to queue and wait for result
    return await mockup_queue.submit(_generate)


async def _generate_ai_mockup_queued(
    ai_prompt: str,
    enhanced_prompt: str,
    num_ai_frames: int,
    location_key: str,
    time_of_day: str,
    finish: str
):
    """
    Wrapper for AI mockup generation (AI creative generation + mockup) through the queue.
    This ensures the entire AI workflow (fetch from OpenAI + image processing + mockup)
    is treated as ONE queued task to prevent memory spikes.

    Args:
        ai_prompt: User's original AI prompt
        enhanced_prompt: Full enhanced system prompt
        num_ai_frames: Number of frames to generate
        location_key: Location identifier
        time_of_day: Time of day variation
        finish: Finish type

    Returns:
        Tuple of (result_path, ai_creative_paths)
    """
    import mockup_generator
    import gc

    logger = config.logger
    logger.info(f"[QUEUE] AI mockup requested for {location_key} ({num_ai_frames} frames)")

    async def _generate():
        try:
            logger.info(f"[QUEUE] Generating {num_ai_frames} AI creative(s) for {location_key}")
            ai_creative_paths = []

            if num_ai_frames > 1:
                # Multi-frame: parse prompt into variations
                prompt_variations = await mockup_generator.parse_prompt_for_multi_frame(ai_prompt, num_ai_frames)

                # Generate each creative
                for i, variation_prompt in enumerate(prompt_variations, 1):
                    logger.info(f"[AI QUEUE] Generating creative {i}/{num_ai_frames}")
                    creative_path = await mockup_generator.generate_ai_creative(
                        prompt=enhanced_prompt.replace(ai_prompt, variation_prompt),
                        location_key=location_key
                    )
                    if not creative_path:
                        raise Exception(f"Failed to generate AI creative {i}/{num_ai_frames}")
                    ai_creative_paths.append(creative_path)
            else:
                # Single frame
                creative_path = await mockup_generator.generate_ai_creative(
                    prompt=enhanced_prompt,
                    location_key=location_key
                )
                if not creative_path:
                    raise Exception("Failed to generate AI creative")
                ai_creative_paths.append(creative_path)

            logger.info(f"[QUEUE] AI creatives ready, generating mockup for {location_key}")

            # Generate mockup with AI creatives
            result_path, metadata = mockup_generator.generate_mockup(
                location_key,
                ai_creative_paths,
                time_of_day=time_of_day,
                finish=finish
            )

            if not result_path:
                raise Exception("Failed to generate mockup")

            logger.info(f"[QUEUE] AI mockup completed for {location_key}")
            gc.collect()
            return result_path, ai_creative_paths

        except Exception as e:
            logger.error(f"[QUEUE] AI mockup failed for {location_key}: {e}")
            raise

    # Submit entire AI workflow to queue as ONE task
    return await mockup_queue.submit(_generate)


async def _persist_location_upload(location_key: str, pptx_path: Path, metadata_text: str) -> None:
    location_dir = config.TEMPLATES_DIR / location_key
    location_dir.mkdir(parents=True, exist_ok=True)
    target_pptx = location_dir / f"{location_key}.pptx"
    target_meta = location_dir / "metadata.txt"
    # Move/copy files
    import shutil
    shutil.move(str(pptx_path), str(target_pptx))
    target_meta.write_text(metadata_text, encoding="utf-8")


async def _handle_booking_order_parse(
    company: str,
    slack_event: Dict[str, Any],
    channel: str,
    status_ts: str,
    user_notes: str,
    user_id: str,
    user_message: str = ""
):
    """Handle booking order parsing workflow"""
    logger = config.logger

    # Extract files from slack event
    files = slack_event.get("files", [])
    if not files and slack_event.get("subtype") == "file_share" and "file" in slack_event:
        files = [slack_event["file"]]

    if not files:
        await config.slack_client.chat_update(
            channel=channel,
            ts=status_ts,
            text=config.markdown_to_slack("❌ No file detected. Please upload a booking order document (Excel, PDF, or image).")
        )
        return

    file_info = files[0]
    logger.info(f"[BOOKING] Processing file: {file_info.get('name')}")

    # Download file
    try:
        tmp_file = await _download_slack_file(file_info)
    except Exception as e:
        logger.error(f"[BOOKING] Download failed: {e}")
        await config.slack_client.chat_update(
            channel=channel,
            ts=status_ts,
            text=config.markdown_to_slack(f"❌ Failed to download file: {e}")
        )
        return

    # Initialize parser
    parser = BookingOrderParser(company=company)
    file_type = parser.detect_file_type(tmp_file)

    # Classify document
    try:
        await config.slack_client.chat_update(channel=channel, ts=status_ts, text="⏳ _Classifying document..._")
    except Exception as e:
        logger.error(f"[SLACK] Failed to update status message while classifying: {e}", exc_info=True)
        # Continue processing - status update failure shouldn't stop the workflow

    classification = await parser.classify_document(tmp_file, user_message=user_message)
    logger.info(f"[BOOKING] Classification: {classification}")

    # Check if it's actually a booking order
    if classification.get("classification") != "BOOKING_ORDER" or classification.get("confidence") in {"low", None}:
        try:
            await config.slack_client.chat_update(
                channel=channel,
                ts=status_ts,
                text=config.markdown_to_slack(
                    f"⚠️ This doesn't look like a booking order (confidence: {classification.get('confidence', 'unknown')}).\n\n"
                    f"Reasoning: {classification.get('reasoning', 'N/A')}\n\n"
                    f"If this is artwork for a mockup, please request a mockup instead."
                )
            )
        except Exception as e:
            logger.error(f"[SLACK] Failed to send classification result to user: {e}", exc_info=True)
        tmp_file.unlink(missing_ok=True)
        return

    # Parse the booking order
    try:
        await config.slack_client.chat_update(channel=channel, ts=status_ts, text="⏳ _Extracting booking order data..._")
    except Exception as e:
        logger.error(f"[SLACK] Failed to update status message while parsing: {e}", exc_info=True)
    try:
        result = await parser.parse_file(tmp_file, file_type, user_message=user_message)
    except Exception as e:
        logger.error(f"[BOOKING] Parsing failed: {e}", exc_info=True)
        try:
            await config.slack_client.chat_update(
                channel=channel,
                ts=status_ts,
                text=config.markdown_to_slack(
                    f"❌ **Error:** Failed to extract data from the booking order.\n\n"
                    f"If you believe this is a bug, please contact the AI team with the timestamp: `{datetime.now().strftime('%Y-%m-%d %H:%M:%S UTC')}`"
                )
            )
        except Exception as slack_error:
            logger.error(f"[SLACK] Failed to send parsing error to user: {slack_error}", exc_info=True)
        tmp_file.unlink(missing_ok=True)
        return

    # NEW FLOW: Generate Excel immediately and send with Approve/Reject buttons
    try:
        await config.slack_client.chat_update(channel=channel, ts=status_ts, text="⏳ _Generating Excel file..._")
    except Exception as e:
        logger.error(f"[SLACK] Failed to update status message while generating Excel: {e}", exc_info=True)

    try:
        # Generate combined PDF (Excel + Original BO concatenated)
        import bo_approval_workflow
        combined_pdf_path = await parser.generate_combined_pdf(result.data, f"DRAFT_{company}_REVIEW", Path(tmp_file))

        # Create approval workflow
        workflow_id = await bo_approval_workflow.create_approval_workflow(
            user_id=user_id,
            company=company,
            data=result.data,
            warnings=result.warnings,
            missing_required=result.missing_required,
            original_file_path=tmp_file,
            original_filename=file_info.get("name"),
            file_type=file_type,
            user_notes=user_notes
        )

        # Get coordinator channel (uses conversations.open to get DM channel ID)
        coordinator_channel = await bo_approval_workflow.get_coordinator_channel(company)
        logger.info(f"[BO APPROVAL] Coordinator channel for {company}: {coordinator_channel}")
        if not coordinator_channel:
            try:
                await config.slack_client.chat_update(
                    channel=channel,
                    ts=status_ts,
                    text=config.markdown_to_slack(f"❌ **Error:** Sales Coordinator for {company} not configured. Please contact the AI team to set this up.")
                )
            except Exception as e:
                logger.error(f"[SLACK] Failed to send config error to user: {e}", exc_info=True)
            return

        # Send Excel + summary + Approve/Reject buttons to coordinator
        submitter_name = await bo_slack_messaging.get_user_real_name(user_id)
        preview_text = f"📋 **New Booking Order - Ready for Approval**\n\n"
        preview_text += f"**Company:** {company.upper()}\n"
        preview_text += f"**Submitted by:** {submitter_name}\n\n"
        preview_text += f"**Client:** {result.data.get('client', 'N/A')}\n"
        preview_text += f"**Campaign:** {result.data.get('brand_campaign', 'N/A')}\n"
        preview_text += f"**BO Number:** {result.data.get('bo_number', 'N/A')}\n"
        preview_text += f"**Net (pre-VAT):** AED {result.data.get('net_pre_vat', 0):,.2f}\n"
        preview_text += f"**VAT (5%):** AED {result.data.get('vat_calc', 0):,.2f}\n"
        preview_text += f"**Gross Total:** AED {result.data.get('gross_calc', 0):,.2f}\n"

        if result.data.get("locations"):
            preview_text += f"\n**Locations:** {len(result.data['locations'])}\n"
            for loc in result.data["locations"][:3]:  # Show first 3
                preview_text += f"  • {loc.get('name', 'Unknown')}: {loc.get('start_date', '?')} to {loc.get('end_date', '?')} (AED {loc.get('net_amount', 0):,.2f})\n"
            if len(result.data["locations"]) > 3:
                preview_text += f"  ...and {len(result.data['locations']) - 3} more\n"

        if result.warnings:
            preview_text += "\n⚠️ **Warnings:**\n" + "\n".join(f"• {w}" for w in result.warnings)

        if result.missing_required:
            preview_text += "\n❌ **Missing Required:**\n" + "\n".join(f"• {m}" for m in result.missing_required)

        if user_notes:
            preview_text += f"\n📝 **Sales Notes:** {user_notes}\n"

        preview_text += "\n\n📎 **Please review the Excel file attached below, then:**\n"
        preview_text += "• Press **Approve** to send to Head of Sales\n"
        preview_text += "• Press **Reject** to request changes in a thread"

        # NEW FLOW: Post notification in main channel, then file + buttons in thread
        try:
            await config.slack_client.chat_update(channel=channel, ts=status_ts, text="⏳ _Sending to coordinator..._")
        except Exception as e:
            logger.error(f"[SLACK] Failed to update status message: {e}", exc_info=True)

        logger.info(f"[BO APPROVAL] Posting notification to coordinator channel: {coordinator_channel}")

        # Get submitter's real name
        submitter_name = await bo_slack_messaging.get_user_real_name(user_id)

        # Step 1: Post notification message in main channel
        notification_text = (
            f"📋 **New Booking Order Submitted**\n\n"
            f"**Client:** {result.data.get('client', 'N/A')}\n"
            f"**Campaign:** {result.data.get('brand_campaign', 'N/A')}\n"
            f"**Gross Total:** AED {result.data.get('gross_calc', 0):,.2f}\n\n"
            f"**Submitted by:** {submitter_name}\n\n"
            f"_Please review the details in the thread below..._"
        )

        notification_msg = await config.slack_client.chat_postMessage(
            channel=coordinator_channel,
            text=config.markdown_to_slack(notification_text)
        )
        notification_ts = notification_msg["ts"]
        logger.info(f"[BO APPROVAL] Posted notification with ts: {notification_ts}")

        # Step 2: Upload combined PDF file as threaded reply
        logger.info(f"[BO APPROVAL] Uploading combined PDF in thread...")
        try:
            file_upload = await config.slack_client.files_upload_v2(
                channel=coordinator_channel,
                file=str(combined_pdf_path),
                title=f"BO Draft - {result.data.get('client', 'Unknown')}",
                initial_comment=config.markdown_to_slack(preview_text),
                thread_ts=notification_ts  # Post in thread
            )
            logger.info(f"[BO APPROVAL] Combined PDF uploaded in thread successfully")
        except Exception as upload_error:
            logger.error(f"[BO APPROVAL] Failed to upload combined PDF: {upload_error}", exc_info=True)
            raise Exception(f"Failed to send combined PDF to coordinator. Channel/User ID: {coordinator_channel}")

        # Wait for file to fully appear in Slack before posting buttons
        logger.info(f"[BO APPROVAL] Waiting 10 seconds for file to render in Slack...")
        await asyncio.sleep(10)

        # Step 3: Post buttons in the same thread
        logger.info(f"[BO APPROVAL] Posting approval buttons in thread...")
        button_blocks = [
            {
                "type": "section",
                "text": {
                    "type": "mrkdwn",
                    "text": "📎 *Please review the PDF above (Excel + Original BO), then:*"
                }
            },
            {
                "type": "actions",
                "elements": [
                    {
                        "type": "button",
                        "text": {"type": "plain_text", "text": "✅ Approve"},
                        "style": "primary",
                        "value": workflow_id,
                        "action_id": "approve_bo_coordinator"
                    },
                    {
                        "type": "button",
                        "text": {"type": "plain_text", "text": "❌ Reject"},
                        "style": "danger",
                        "value": workflow_id,
                        "action_id": "reject_bo_coordinator"
                    }
                ]
            }
        ]

        button_msg = await config.slack_client.chat_postMessage(
            channel=coordinator_channel,
            thread_ts=notification_ts,  # Post in same thread
            text="Please review and approve or reject:",
            blocks=button_blocks
        )
        coordinator_msg_ts = button_msg["ts"]
        logger.info(f"[BO APPROVAL] Posted buttons in thread with ts: {coordinator_msg_ts}")

        # Update workflow with coordinator message info
        await bo_approval_workflow.update_workflow(workflow_id, {
            "coordinator_thread_channel": coordinator_channel,
            "coordinator_thread_ts": notification_ts,  # The notification message is the thread root
            "coordinator_msg_ts": coordinator_msg_ts,  # The button message
            "combined_pdf_path": str(combined_pdf_path)
        })

        # Notify sales person
        try:
            await config.slack_client.chat_update(
                channel=channel,
                ts=status_ts,
                text=config.markdown_to_slack(
                    f"✅ **Booking Order Submitted**\n\n"
                    f"**Client:** {result.data.get('client', 'N/A')}\n"
                    f"**Campaign:** {result.data.get('brand_campaign', 'N/A')}\n"
                    f"**Gross Total:** AED {result.data.get('gross_calc', 0):,.2f}\n\n"
                    f"Your booking order has been sent to the Sales Coordinator with a combined PDF (parsed data + original BO) for immediate review. "
                    f"You'll be notified once the approval process is complete."
                )
            )
        except Exception as e:
            logger.error(f"[SLACK] Failed to send success message to user: {e}", exc_info=True)

        logger.info(f"[BO APPROVAL] Sent {workflow_id} to coordinator with combined PDF and approval buttons")

    except Exception as e:
        logger.error(f"[BO APPROVAL] Error creating workflow: {e}", exc_info=True)
        try:
            await config.slack_client.chat_update(
                channel=channel,
                ts=status_ts,
                text=config.markdown_to_slack(
                    f"❌ **Error:** Failed to start the approval workflow.\n\n"
                    f"If you believe this is a bug, please contact the AI team with the timestamp: `{datetime.now().strftime('%Y-%m-%d %H:%M:%S UTC')}`"
                )
            )
        except Exception as slack_error:
            logger.error(f"[SLACK] Failed to send workflow error to user: {slack_error}", exc_info=True)


async def handle_booking_order_edit_flow(channel: str, user_id: str, user_input: str) -> str:
    """Handle booking order edit flow with structured LLM response"""
    try:
        edit_data = pending_booking_orders.get(user_id, {})
        current_data = edit_data.get("data", {})
        warnings = edit_data.get("warnings", [])
        missing_required = edit_data.get("missing_required", [])

        # Build system prompt for LLM to parse user intent
        system_prompt = f"""
You are an intelligent assistant helping review and edit a booking order draft for a billboard/outdoor advertising company.

**TAKE YOUR TIME AND BE INTELLIGENT:**
When the user requests changes, think carefully about ALL fields that need to be updated. Don't just update what they explicitly mention - understand the cascading effects and update ALL related fields automatically.

**USER SAID:** "{user_input}"

**BILLBOARD INDUSTRY CONTEXT (Critical for intelligent edits):**

1. **Multiple Locations Under One Payment:**
   - Clients buy multiple billboard locations, sometimes bundled under one payment
   - If user adds a location, you may need to adjust payment splits across locations

2. **Fee Types (Municipality, Upload, Production):**
   - Municipality fee applies to ALL locations
   - Upload fee is for DIGITAL locations only
   - Production fee is for STATIC locations only
   - These are separate from location rentals

3. **SLA (Service Level Agreement) %:**
   - This is a DEDUCTION applied to net rental (rental + fees + municipality)
   - Applied BEFORE VAT
   - Formula: Net after SLA = Net rental - (Net rental × SLA%)
   - VAT is then: (Net after SLA) × 0.05
   - Gross = Net after SLA + VAT

4. **Cascading Calculations:**
   When user changes certain fields, you MUST automatically update related fields:

   **If they add/remove a location:**
   - Update "asset" list to include/exclude the location code
   - Update "locations" array with full location details (name, dates, costs, etc.)
   - Recalculate "net_pre_vat" (sum of all location rentals + fees)
   - Recalculate "vat_calc" (net × 0.05 or apply SLA first if present)
   - Recalculate "gross_calc" (net + VAT)

   **If they change any fee (rental, production, upload, municipality):**
   - Recalculate "net_pre_vat"
   - Recalculate SLA deduction if SLA% exists
   - Recalculate "vat_calc"
   - Recalculate "gross_calc"

   **If they change SLA%:**
   - Recalculate SLA deduction
   - Recalculate "vat_calc"
   - Recalculate "gross_calc"

**YOUR TASK:**
Determine their intent and intelligently parse field updates with cascading changes.

**ACTIONS:**
- If they want to approve/save/confirm/submit: action = 'approve'
- If they want to cancel/discard/abort: action = 'cancel'
- If they want to see current values: action = 'view'
- If they're making changes/corrections: action = 'edit' and parse ALL field updates (including cascading updates)

**CURRENT BOOKING ORDER DATA:**
{json.dumps(current_data, indent=2)}

**WARNINGS:** {warnings}
**MISSING REQUIRED FIELDS:** {missing_required}

**FIELD MAPPING (use these exact keys):**
- Client/client name/customer → "client"
- Campaign/campaign name/brand → "brand_campaign"
- BO number/booking order number → "bo_number"
- BO date/booking order date → "bo_date"
- Net/net amount/net pre-VAT → "net_pre_vat"
- VAT/vat amount → "vat_calc"
- Gross/gross amount/total → "gross_calc"
- Agency/agency name → "agency"
- Sales person/salesperson → "sales_person"
- SLA percentage → "sla_pct"
- Payment terms → "payment_terms"
- Commission percentage → "commission_pct"
- Notes → "notes"
- Category → "category"
- Asset → "asset" (list of location codes)
- Locations → "locations" (array of location objects with name, dates, costs, etc.)

**FOR LOCATION UPDATES:**
- "Add location X with rental Y" → Append to "locations" AND "asset", recalculate totals
- "Remove location Y" → Remove from "locations" AND "asset", recalculate totals
- "Change location 1 start date to X" → Update by index in locations array

**IMPORTANT:**
- When editing, include ALL fields that need to change (direct changes + cascading updates)
- Be intelligent about understanding what needs to update together
- Calculate new values for financial fields when underlying data changes
- Use natural language in your message. Be friendly and conversational.

Return JSON with: action, fields (ALL changed fields including cascading updates), message (explain what you updated).
"""

        res = await config.openai_client.responses.create(
            model=config.OPENAI_MODEL,
            input=[{"role": "system", "content": system_prompt}],
            text={
                'format': {
                    'type': 'json_schema',
                    'name': 'booking_order_edit_response',
                    'strict': False,
                    'schema': {
                        'type': 'object',
                        'properties': {
                            'action': {'type': 'string', 'enum': ['approve', 'cancel', 'edit', 'view']},
                            'fields': {
                                'type': 'object',
                                'properties': {
                                    'client': {'type': 'string'},
                                    'brand_campaign': {'type': 'string'},
                                    'bo_number': {'type': 'string'},
                                    'bo_date': {'type': 'string'},
                                    'net_pre_vat': {'type': 'number'},
                                    'vat_calc': {'type': 'number'},
                                    'gross_calc': {'type': 'number'},
                                    'agency': {'type': 'string'},
                                    'sales_person': {'type': 'string'},
                                    'sla_pct': {'type': 'number'},
                                    'payment_terms': {'type': 'string'},
                                    'commission_pct': {'type': 'number'},
                                    'notes': {'type': 'string'},
                                    'category': {'type': 'string'},
                                    'asset': {'type': 'string'}
                                },
                                'additionalProperties': True  # Allow locations and other fields
                            },
                            'message': {'type': 'string'}
                        },
                        'required': ['action'],
                        'additionalProperties': False
                    }
                }
            },
            store=False
        )

        decision = json.loads(res.output[0].content[-1].text)
        action = decision.get('action')
        message = decision.get('message', '')

        if action == 'approve':
            # Start approval workflow - send directly to Sales Coordinator (admin is HoS)
            try:
                import bo_approval_workflow

                # Generate combined PDF for coordinator review
                parser = BookingOrderParser(company=edit_data.get("company"))
                combined_pdf = await parser.generate_combined_pdf(
                    current_data,
                    f"DRAFT_{edit_data.get('company')}",
                    Path(edit_data.get("original_file_path"))
                )

                # Create approval workflow (start at coordinator stage since admin is HoS)
                workflow_id = await bo_approval_workflow.create_approval_workflow(
                    user_id=user_id,
                    company=edit_data.get("company"),
                    data=current_data,
                    warnings=warnings,
                    missing_required=missing_required,
                    original_file_path=edit_data.get("original_file_path"),
                    original_filename=edit_data.get("original_filename"),
                    file_type=edit_data.get("file_type"),
                    user_notes=edit_data.get("user_notes", "")
                )

                # Get coordinator channel (uses conversations.open to get DM channel ID)
                coordinator_channel = await bo_approval_workflow.get_coordinator_channel(edit_data.get("company"))
                if not coordinator_channel:
                    return f"❌ **Error:** Sales Coordinator for {edit_data.get('company')} not configured. Please update hos_config.json"

                # Send to Sales Coordinator with buttons
                result = await bo_slack_messaging.send_to_coordinator(
                    channel=coordinator_channel,
                    workflow_id=workflow_id,
                    company=edit_data.get("company"),
                    data=current_data,
                    combined_pdf_path=str(combined_pdf)
                )

                # Update workflow with coordinator message info
                await bo_approval_workflow.update_workflow(workflow_id, {
                    "coordinator_msg_ts": result["message_id"]
                })

                # Clean up edit session
                del pending_booking_orders[user_id]

                return f"✅ **Booking Order Submitted for Approval**\n\n**Client:** {current_data.get('client', 'N/A')}\n**Campaign:** {current_data.get('brand_campaign', 'N/A')}\n**Gross Total:** AED {current_data.get('gross_calc', 0):,.2f}\n\nSent to Sales Coordinator for approval. You'll be notified once the approval process is complete."

            except Exception as e:
                config.logger.error(f"[BOOKING ORDER] Error starting approval workflow: {e}")
                return f"❌ **Error starting approval workflow:** {str(e)}\n\nPlease try again or say 'cancel' to discard."

        elif action == 'cancel':
            # Clean up temp file and session
            try:
                original_file_path = edit_data.get("original_file_path")
                if original_file_path and Path(original_file_path).exists():
                    Path(original_file_path).unlink()
            except Exception as e:
                config.logger.error(f"[BOOKING ORDER] Error deleting temp file: {e}")

            del pending_booking_orders[user_id]
            return message or "❌ **Booking order draft discarded.**"

        elif action == 'view':
            # Show current draft
            preview = "📋 **Current Booking Order Draft**\n\n"

            # Core fields
            preview += f"**Client:** {current_data.get('client', 'N/A')}\n"
            preview += f"**Campaign:** {current_data.get('brand_campaign', 'N/A')}\n"
            preview += f"**BO Number:** {current_data.get('bo_number', 'N/A')}\n"
            preview += f"**BO Date:** {current_data.get('bo_date', 'N/A')}\n"
            preview += f"**Net (pre-VAT):** AED {current_data.get('net_pre_vat', 0):,.2f}\n"
            preview += f"**VAT (5%):** AED {current_data.get('vat_calc', 0):,.2f}\n"
            preview += f"**Gross Total:** AED {current_data.get('gross_calc', 0):,.2f}\n\n"

            # Locations
            locations = current_data.get('locations', [])
            if locations:
                preview += f"**Locations ({len(locations)}):**\n"
                for i, loc in enumerate(locations, 1):
                    preview += f"{i}. {loc.get('name', 'Unknown')}: {loc.get('start_date', '?')} to {loc.get('end_date', '?')} (AED {loc.get('net_amount', 0):,.2f})\n"

            if warnings:
                preview += f"\n⚠️ **Warnings ({len(warnings)}):**\n"
                for w in warnings[:3]:
                    preview += f"• {w}\n"

            if missing_required:
                preview += f"\n❗ **Missing Required Fields:** {', '.join(missing_required)}\n"

            preview += "\n**What would you like to do?**\n"
            preview += "• Tell me any corrections\n"
            preview += "• Say 'approve' to save\n"
            preview += "• Say 'cancel' to discard"

            return preview

        elif action == 'edit':
            # Apply field updates
            fields = decision.get('fields', {})
            if fields:
                # Update the draft data
                for field, value in fields.items():
                    current_data[field] = value

                # Recalculate VAT and gross if net changed
                if 'net_pre_vat' in fields:
                    current_data['vat_calc'] = round(current_data['net_pre_vat'] * 0.05, 2)
                    current_data['gross_calc'] = round(current_data['net_pre_vat'] + current_data['vat_calc'], 2)

                # Save updated draft
                pending_booking_orders[user_id]["data"] = current_data

                response = message or "✅ **Changes applied:**\n"
                for field, value in fields.items():
                    response += f"• {field}: {value}\n"
                response += "\nSay 'approve' to save or continue editing."
                return response
            else:
                return message or "I didn't catch any changes. What would you like to update?"

        else:
            return "I didn't understand. Please tell me what to change, or say 'approve' to save or 'cancel' to discard."

    except Exception as e:
        config.logger.error(f"[BOOKING ORDER] Error in edit flow: {e}")
        return f"❌ **Error processing your request:** {str(e)}\n\nPlease try again."


async def main_llm_loop(channel: str, user_id: str, user_input: str, slack_event: Dict[str, Any] = None):
    logger = config.logger
    
    # Debug logging
    logger.info(f"[MAIN_LLM] Starting for user {user_id}, pending_adds: {list(pending_location_additions.keys())}")
    if slack_event:
        logger.info(f"[MAIN_LLM] Slack event keys: {list(slack_event.keys())}")
        if "files" in slack_event:
            logger.info(f"[MAIN_LLM] Files found: {len(slack_event['files'])}")
    
    # Check if message is in a coordinator thread (before sending status message)
    thread_ts = slack_event.get("thread_ts") if slack_event else None
    if thread_ts:
        # Check if this thread is an active coordinator thread
        import bo_approval_workflow

        # Find workflow with matching coordinator thread
        logger.info(f"[BO APPROVAL] Checking {len(bo_approval_workflow.approval_workflows)} workflows for thread {thread_ts}")
        for workflow_id, workflow in bo_approval_workflow.approval_workflows.items():
            coordinator_thread = workflow.get("coordinator_thread_ts")
            logger.info(f"[BO APPROVAL] Workflow {workflow_id}: coordinator_thread={coordinator_thread}, status={workflow.get('status')}")

            # Check if message is in a coordinator thread (even if not active for editing yet)
            if coordinator_thread == thread_ts:
                # Check if thread is active for editing (after rejection)
                if bo_approval_workflow.is_coordinator_thread_active(workflow, thread_ts):
                    logger.info(f"[BO APPROVAL] Message in active coordinator thread for {workflow_id}")
                    try:
                        answer = await bo_approval_workflow.handle_coordinator_thread_message(
                            workflow_id=workflow_id,
                            user_id=user_id,
                            user_input=user_input,
                            channel=channel,
                            thread_ts=thread_ts
                        )
                        # Only send message if there's an answer (execute action returns None)
                        if answer is not None:
                            await config.slack_client.chat_postMessage(
                                channel=channel,
                                thread_ts=thread_ts,
                                text=config.markdown_to_slack(answer)
                            )
                    except Exception as e:
                        logger.error(f"[BO APPROVAL] Error in coordinator thread handler: {e}", exc_info=True)
                        await config.slack_client.chat_postMessage(
                            channel=channel,
                            thread_ts=thread_ts,
                            text=config.markdown_to_slack(f"❌ **Error:** {str(e)}")
                        )
                    return  # Exit early, don't process as normal message
                else:
                    # Thread exists but not active yet (user hasn't rejected)
                    logger.info(f"[BO APPROVAL] Message in coordinator thread for {workflow_id} but thread not active - reminding user to use buttons")
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        thread_ts=thread_ts,
                        text=config.markdown_to_slack(
                            "⚠️ **Please use the Approve or Reject buttons first**\n\n"
                            "To make edits to this booking order, click the **❌ Reject** button above. "
                            "This will open the thread for editing.\n\n"
                            "If the BO looks good, click **✅ Approve** to send it to the next stage."
                        )
                    )
                    return  # Exit early, don't process as normal message

    # Send initial status message
    status_message = await config.slack_client.chat_postMessage(
        channel=channel,
        text="⏳ _Please wait..._"
    )
    status_ts = status_message.get("ts")

    # OLD EDIT FLOW REMOVED - Now coordinators edit directly in threads

    # Check if user has a pending location addition or mockup request and uploaded a file
    # Also check for file_share events which Slack sometimes uses
    has_files = slack_event and ("files" in slack_event or (slack_event.get("subtype") == "file_share"))

    # Note: Mockup generation is now handled in one step within the tool handler
    # No need for pending state - users must upload image WITH request or provide AI prompt

    if user_id in pending_location_additions and has_files:
        pending_data = pending_location_additions[user_id]

        # Check if pending request is still valid (10 minute window)
        timestamp = pending_data.get("timestamp")
        if timestamp and (datetime.now() - timestamp) > timedelta(minutes=10):
            del pending_location_additions[user_id]
            logger.warning(f"[LOCATION_ADD] Pending location expired for user {user_id}")
            await config.slack_client.chat_postMessage(
                channel=channel,
                text=config.markdown_to_slack("❌ **Error:** Location upload session expired (10 minute limit). Please restart the location addition process.")
            )
            return

        logger.info(f"[LOCATION_ADD] Found pending location for user {user_id}: {pending_data['location_key']}")
        logger.info(f"[LOCATION_ADD] Files in event: {len(slack_event.get('files', []))}")
        
        # Check if any of the files is a PDF (we'll convert it to PPTX)
        pptx_file = None
        files = slack_event.get("files", [])

        # If it's a file_share event, files might be structured differently
        if not files and slack_event.get("subtype") == "file_share" and "file" in slack_event:
            files = [slack_event["file"]]
            logger.info(f"[LOCATION_ADD] Using file from file_share event")

        for f in files:
            logger.info(f"[LOCATION_ADD] Checking file: name={f.get('name')}, filetype={f.get('filetype')}, mimetype={f.get('mimetype')}")

            # Accept PDF files (new) - will be converted to PPTX
            if f.get("filetype") == "pdf" or f.get("mimetype") == "application/pdf" or f.get("name", "").lower().endswith(".pdf"):
                try:
                    pdf_file = await _download_slack_file(f)
                except Exception as e:
                    logger.error(f"Failed to download PDF file: {e}")
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack("❌ **Error:** Failed to download the PDF file. Please try again.")
                    )
                    return

                # Validate it's actually a PDF file
                if not _validate_pdf_file(pdf_file):
                    logger.error(f"Invalid PDF file: {f.get('name')}")
                    try:
                        os.unlink(pdf_file)
                    except:
                        pass
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack("❌ **Error:** The uploaded file is not a valid PDF. Please upload a .pdf file.")
                    )
                    return

                # Post status message about conversion
                conversion_status = await config.slack_client.chat_postMessage(
                    channel=channel,
                    text="⏳ _Converting PDF to PowerPoint with maximum quality (300 DPI)..._"
                )

                # Convert PDF to PPTX
                logger.info(f"[LOCATION_ADD] Converting PDF to PPTX...")
                pptx_file = await _convert_pdf_to_pptx(pdf_file)

                # Clean up original PDF
                try:
                    os.unlink(pdf_file)
                except:
                    pass

                # Delete conversion status message
                await config.slack_client.chat_delete(channel=channel, ts=conversion_status["ts"])

                if not pptx_file:
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack("❌ **Error:** Failed to convert PDF to PowerPoint. Please try again or contact support.")
                    )
                    return

                logger.info(f"[LOCATION_ADD] ✓ PDF converted to PPTX: {pptx_file}")
                break
        
        if pptx_file:
            # Build metadata.txt content matching exact format of existing files
            metadata_lines = []
            metadata_lines.append(f"Location Name: {pending_data['display_name']}")
            metadata_lines.append(f"Display Name: {pending_data['display_name']}")
            metadata_lines.append(f"Display Type: {pending_data['display_type']}")
            metadata_lines.append(f"Number of Faces: {pending_data['number_of_faces']}")
            
            # For digital locations, add digital-specific fields in the correct order
            if pending_data['display_type'] == 'Digital':
                metadata_lines.append(f"Spot Duration: {pending_data['spot_duration']}")
                metadata_lines.append(f"Loop Duration: {pending_data['loop_duration']}")
                metadata_lines.append(f"SOV: {pending_data['sov']}")
                if pending_data['upload_fee'] is not None:
                    metadata_lines.append(f"Upload Fee: {pending_data['upload_fee']}")
            
            # Series, Height, Width come after digital fields
            metadata_lines.append(f"Series: {pending_data['series']}")
            metadata_lines.append(f"Height: {pending_data['height']}")
            metadata_lines.append(f"Width: {pending_data['width']}")
            
            metadata_text = "\n".join(metadata_lines)
            
            try:
                # Save the location
                await _persist_location_upload(pending_data['location_key'], pptx_file, metadata_text)
                
                # Clean up
                del pending_location_additions[user_id]
                
                # Refresh templates
                config.refresh_templates()
                
                # Delete status message
                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                
                await config.slack_client.chat_postMessage(
                    channel=channel,
                    text=config.markdown_to_slack(
                        f"✅ **Successfully added location `{pending_data['location_key']}`**\n\n"
                        f"The location is now available for use in proposals."
                    )
                )
                return
            except Exception as e:
                logger.error(f"Failed to save location: {e}")
                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                await config.slack_client.chat_postMessage(
                    channel=channel,
                    text=config.markdown_to_slack("❌ **Error:** Failed to save the location. Please try again.")
                )
                # Clean up the temporary file
                try:
                    os.unlink(pptx_file)
                except:
                    pass
                return
        else:
            # No PPT file found, cancel the addition
            del pending_location_additions[user_id]
            await config.slack_client.chat_delete(channel=channel, ts=status_ts)
            await config.slack_client.chat_postMessage(
                channel=channel,
                text=config.markdown_to_slack(
                    "❌ **Location addition cancelled.**\n\n"
                    "No PowerPoint file was found in your message. Please start over with 'add location' if you want to try again."
                )
            )
            return

    # Check for location deletion confirmation
    if user_input.strip().lower().startswith("confirm delete ") and config.is_admin(user_id):
        location_key = user_input.strip().lower().replace("confirm delete ", "").strip()

        if location_key in config.LOCATION_METADATA:
            location_dir = config.TEMPLATES_DIR / location_key
            display_name = config.LOCATION_METADATA[location_key].get('display_name', location_key)

            try:
                # Delete the location directory and all its contents
                import shutil
                import mockup_generator

                # Delete PowerPoint templates
                if location_dir.exists():
                    shutil.rmtree(location_dir)
                    logger.info(f"[LOCATION_DELETE] Deleted location directory: {location_dir}")

                # Delete all mockup photos and database entries for this location
                mockup_dir = mockup_generator.MOCKUPS_DIR / location_key
                if mockup_dir.exists():
                    shutil.rmtree(mockup_dir)
                    logger.info(f"[LOCATION_DELETE] Deleted mockup directory: {mockup_dir}")

                # Delete all mockup frame data from database
                import db
                conn = db._connect()
                try:
                    result = conn.execute("DELETE FROM mockup_frames WHERE location_key = ?", (location_key,))
                    deleted_count = result.rowcount
                    conn.commit()
                    logger.info(f"[LOCATION_DELETE] Deleted {deleted_count} mockup frame entries from database")
                finally:
                    conn.close()

                # Refresh templates to remove from cache
                config.refresh_templates()

                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                await config.slack_client.chat_postMessage(
                    channel=channel,
                    text=config.markdown_to_slack(
                        f"✅ **Location `{location_key}` successfully deleted**\n\n"
                        f"📍 **Removed:** {display_name}\n"
                        f"🗑️ **Files deleted:** PowerPoint template, metadata, and {deleted_count} mockup frames\n"
                        f"🔄 **Templates refreshed:** Location no longer available for proposals"
                    )
                )
                return
            except Exception as e:
                logger.error(f"[LOCATION_DELETE] Failed to delete location {location_key}: {e}")
                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                await config.slack_client.chat_postMessage(
                    channel=channel,
                    text=config.markdown_to_slack(f"❌ **Error:** Failed to delete location `{location_key}`. Please try again or check server logs.")
                )
                return
        else:
            await config.slack_client.chat_delete(channel=channel, ts=status_ts)
            await config.slack_client.chat_postMessage(
                channel=channel,
                text=config.markdown_to_slack(f"❌ **Error:** Location `{location_key}` not found. Deletion cancelled.")
            )
            return

    # Handle cancellation
    if user_input.strip().lower() == "cancel" and config.is_admin(user_id):
        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
        await config.slack_client.chat_postMessage(
            channel=channel,
            text=config.markdown_to_slack("✅ **Operation cancelled.**")
        )
        return

    # Clean up old pending additions (older than 10 minutes)
    cutoff = datetime.now() - timedelta(minutes=10)
    expired_users = [
        uid for uid, data in pending_location_additions.items()
        if data.get("timestamp", datetime.now()) < cutoff
    ]
    for uid in expired_users:
        del pending_location_additions[uid]

    available_names = ", ".join(config.available_location_names())
    
    # Get static and digital locations for the prompt
    static_locations = []
    digital_locations = []
    for key, meta in config.LOCATION_METADATA.items():
        display_name = meta.get('display_name', key)
        if meta.get('display_type', '').lower() == 'static':
            static_locations.append(f"{display_name} ({key})")
        elif meta.get('display_type', '').lower() == 'digital':
            digital_locations.append(f"{display_name} ({key})")

    static_list = ", ".join(static_locations) if static_locations else "None"
    digital_list = ", ".join(digital_locations) if digital_locations else "None"

    # Check if user is admin for system prompt and tool filtering
    is_admin = config.is_admin(user_id)

    prompt = (
        f"You are an AI sales assistant for BackLite Media. You provide comprehensive sales support tools including:\n"
        f"• Financial proposal generation for advertising locations\n"
        f"• Billboard mockup visualization (upload-based or AI-generated)\n"
        f"• Booking order parsing and management\n"
        f"• Location database management\n"
        f"• Sales analytics and reporting\n"
        f"• Code interpreter for calculations and data analysis\n\n"
        f"CRITICAL INSTRUCTION:\n"
        f"You are an INTERFACE to tools, NOT the executor. When users request actions:\n"
        f"- DO NOT say 'Generating now...', 'Creating proposal...', or 'Building mockup...'\n"
        f"- DO call the appropriate tool/function immediately\n"
        f"- Let the TOOL handle the actual execution\n"
        f"- Only respond with text AFTER the tool completes or if asking clarifying questions\n\n"
        f"⚠️ CONTEXT SWITCHING - CRITICAL:\n"
        f"Users frequently switch between different task types. ALWAYS analyze the CURRENT message to determine what they want NOW:\n\n"
        f"🔴 IMMEDIATE CONTEXT RESET RULES:\n"
        f"1. If user mentions NEW location names → FORGET previous proposal, generate NEW proposal with NEW locations\n"
        f"   Example: Just made 'gateway' proposal, now user says 'jawhara' → Generate JAWHARA proposal (not gateway again)\n"
        f"2. If user mentions NEW client name → FORGET previous client, use NEW client name\n"
        f"3. If user mentions NEW dates/rates → FORGET previous values, use NEW values\n"
        f"4. If user uploads a PDF after proposal request → They want BOOKING ORDER parsing, NOT proposal\n"
        f"5. If user uploads image after any request → They want MOCKUP generation, NOT previous task\n"
        f"6. Each message is a FRESH request - extract ALL parameters from CURRENT message only\n\n"
        f"🟢 HOW TO DETECT NEW REQUESTS:\n"
        f"- Look for location names in CURRENT message (gateway, jawhara, landmark, etc.)\n"
        f"- Look for proposal keywords (make, create, generate, proposal)\n"
        f"- Look for dates, rates, durations in CURRENT message\n"
        f"- If CURRENT message has complete info → Call tool immediately with NEW data\n"
        f"- If CURRENT message missing info → Ask for missing info ONLY (don't repeat previous request)\n\n"
        f"🔴 NEVER DO THIS:\n"
        f"- Don't say 'generating the same proposal as before'\n"
        f"- Don't ask 'did you mean the previous location?'\n"
        f"- Don't use location names from previous messages if CURRENT message has different locations\n"
        f"- Don't require user to explicitly say 'new proposal' or 'different location'\n\n"
        f"✅ ALWAYS DO THIS:\n"
        f"- Parse CURRENT message for ALL parameters (locations, dates, rates, client)\n"
        f"- If CURRENT message has different locations than last → Use CURRENT locations\n"
        f"- Trust CURRENT message over conversation history\n"
        f"- Call tool with parameters from CURRENT message ONLY\n\n"
        f"Today's date is: {datetime.now().strftime('%B %d, %Y')} ({datetime.now().strftime('%A')})\n"
        f"Use this date to understand relative dates like 'tomorrow', 'next week', 'next month', etc.\n\n"
        f"═══════════════════════════════════════════════════════════════════\n"
        f"📊 PROPOSAL GENERATION\n"
        f"═══════════════════════════════════════════════════════════════════\n\n"
        f"You can handle SINGLE or MULTIPLE location proposals in one request.\n\n"
        f"PACKAGE TYPES:\n"
        f"1. SEPARATE PACKAGE (default): Each location gets its own proposal slide, multiple durations/rates allowed per location\n"
        f"2. COMBINED PACKAGE: All locations in ONE proposal slide, single duration per location, one combined net rate\n\n"
        
        f"LOCATION TYPES - CRITICAL TO UNDERSTAND:\n\n"
        f"🔴 DIGITAL LOCATIONS (LED screens with rotating ads):\n"
        f"   Features: Multiple advertisers share screen time, ads rotate in loops\n"
        f"   Fee Structure: NET RATE + PRE-CONFIGURED UPLOAD FEE (automatically added)\n"
        f"   Examples: {digital_list}\n"
        f"   Upload Fee: System automatically adds the correct upload fee for each digital location\n\n"

        f"🔵 STATIC LOCATIONS (Traditional billboards, prints, physical displays):\n"
        f"   Features: Single advertiser has exclusive display, no rotation\n"
        f"   Fee Structure: NET RATE + PRODUCTION FEE (must be collected from user)\n"
        f"   Examples: {static_list}\n"
        f"   Production Fee: REQUIRED - ask user for production fee amount (e.g., 'AED 5,000')\n\n"

        f"CRITICAL RULES:\n"
        f"- DIGITAL = Upload fee (automatic) | STATIC = Production fee (ask user)\n"
        f"- NEVER ask for production fee on digital locations\n"
        f"- NEVER skip production fee on static locations\n"
        f"- If user mentions 'upload fee' for static locations, correct them to 'production fee'\n\n"
        
        f"REQUIRED INFORMATION:\n"
        f"For SEPARATE PACKAGE (each location):\n"
        f"1. Location (must match from lists above - intelligently infer if user says 'gateway'→'dubai_gateway', 'jawhara'→'dubai_jawhara', 'the landmark'→'landmark', etc.)\n"
        f"2. Start Date\n"
        f"3. Duration Options (multiple allowed)\n"
        f"4. Net Rates for EACH duration\n"
        f"5. Fees - CHECK LOCATION TYPE:\n"
        f"   • DIGITAL locations: NO FEE NEEDED (upload fee auto-added)\n"
        f"   • STATIC locations: ASK for production fee (e.g., 'AED 5,000')\n"
        f"6. Client Name (required)\n"
        f"7. Submitted By (optional - defaults to current user)\n\n"
        f"For COMBINED PACKAGE:\n"
        f"1. All Locations (mix of digital/static allowed - intelligently infer names from available list)\n"
        f"2. Start Date for EACH location\n"
        f"3. ONE Duration per location\n"
        f"4. ONE Combined Net Rate for entire package\n"
        f"5. Fees - CHECK EACH LOCATION TYPE:\n"
        f"   • DIGITAL locations: NO FEE NEEDED (upload fees auto-added)\n"
        f"   • STATIC locations: ASK for production fee for EACH static location\n"
        f"6. Client Name (required)\n"
        f"7. Submitted By (optional - defaults to current user)\n\n"
        
        f"MULTIPLE PROPOSALS RULES:\n"
        f"- User can request proposals for multiple locations at once\n"
        f"- EACH location must have its own complete set of information\n"
        f"- EACH location must have matching number of durations and net rates\n"
        f"- Different locations can have different durations/rates\n"
        f"- Multiple proposals will be combined into a single PDF document\n\n"
        
        f"VALIDATION RULES:\n"
        f"- For EACH location, durations count MUST equal net rates count\n"
        f"- If a location has 3 duration options, it MUST have exactly 3 net rates\n"
        f"- DO NOT proceed until ALL locations have complete information\n"
        f"- Ask follow-up questions for any missing information\n"
        f"- ALWAYS ask for client name if not provided\n\n"
        
        f"PARSING EXAMPLES:\n"
        f"User: 'jawhara, oryx and triple crown special combined deal 2 mil, 2, 4 and 6 weeks respectively, 1st jan 2026, 2nd jan 2026 and 3rd'\n"
        f"Parse as: Combined package with Jawhara (2 weeks, Jan 1), Oryx (4 weeks, Jan 2), Triple Crown (6 weeks, Jan 3), total 2 million AED\n\n"
        
        f"SINGLE LOCATION EXAMPLE:\n"
        f"User: 'Proposal for landmark, Jan 1st, 2 weeks at 1.5M'\n"
        f"Bot confirms and generates one proposal\n\n"
        
        f"MULTIPLE LOCATIONS EXAMPLE:\n"
        f"User: 'I need proposals for landmark and gateway'\n"
        f"Bot: 'I'll help you create proposals for The Landmark and The Gateway. Let me get the details for each:\n\n"
        f"For THE LANDMARK:\n"
        f"- What's the campaign start date?\n"
        f"- What duration options do you want?\n"
        f"- What are the net rates for each duration?\n\n"
        f"For THE GATEWAY:\n"
        f"- What's the campaign start date?\n"
        f"- What duration options do you want?\n"
        f"- What are the net rates for each duration?'\n\n"
        
        f"COMBINED PACKAGE EXAMPLE:\n"
        f"User: 'I need a combined package for landmark, gateway, and oryx at 5 million total'\n"
        f"Bot: 'I'll create a combined package proposal. Let me confirm the details:\n\n"
        f"COMBINED PACKAGE:\n"
        f"- Locations: The Landmark, The Gateway, The Oryx\n"
        f"- Package Net Rate: AED 5,000,000\n\n"
        f"For each location, I need:\n"
        f"- Start date\n"
        f"- Duration (one per location for combined packages)\n\n"
        f"Please provide these details.'\n\n"

        f"═══════════════════════════════════════════════════════════════════\n"
        f"🎨 MOCKUP GENERATION\n"
        f"═══════════════════════════════════════════════════════════════════\n\n"
        f"MOCKUP SETUP WEBSITE: {os.getenv('RENDER_EXTERNAL_URL', 'http://localhost:3000')}/mockup\n"
        f"(Share this URL when users ask about setting up mockup frames or uploading billboard photos)\n\n"
        f"You can GENERATE MOCKUPS: Create billboard mockups with uploaded or AI-generated creatives:\n"
        f"  TWO MODES (everything must be in ONE message):\n"
        f"  A) USER UPLOAD MODE (requires image attachment):\n"
        f"     1. User UPLOADS image(s) WITH mockup request in same message\n"
        f"     2. System detects images and generates mockup immediately\n"
        f"     3. INTELLIGENT TEMPLATE SELECTION:\n"
        f"        • 1 image uploaded → Selects ANY template (image duplicated across all frames)\n"
        f"        • N images uploaded → ONLY selects templates with EXACTLY N frames\n"
        f"        • Example: Upload 3 images → system finds template with 3 frames only\n"
        f"     CRITICAL: If you see '[User uploaded X image file(s): ...]' in the message, call generate_mockup IMMEDIATELY\n"
        f"     DO NOT ask for clarification - the images are already uploaded!\n"
        f"  B) AI GENERATION MODE (NO upload needed):\n"
        f"     1. User provides location AND creative description in request\n"
        f"     2. System generates creative using gpt-image-1 model (NO upload needed)\n"
        f"     3. MULTI-FRAME AI SUPPORT:\n"
        f"        • Default: Generates 1 artwork (duplicated across frames if multi-frame template)\n"
        f"        • User can request multiple variations: 'dual frame mockup', 'triple frame with 3 different ads'\n"
        f"        • Set num_ai_frames parameter to generate N DIFFERENT artworks for N frames\n"
        f"        • System uses prompt parser to create variations, then matches to N-frame template\n"
        f"        • Example: 'triple crown with 3 different nike ads' → num_ai_frames=3\n"
        f"     4. System applies AI creative(s) to billboard and returns mockup\n"
        f"     IMPORTANT: If description provided = AI mode, ignore any uploaded images\n"
        f"  Decision Logic:\n"
        f"  - Has creative description? → Use AI mode (ignore uploads)\n"
        f"  - No description but has upload? → Use upload mode (DO NOT ASK FOR CLARIFICATION)\n"
        f"  - No description and no upload? → ERROR\n"
        f"  Examples:\n"
        f"  - [uploads creative.jpg] + 'mockup for Dubai Gateway' → uses uploaded image (IMMEDIATE)\n"
        f"  - [uploads 3 images] + 'triple crown mockup' → matches to 3-frame template (IMMEDIATE)\n"
        f"  - 'put this on triple crown' + [User uploaded 1 image file(s): test.jpg] → IMMEDIATE mockup\n"
        f"  - 'mockup for Oryx with luxury watch ad, gold and elegant' → AI generates 1 creative\n"
        f"  - 'triple crown with 3 different nike shoe ads' → AI generates 3 variations (num_ai_frames=3)\n"
        f"  - 'mockup for Gateway' (no upload, no description) → ERROR: missing creative\n"
        f"  Keywords: 'mockup', 'mock up', 'billboard preview', 'show my ad on', 'put this on', 'dual frame', 'triple frame'\n\n"

        f"═══════════════════════════════════════════════════════════════════\n"
        f"🗄️ DATABASE & LOCATION MANAGEMENT\n"
        f"═══════════════════════════════════════════════════════════════════\n\n"
        f"- ADD NEW LOCATIONS (admin only):\n"
        f"  • Admin provides ALL metadata: location_key, display_name, display_type, height, width, number_of_faces, sov, series, spot_duration, loop_duration, upload_fee (for digital)\n"
        f"  • Once validated, admin uploads the PPT template file\n"
        f"  • Location becomes immediately available for proposals\n\n"
        f"- DELETE LOCATIONS (admin only): Requires double confirmation to prevent accidents\n"
        f"- REFRESH TEMPLATES: Reload available locations from disk\n"
        f"- LIST LOCATIONS: Show all available advertising locations\n\n"

        f"═══════════════════════════════════════════════════════════════════\n"
        f"📈 ANALYTICS & REPORTING\n"
        f"═══════════════════════════════════════════════════════════════════\n\n"
        f"- EXPORT DATABASE: Export all proposals to Excel (admin only - triggered by 'excel backend' or similar)\n"
        f"- GET STATISTICS: View proposal generation summary and recent activity\n"
        f"- EDIT TASKS: Modify task management workflows\n\n"

        f"═══════════════════════════════════════════════════════════════════\n"
        f"👤 USER PERMISSIONS\n"
        f"═══════════════════════════════════════════════════════════════════\n\n"
        f"Current User: {'ADMIN' if is_admin else 'STANDARD USER'}\n\n"
        f"{'✅ ADMIN TOOLS AVAILABLE:' if is_admin else '❌ ADMIN TOOLS NOT AVAILABLE:'}\n"
        f"- Location Management (add_location, delete_location)\n"
        f"- Database Export (export_proposals_to_excel, export_booking_orders_to_excel)\n"
        f"- Fetch Booking Orders (fetch_booking_order)\n\n"
        f"✅ AVAILABLE TO ALL USERS:\n"
        f"- Booking Order Upload & Parsing (parse_booking_order)\n"
        f"- Any sales person can upload booking orders for approval\n\n"
        f"{'You have access to all admin-only tools listed above.' if is_admin else 'You do NOT have access to admin-only tools like location management or database export, but you CAN upload and parse booking orders.'}\n\n"

        f"═══════════════════════════════════════════════════════════════════\n"
        f"⚙️ SYSTEM GUIDELINES\n"
        f"═══════════════════════════════════════════════════════════════════\n\n"
        f"IMPORTANT:\n"
        f"- Use get_separate_proposals for individual location proposals with multiple duration/rate options\n"
        f"- Use get_combined_proposal for special package deals with one total price\n"
        f"- For SEPARATE packages: each location gets its own proposal slide\n"
        f"- For COMBINED packages: all locations in ONE proposal slide with ONE net rate\n"
        f"- Single location always uses get_separate_proposals\n"
        f"- When user mentions 'combined deal' or 'special package' with total price, use get_combined_proposal\n"
        f"- Format all rates as 'AED X,XXX,XXX'\n"
        f"- Parse 'mil' or 'million' as 000,000 (e.g., '2 mil' = 'AED 2,000,000')\n"
        f"- Number of spots defaults to 1 if not specified\n"
        f"FEE COLLECTION RULES (CRITICAL):\n"
        f"- DIGITAL locations: NEVER ask for fees - upload fees are automatic\n"
        f"- STATIC locations: ALWAYS ask for production fee - it's mandatory\n"
        f"- Mixed packages: Ask production fees only for static locations\n"
        f"- If confused about location type, check the lists above\n"
        f"- ALWAYS collect client name - it's required for tracking\n\n"

        f"═══════════════════════════════════════════════════════════════════\n"
        f"🎨 BILLBOARD MOCKUP GENERATION\n"
        f"═══════════════════════════════════════════════════════════════════\n\n"

        f"MOCKUP MEMORY SYSTEM (30-Minute Creative Storage):\n"
        f"When a user generates a mockup, the system stores their creative files (NOT the final mockup) for 30 minutes.\n"
        f"This enables FOLLOW-UP REQUESTS where users can apply the same creatives to different locations.\n\n"

        f"FOLLOW-UP REQUEST DETECTION:\n"
        f"If a user recently generated a mockup (within 30 min) and asks to see it on another location WITHOUT uploading new images or providing AI prompt:\n"
        f"- Examples: 'show me this on Dubai Gateway', 'apply to The Landmark', 'how would it look at Oryx'\n"
        f"- Just call generate_mockup with the new location name - the system automatically reuses stored creatives\n"
        f"- DO NOT ask them to re-upload images or provide AI prompt again\n"
        f"- The system validates frame count compatibility (3-frame creatives can't be used on 1-frame locations)\n\n"

        f"FRAME COUNT VALIDATION:\n"
        f"- Multi-frame locations (2, 3, or more frames) require matching number of creatives\n"
        f"- If user has 3-frame creatives in memory but requests 1-frame location → system shows error automatically\n"
        f"- If frame mismatch error occurs, explain user needs to upload correct number of images OR use AI generation\n\n"

        f"MOCKUP GENERATION MODES:\n"
        f"1. UPLOAD MODE: User uploads image file(s) → Call generate_mockup IMMEDIATELY, no questions\n"
        f"   - Takes priority over everything else\n"
        f"   - Replaces any stored creatives with new upload\n"
        f"   - DO NOT ask for clarification if user uploads images with location mention\n\n"

        f"2. AI MODE: User provides creative description (no upload) → Call generate_mockup with ai_prompt\n"
        f"   - Example: 'mockup for Dubai Gateway with luxury watch ad, gold and elegant'\n"
        f"   - For multi-frame locations, specify num_ai_frames parameter\n"
        f"   - System generates flat artwork designs (NOT photos of billboards)\n\n"

        f"3. FOLLOW-UP MODE: User requests different location (no upload, no AI, within 30 min)\n"
        f"   - Example: 'show me this on The Landmark'\n"
        f"   - Just call generate_mockup with new location - system handles rest\n"
        f"   - User doesn't need to specify they want to reuse creatives\n\n"

        f"CRITICAL MOCKUP RULES:\n"
        f"- If user uploads images AND mentions location → Call generate_mockup IMMEDIATELY\n"
        f"- Don't ask 'which mockup' or 'which creative' for follow-ups - system knows\n"
        f"- Frame count errors are handled automatically - just relay system message\n"
        f"- After 30 minutes, stored creatives expire - user must upload/generate again"
    )

    # Check if user uploaded files and append to message
    user_message_content = user_input
    image_files = []  # Initialize outside conditional block
    document_files = []  # For PDFs, Excel, etc.

    if has_files and slack_event:
        files = slack_event.get("files", [])
        if not files and slack_event.get("subtype") == "file_share" and "file" in slack_event:
            files = [slack_event["file"]]

        # Check for image files and document files
        for f in files:
            filetype = f.get("filetype", "")
            mimetype = f.get("mimetype", "")
            filename = f.get("name", "").lower()

            # Image files (for mockups)
            if (filetype in ["jpg", "jpeg", "png", "gif", "bmp"] or
                mimetype.startswith("image/") or
                any(filename.endswith(ext) for ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp"])):
                image_files.append(f.get("name", "image"))
            # Document files (for booking orders, proposals, etc.)
            elif (filetype in ["pdf", "xlsx", "xls", "csv", "docx", "doc"] or
                  mimetype in ["application/pdf", "application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"] or
                  any(filename.endswith(ext) for ext in [".pdf", ".xlsx", ".xls", ".csv", ".docx", ".doc"])):
                document_files.append(f.get("name", "document"))

        # PRE-ROUTING CLASSIFIER: Classify and route files before LLM
        if len(files) == 1:
            logger.info(f"[PRE-ROUTER] Single file upload detected, running classification...")

            try:
                file_info = files[0]

                # Download file
                tmp_file = await _download_slack_file(file_info)
                logger.info(f"[PRE-ROUTER] Downloaded: {tmp_file}")

                # Classify using existing classifier (converts to PDF, sends to OpenAI, returns classification)
                from booking_parser import BookingOrderParser
                parser = BookingOrderParser(company="backlite")  # Company will be determined by classifier
                classification = await parser.classify_document(tmp_file, user_message=user_input)

                logger.info(f"[PRE-ROUTER] Classification: {classification}")

                # Route based on HIGH confidence only
                if classification.get("classification") == "BOOKING_ORDER" and classification.get("confidence") == "high":
                    company = classification.get("company", "backlite")  # Get company from classifier
                    logger.info(f"[PRE-ROUTER] HIGH CONFIDENCE BOOKING ORDER ({company}) - routing directly")

                    # Route to booking order parser
                    await _handle_booking_order_parse(
                        company=company,
                        slack_event=slack_event,
                        channel=channel,
                        status_ts=status_ts,
                        user_notes="",
                        user_id=user_id,
                        user_message=user_input
                    )
                    return  # Exit early - don't call LLM

                elif classification.get("classification") == "ARTWORK" and classification.get("confidence") == "high":
                    logger.info(f"[PRE-ROUTER] HIGH CONFIDENCE ARTWORK - letting LLM handle mockup")
                    tmp_file.unlink(missing_ok=True)
                    # Clear document_files and set as image for LLM to handle as mockup
                    document_files.clear()
                    if not image_files:  # If not already marked as image
                        image_files.append(file_info.get("name", "artwork"))
                    # Fall through to LLM for mockup generation

                else:
                    logger.info(f"[PRE-ROUTER] Low/medium confidence - letting LLM decide")
                    tmp_file.unlink(missing_ok=True)
                    # Fall through to LLM

            except Exception as e:
                logger.error(f"[PRE-ROUTER] Classification/routing failed: {e}", exc_info=True)
                # Fall through to LLM on error

        # Inform LLM about uploaded files (only if pre-router didn't handle it)
        if image_files:
            user_message_content = f"{user_input}\n\n[User uploaded {len(image_files)} image file(s): {', '.join(image_files)}]"
            logger.info(f"[LLM] Detected {len(image_files)} uploaded image(s), informing LLM")
        elif document_files:
            user_message_content = f"{user_input}\n\n[User uploaded {len(document_files)} document file(s): {', '.join(document_files)}]"
            logger.info(f"[LLM] Detected {len(document_files)} uploaded document(s), informing LLM")

    # Inject mockup history context ONLY if user did NOT upload ANY files (to avoid confusion)
    # Don't inject mockup history when user uploads documents (BO PDFs, Excel, etc.) or images
    if not image_files and not document_files:
        mockup_hist = get_mockup_history(user_id)
        if mockup_hist:
            metadata = mockup_hist.get("metadata", {})
            stored_location = metadata.get("location_name", "unknown")
            stored_frames = metadata.get("num_frames", 1)
            mode = metadata.get("mode", "unknown")

            # Calculate time remaining
            timestamp = mockup_hist.get("timestamp")
            if timestamp:
                time_remaining = 30 - int((datetime.now() - timestamp).total_seconds() / 60)
                time_remaining = max(0, time_remaining)

                user_message_content = (
                    f"{user_input}\n\n"
                    f"[SYSTEM: User has {stored_frames}-frame creative(s) in memory from '{stored_location}' ({mode}). "
                    f"Expires in {time_remaining} minutes. Can reuse for follow-up mockup requests on locations with {stored_frames} frame(s).]"
                )
                logger.info(f"[LLM] Injected mockup history context: {stored_frames} frames from {stored_location}, {time_remaining}min remaining")

    history = user_history.get(user_id, [])
    history.append({"role": "user", "content": user_message_content, "timestamp": datetime.now().isoformat()})
    history = history[-10:]
    # Remove timestamp from messages sent to OpenAI
    messages_for_openai = [{"role": msg["role"], "content": msg["content"]} for msg in history if "role" in msg and "content" in msg]
    messages = [{"role": "developer", "content": prompt}] + messages_for_openai

    # Base tools available to all users
    tools = [
        {
            "type": "function",
            "name": "get_separate_proposals",
            "description": "Generate SEPARATE proposals - each location gets its own proposal slide with multiple duration/rate options. Use this when user asks to 'make', 'create', or 'generate' proposals for specific locations. Returns individual PPTs and combined PDF.",
            "parameters": {
                "type": "object",
                "properties": {
                    "proposals": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "location": {"type": "string", "description": "The location name - intelligently match to available locations. If user says 'gateway' or 'the gateway', match to 'dubai_gateway'. If user says 'jawhara', match to 'dubai_jawhara'. Use your best judgment to infer the correct location from the available list even if the name is abbreviated or has 'the' prefix."},
                                "start_date": {"type": "string", "description": "Start date for the campaign (e.g., 1st December 2025)"},
                                "durations": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "description": "List of duration options (e.g., ['2 Weeks', '4 Weeks', '6 Weeks'])"
                                },
                                "net_rates": {
                                    "type": "array", 
                                    "items": {"type": "string"},
                                    "description": "List of net rates corresponding to each duration (e.g., ['AED 1,250,000', 'AED 2,300,000', 'AED 3,300,000'])"
                                },
                                "spots": {"type": "integer", "description": "Number of spots (default: 1)", "default": 1},
                                "production_fee": {"type": "string", "description": "Production fee for static locations (e.g., 'AED 5,000'). Required for static locations."}
                            },
                            "required": ["location", "start_date", "durations", "net_rates"]
                        },
                        "description": "Array of proposal objects. Each location can have multiple duration/rate options."
                    },
                    "client_name": {
                        "type": "string",
                        "description": "Name of the client (required)"
                    }
                },
                "required": ["proposals", "client_name"]
            }
        },
        {
            "type": "function",
            "name": "get_combined_proposal",
            "description": "Generate COMBINED package proposal - all locations in ONE slide with single net rate. Use this when user asks for a 'package', 'bundle', or 'combined' deal with multiple locations sharing one total rate.",
            "parameters": {
                "type": "object",
                "properties": {
                    "proposals": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "location": {"type": "string", "description": "The location name - intelligently match to available locations. If user says 'gateway' or 'the gateway', match to 'dubai_gateway'. If user says 'jawhara', match to 'dubai_jawhara'. Use your best judgment to infer the correct location from the available list even if the name is abbreviated or has 'the' prefix."},
                                "start_date": {"type": "string", "description": "Start date for this location (e.g., 1st January 2026)"},
                                "duration": {"type": "string", "description": "Duration for this location (e.g., '2 Weeks')"},
                                "spots": {"type": "integer", "description": "Number of spots (default: 1)", "default": 1},
                                "production_fee": {"type": "string", "description": "Production fee for static locations (e.g., 'AED 5,000'). Required for static locations."}
                            },
                            "required": ["location", "start_date", "duration"]
                        },
                        "description": "Array of locations with their individual durations and start dates"
                    },
                    "combined_net_rate": {
                        "type": "string",
                        "description": "The total net rate for the entire package (e.g., 'AED 2,000,000')"
                    },
                    "client_name": {
                        "type": "string",
                        "description": "Name of the client (required)"
                    }
                },
                "required": ["proposals", "combined_net_rate", "client_name"]
            }
        },
        {"type": "function", "name": "refresh_templates", "parameters": {"type": "object", "properties": {}}},
        {"type": "function", "name": "edit_task_flow", "parameters": {"type": "object", "properties": {"task_number": {"type": "integer"}, "task_data": {"type": "object"}}, "required": ["task_number", "task_data"]}},
        {
            "type": "function",
            "name": "add_location",
            "description": "Add a new location. Admin must provide ALL required metadata upfront. Digital locations require: sov, spot_duration, loop_duration, upload_fee. Static locations don't need these fields. ADMIN ONLY.", 
            "parameters": {
                "type": "object", 
                "properties": {
                    "location_key": {"type": "string", "description": "Folder/key name (lowercase, underscores for spaces, e.g., 'dubai_gateway')"},
                    "display_name": {"type": "string", "description": "Display name shown to users (e.g., 'The Dubai Gateway')"},
                    "display_type": {"type": "string", "enum": ["Digital", "Static"], "description": "Display type - determines which fields are required"},
                    "height": {"type": "string", "description": "Height with unit (e.g., '6m', '14m')"},
                    "width": {"type": "string", "description": "Width with unit (e.g., '12m', '7m')"},
                    "number_of_faces": {"type": "integer", "description": "Number of display faces (e.g., 1, 2, 4, 6)", "default": 1},
                    "series": {"type": "string", "description": "Series name (e.g., 'The Landmark Series', 'Digital Icons')"},
                    "sov": {"type": "string", "description": "Share of voice percentage - REQUIRED for Digital only (e.g., '16.6%', '12.5%')"},
                    "spot_duration": {"type": "integer", "description": "Duration of each spot in seconds - REQUIRED for Digital only (e.g., 10, 12, 16)"},
                    "loop_duration": {"type": "integer", "description": "Total loop duration in seconds - REQUIRED for Digital only (e.g., 96, 100)"},
                    "upload_fee": {"type": "integer", "description": "Upload fee in AED - REQUIRED for Digital only (e.g., 1000, 1500, 2000, 3000)"}
                }, 
                "required": ["location_key", "display_name", "display_type", "height", "width", "series"]
            }
        },
        {"type": "function", "name": "list_locations", "description": "ONLY call this when user explicitly asks to SEE or LIST available locations (e.g., 'what locations do you have?', 'show me locations', 'list all locations'). DO NOT call this when user mentions specific location names in a proposal request.", "parameters": {"type": "object", "properties": {}}},
        {
            "type": "function",
            "name": "delete_location",
            "description": "Delete an existing location (admin only, requires confirmation). ADMIN ONLY.",
            "parameters": {
                "type": "object",
                "properties": {
                    "location_key": {"type": "string", "description": "The location key or display name to delete"}
                },
                "required": ["location_key"]
            }
        },
        {"type": "function", "name": "export_proposals_to_excel", "description": "Export all proposals from the backend database to Excel and send to user. ADMIN ONLY.", "parameters": {"type": "object", "properties": {}}},
        {"type": "function", "name": "get_proposals_stats", "description": "Get summary statistics of proposals from the database", "parameters": {"type": "object", "properties": {}}},
        {"type": "function", "name": "export_booking_orders_to_excel", "description": "Export all booking orders from the backend database to Excel and send to user. Shows BO ref, client, campaign, gross total, status, dates, etc. ADMIN ONLY.", "parameters": {"type": "object", "properties": {}}},
        {
            "type": "function",
            "name": "fetch_booking_order",
            "description": "Fetch a booking order by its BO number from the original document (e.g., BL-001, VL-042, ABC123, etc). This is the BO number that appears in the client's booking order document. Returns the BO data and combined PDF file. If the BO exists but was created with outdated schema/syntax, it will be automatically regenerated with the latest format. ADMIN ONLY.",
            "parameters": {
                "type": "object",
                "properties": {
                    "bo_number": {"type": "string", "description": "The booking order number from the original document (e.g., 'BL-001', 'VL-042', 'ABC123')"}
                },
                "required": ["bo_number"]
            }
        },
        {
            "type": "function",
            "name": "generate_mockup",
            "description": "Generate a billboard mockup. IMPORTANT: If user uploads image file(s) and mentions a location for mockup, call this function IMMEDIATELY - do not ask for clarification. User can upload image(s) OR provide a text prompt for AI generation OR reuse creatives from recent mockup (within 30 min) by just specifying new location. System stores creative files for 30 minutes enabling follow-up requests on different locations. Supports multiple frames: 1 creative = duplicate across all, N creatives = match to N frames. System validates frame count compatibility automatically. Billboard variations can be specified with time_of_day (day/night/all) and finish (gold/silver/all). Use 'all' or omit to randomly select from all available variations.",
            "parameters": {
                "type": "object",
                "properties": {
                    "location": {"type": "string", "description": "The location name - intelligently match to available locations. If user says 'gateway' or 'the gateway', match to 'dubai_gateway'. If user says 'jawhara', match to 'dubai_jawhara'. Use your best judgment to infer the correct location from the available list."},
                    "time_of_day": {"type": "string", "description": "Optional time of day: 'day', 'night', or 'all' (default). Use 'all' for random selection from all time variations.", "enum": ["day", "night", "all"]},
                    "finish": {"type": "string", "description": "Optional billboard finish: 'gold', 'silver', or 'all' (default). Use 'all' for random selection from all finish variations.", "enum": ["gold", "silver", "all"]},
                    "ai_prompt": {"type": "string", "description": "Optional: AI prompt to generate billboard-ready ARTWORK ONLY (flat advertisement design, NO billboards/signs/streets in the image). System will automatically place the artwork onto the billboard. Example: 'A luxury watch advertisement with gold accents and elegant typography' - this creates the ad design itself, not a photo of a billboard"},
                    "num_ai_frames": {"type": "integer", "description": "Optional: For AI generation only - specify number of different artworks to generate (e.g., 2 for dual frames, 3 for triple frames). System will use prompt parser to create N variations and match to template with N frames. Default is 1 (single artwork)."}
                },
                "required": ["location"]
            }
        },
        {
            "type": "code_interpreter",
            "container": {"type": "auto"}
        }
    ]

    # Booking order parsing - Available to all users
    tools.append({
        "type": "function",
        "name": "parse_booking_order",
        "description": "Parse a booking order document (Excel, PDF, or image) for Backlite or Viola. Available to ALL users. Extracts client, campaign, locations, pricing, dates, and financial data. Infer the company from document content (e.g., letterhead, branding, or 'BackLite'/'Viola' text) - default to 'backlite' if unclear. Biased toward classifying uploads as ARTWORK unless clearly a booking order.",
        "parameters": {
            "type": "object",
            "properties": {
                "company": {
                    "type": "string",
                    "enum": ["backlite", "viola"],
                    "description": "Company name - either 'backlite' or 'viola'. Infer from document branding/letterhead. Default to 'backlite' if unclear."
                },
                "user_notes": {
                    "type": "string",
                    "description": "Optional notes or instructions from user about the booking order"
                }
            },
            "required": ["company"]
        }
    })

    # Admin-only tools
    if is_admin:
        admin_tools = []
        tools.extend(admin_tools)
        logger.info(f"[LLM] Admin user {user_id} - added {len(admin_tools)} admin-only tools")

    try:
        res = await config.openai_client.responses.create(model=config.OPENAI_MODEL, input=messages, tools=tools, tool_choice="auto")

        if not res.output or len(res.output) == 0:
            await config.slack_client.chat_delete(channel=channel, ts=status_ts)
            await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("I can help with proposals or add locations. Say 'add location'."))
            return

        logger.info(f"[LLM] FULL RESPONSE: {res}")
        logger.info(f"[LLM] Output items: {len(res.output)}, Types: {[item.type for item in res.output]}")

        msg = res.output[0]
        logger.info(f"[LLM] First item type: {msg.type}, hasattr name: {hasattr(msg, 'name')}")
        if hasattr(msg, 'name'):
            logger.info(f"[LLM] Function name: {msg.name}")
        if msg.type == "function_call":
            if msg.name == "get_separate_proposals":
                # Update status to Building Proposal
                await config.slack_client.chat_update(
                    channel=channel,
                    ts=status_ts,
                    text="⏳ _Building Proposal..._"
                )
                
                args = json.loads(msg.arguments)
                proposals_data = args.get("proposals", [])
                client_name = args.get("client_name") or "Unknown Client"
                
                logger.info(f"[SEPARATE] Raw args: {args}")
                logger.info(f"[SEPARATE] Proposals data: {proposals_data}")
                logger.info(f"[SEPARATE] Client: {client_name}, User: {user_id}")

                if not proposals_data:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** No proposals data provided"))
                    return
                
                result = await process_proposals(proposals_data, "separate", None, user_id, client_name)
            elif msg.name == "get_combined_proposal":
                # Update status to Building Proposal
                await config.slack_client.chat_update(
                    channel=channel,
                    ts=status_ts,
                    text="⏳ _Building Proposal..._"
                )
                
                args = json.loads(msg.arguments)
                proposals_data = args.get("proposals", [])
                combined_net_rate = args.get("combined_net_rate", None)
                client_name = args.get("client_name") or "Unknown Client"
                
                logger.info(f"[COMBINED] Raw args: {args}")
                logger.info(f"[COMBINED] Proposals data: {proposals_data}")
                logger.info(f"[COMBINED] Combined rate: {combined_net_rate}")
                logger.info(f"[COMBINED] Client: {client_name}, User: {user_id}")

                if not proposals_data:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** No proposals data provided"))
                    return
                elif not combined_net_rate:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** Combined package requires a combined net rate"))
                    return
                elif len(proposals_data) < 2:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** Combined package requires at least 2 locations"))
                    return
                
                # Transform proposals data for combined package (add durations as list with single item)
                for proposal in proposals_data:
                    if "duration" in proposal:
                        proposal["durations"] = [proposal.pop("duration")]
                        logger.info(f"[COMBINED] Transformed proposal: {proposal}")
                        
                result = await process_proposals(proposals_data, "combined", combined_net_rate, user_id, client_name)
            
            # Handle result for both get_separate_proposals and get_combined_proposal
            if msg.name in ["get_separate_proposals", "get_combined_proposal"] and 'result' in locals():
                logger.info(f"[RESULT] Processing result: {result}")
                if result["success"]:
                    # Delete status message before uploading files
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    
                    if result.get("is_combined"):
                        logger.info(f"[RESULT] Combined package - PDF: {result.get('pdf_filename')}")
                        await config.slack_client.files_upload_v2(channel=channel, file=result["pdf_path"], filename=result["pdf_filename"], initial_comment=config.markdown_to_slack(f"📦 **Combined Package Proposal**\n📍 Locations: {result['locations']}"))
                        try: os.unlink(result["pdf_path"])  # type: ignore
                        except: pass
                    elif result.get("is_single"):
                        logger.info(f"[RESULT] Single proposal - Location: {result.get('location')}")
                        await config.slack_client.files_upload_v2(channel=channel, file=result["pptx_path"], filename=result["pptx_filename"], initial_comment=config.markdown_to_slack(f"📊 **PowerPoint Proposal**\n📍 Location: {result['location']}"))
                        await config.slack_client.files_upload_v2(channel=channel, file=result["pdf_path"], filename=result["pdf_filename"], initial_comment=config.markdown_to_slack(f"📄 **PDF Proposal**\n📍 Location: {result['location']}"))
                        try:
                            os.unlink(result["pptx_path"])  # type: ignore
                            os.unlink(result["pdf_path"])  # type: ignore
                        except: pass
                    else:
                        logger.info(f"[RESULT] Multiple separate proposals - Count: {len(result.get('individual_files', []))}")
                        for f in result["individual_files"]:
                            await config.slack_client.files_upload_v2(channel=channel, file=f["path"], filename=f["filename"], initial_comment=config.markdown_to_slack(f"📊 **PowerPoint Proposal**\n📍 Location: {f['location']}"))
                        await config.slack_client.files_upload_v2(channel=channel, file=result["merged_pdf_path"], filename=result["merged_pdf_filename"], initial_comment=config.markdown_to_slack(f"📄 **Combined PDF**\n📍 All Locations: {result['locations']}"))
                        try:
                            for f in result["individual_files"]: os.unlink(f["path"])  # type: ignore
                            os.unlink(result["merged_pdf_path"])  # type: ignore
                        except: pass
                else:
                    logger.error(f"[RESULT] Error: {result.get('error')}")
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack(f"❌ **Error:** {result['error']}"))

            elif msg.name == "refresh_templates":
                config.refresh_templates()
                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("✅ Templates refreshed successfully."))

            elif msg.name == "add_location":
                # Admin permission gate
                if not config.is_admin(user_id):
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** You need admin privileges to add locations."))
                    return

                args = json.loads(msg.arguments)
                location_key = args.get("location_key", "").strip().lower().replace(" ", "_")
                
                if not location_key:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** Location key is required."))
                    return

                # Check if location already exists (filesystem + cache check for security)
                # SECURITY FIX: Previous vulnerability allowed duplicate locations when cache was stale
                # Now we check both filesystem (authoritative) and cache (fallback) to prevent bypass
                location_dir = config.TEMPLATES_DIR / location_key
                mapping = config.get_location_mapping()

                # Dual check: filesystem (primary) + cache (secondary) to prevent bypass
                if location_dir.exists() or location_key in mapping:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    if location_dir.exists():
                        logger.warning(f"[SECURITY] Duplicate location attempt blocked - filesystem check: {location_key}")
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack(f"⚠️ Location `{location_key}` already exists. Please use a different key."))
                    return
                
                # All metadata must be provided upfront
                display_name = args.get("display_name")
                display_type = args.get("display_type")
                height = args.get("height")
                width = args.get("width")
                number_of_faces = args.get("number_of_faces", 1)
                sov = args.get("sov")
                series = args.get("series")
                spot_duration = args.get("spot_duration")
                loop_duration = args.get("loop_duration")
                upload_fee = args.get("upload_fee")
                
                # Clean duration values - remove any non-numeric suffixes
                if spot_duration is not None:
                    # Convert to string first to handle the cleaning
                    spot_str = str(spot_duration).strip()
                    # Remove common suffixes like 's', 'sec', 'seconds', '"'
                    spot_str = spot_str.rstrip('s"').rstrip('sec').rstrip('seconds').strip()
                    try:
                        spot_duration = int(spot_str)
                        if spot_duration <= 0:
                            await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                            await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack(f"❌ **Error:** Spot duration must be greater than 0 seconds. Got: {spot_duration}"))
                            return
                    except ValueError:
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack(f"❌ **Error:** Invalid spot duration '{spot_duration}'. Please provide a number in seconds (e.g., 10, 12, 16)."))
                        return
                
                if loop_duration is not None:
                    # Convert to string first to handle the cleaning
                    loop_str = str(loop_duration).strip()
                    # Remove common suffixes like 's', 'sec', 'seconds', '"'
                    loop_str = loop_str.rstrip('s"').rstrip('sec').rstrip('seconds').strip()
                    try:
                        loop_duration = int(loop_str)
                        if loop_duration <= 0:
                            await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                            await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack(f"❌ **Error:** Loop duration must be greater than 0 seconds. Got: {loop_duration}"))
                            return
                    except ValueError:
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack(f"❌ **Error:** Invalid loop duration '{loop_duration}'. Please provide a number in seconds (e.g., 96, 100)."))
                        return
                
                # Validate required fields
                missing = []
                if not display_name:
                    missing.append("display_name")
                if not display_type:
                    missing.append("display_type")
                if not height:
                    missing.append("height")
                if not width:
                    missing.append("width")
                if not series:
                    missing.append("series")
                
                # For digital locations only, these fields are required
                if display_type == "Digital":
                    if not sov:
                        missing.append("sov")
                    if not spot_duration:
                        missing.append("spot_duration")
                    if not loop_duration:
                        missing.append("loop_duration")
                    if upload_fee is None:
                        missing.append("upload_fee")
                
                if missing:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack(f"❌ **Error:** Missing required fields: {', '.join(missing)}")
                    )
                    return
                
                # Store the pending location data
                pending_location_additions[user_id] = {
                    "location_key": location_key,
                    "display_name": display_name,
                    "display_type": display_type,
                    "height": height,
                    "width": width,
                    "number_of_faces": number_of_faces,
                    "sov": sov,
                    "series": series,
                    "spot_duration": spot_duration,
                    "loop_duration": loop_duration,
                    "upload_fee": upload_fee,
                    "timestamp": datetime.now()
                }
                
                logger.info(f"[LOCATION_ADD] Stored pending location for user {user_id}: {location_key}")
                logger.info(f"[LOCATION_ADD] Current pending additions: {list(pending_location_additions.keys())}")
                
                # Ask for PPT file
                summary_text = (
                    f"✅ **Location metadata validated for `{location_key}`**\n\n"
                    f"📋 **Summary:**\n"
                    f"• Display Name: {display_name}\n"
                    f"• Display Type: {display_type}\n"
                    f"• Dimensions: {height} x {width}\n"
                    f"• Faces: {number_of_faces}\n"
                    f"• Series: {series}\n"
                )
                
                # Add digital-specific fields only for digital locations
                if display_type == "Digital":
                    summary_text += (
                        f"• SOV: {sov}\n"
                        f"• Spot Duration: {spot_duration}s\n"
                        f"• Loop Duration: {loop_duration}s\n"
                        f"• Upload Fee: AED {upload_fee}\n"
                    )
                
                summary_text += "\n📎 **Please upload the PDF template file now.** (Will be converted to PowerPoint at maximum quality)\n\n⏱️ _You have 10 minutes to upload the file._"
                
                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                await config.slack_client.chat_postMessage(
                    channel=channel,
                    text=config.markdown_to_slack(summary_text)
                )
                return

            elif msg.name == "list_locations":
                names = config.available_location_names()
                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                if not names:
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("📍 No locations available. Use **'add location'** to add one."))
                else:
                    listing = "\n".join(f"• {n}" for n in names)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack(f"📍 **Current locations:**\n{listing}"))

            elif msg.name == "delete_location":
                # Admin permission gate
                if not config.is_admin(user_id):
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** You need admin privileges to delete locations."))
                    return

                args = json.loads(msg.arguments)
                location_input = args.get("location_key", "").strip()

                if not location_input:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** Please specify which location to delete."))
                    return

                # Find the actual location key - check if it's a display name or direct key
                location_key = None
                display_name = None

                # First try direct key match
                if location_input.lower().replace(" ", "_") in config.LOCATION_METADATA:
                    location_key = location_input.lower().replace(" ", "_")
                    display_name = config.LOCATION_METADATA[location_key].get('display_name', location_key)
                else:
                    # Try to match by display name
                    for key, meta in config.LOCATION_METADATA.items():
                        if meta.get('display_name', '').lower() == location_input.lower():
                            location_key = key
                            display_name = meta.get('display_name', key)
                            break

                if not location_key:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    available = ", ".join(config.available_location_names())
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack(f"❌ **Error:** Location '{location_input}' not found.\n\n**Available locations:** {available}")
                    )
                    return

                # Double confirmation - show location details and ask for confirmation
                location_dir = config.TEMPLATES_DIR / location_key
                meta = config.LOCATION_METADATA[location_key]

                confirmation_text = (
                    f"⚠️ **CONFIRM LOCATION DELETION**\n\n"
                    f"📍 **Location:** {display_name} (`{location_key}`)\n"
                    f"📊 **Type:** {meta.get('display_type', 'Unknown')}\n"
                    f"📐 **Size:** {meta.get('height')} x {meta.get('width')}\n"
                    f"🎯 **Series:** {meta.get('series', 'Unknown')}\n\n"
                    f"🚨 **WARNING:** This will permanently delete:\n"
                    f"• PowerPoint template file\n"
                    f"• Location metadata\n"
                    f"• Remove location from all future proposals\n\n"
                    f"❓ **To confirm deletion, reply with:** `confirm delete {location_key}`\n"
                    f"❓ **To cancel, reply with:** `cancel` or ignore this message"
                )

                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                await config.slack_client.chat_postMessage(
                    channel=channel,
                    text=config.markdown_to_slack(confirmation_text)
                )
                return
            
            elif msg.name == "export_proposals_to_excel":
                # Admin permission gate
                logger.info(f"[EXCEL_EXPORT] Checking admin privileges for user: {user_id}")
                is_admin_user = config.is_admin(user_id)
                logger.info(f"[EXCEL_EXPORT] User {user_id} admin status: {is_admin_user}")
                
                if not is_admin_user:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** You need admin privileges to export the database."))
                    return
                    
                logger.info("[EXCEL_EXPORT] User requested Excel export")
                try:
                    excel_path = db.export_to_excel()
                    logger.info(f"[EXCEL_EXPORT] Created Excel file at {excel_path}")
                    
                    # Get file size for display
                    file_size = os.path.getsize(excel_path)
                    size_mb = file_size / (1024 * 1024)
                    
                    # Delete status message before uploading file
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    
                    await config.slack_client.files_upload_v2(
                        channel=channel,
                        file=excel_path,
                        filename=f"proposals_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                        initial_comment=config.markdown_to_slack(
                            f"📊 **Proposals Database Export**\n"
                            f"📅 Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
                            f"📁 Size: {size_mb:.2f} MB"
                        )
                    )
                    
                    # Clean up temp file
                    try:
                        os.unlink(excel_path)
                    except:
                        pass
                        
                except Exception as e:
                    logger.error(f"[EXCEL_EXPORT] Error: {e}", exc_info=True)
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack("❌ **Error:** Failed to export database to Excel. Please try again.")
                    )
            
            elif msg.name == "export_booking_orders_to_excel":
                # Admin permission gate
                logger.info(f"[BO_EXPORT] Checking admin privileges for user: {user_id}")
                is_admin_user = config.is_admin(user_id)
                logger.info(f"[BO_EXPORT] User {user_id} admin status: {is_admin_user}")

                if not is_admin_user:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** You need admin privileges to export booking orders."))
                    return

                logger.info("[BO_EXPORT] User requested booking orders Excel export")
                try:
                    import db
                    excel_path = db.export_booking_orders_to_excel()
                    logger.info(f"[BO_EXPORT] Created Excel file at {excel_path}")

                    # Get file size for display
                    file_size = os.path.getsize(excel_path)
                    size_mb = file_size / (1024 * 1024)

                    # Delete status message before uploading file
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)

                    await config.slack_client.files_upload_v2(
                        channel=channel,
                        file=excel_path,
                        filename=f"booking_orders_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                        initial_comment=config.markdown_to_slack(
                            f"📋 **Booking Orders Database Export**\n"
                            f"📅 Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
                            f"📁 Size: {size_mb:.2f} MB"
                        )
                    )

                    # Clean up temp file
                    try:
                        os.unlink(excel_path)
                    except:
                        pass

                except Exception as e:
                    logger.error(f"[BO_EXPORT] Error: {e}", exc_info=True)
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack("❌ **Error:** Failed to export booking orders to Excel. Please try again.")
                    )

            elif msg.name == "fetch_booking_order":
                # Admin permission gate
                logger.info(f"[BO_FETCH] Checking admin privileges for user: {user_id}")
                is_admin_user = config.is_admin(user_id)
                logger.info(f"[BO_FETCH] User {user_id} admin status: {is_admin_user}")

                if not is_admin_user:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** You need admin privileges to fetch booking orders."))
                    return

                args = json.loads(msg.arguments)
                bo_number = args.get("bo_number")
                logger.info(f"[BO_FETCH] User requested BO by number: '{bo_number}' (type: {type(bo_number)}, len: {len(bo_number) if bo_number else 0})")

                try:
                    import db
                    from booking_parser import BookingOrderParser, sanitize_filename

                    # Fetch BO from database by bo_number (user-facing identifier)
                    # Query is case-insensitive and trims whitespace
                    bo_data = db.get_booking_order_by_number(bo_number)
                    logger.info(f"[BO_FETCH] Database query result: {'Found' if bo_data else 'Not found'}")

                    if not bo_data:
                        # Try to list similar BOs for debugging
                        conn = db._connect()
                        sample_bos = conn.execute("SELECT bo_number FROM booking_orders LIMIT 10").fetchall()
                        conn.close()
                        logger.info(f"[BO_FETCH] Sample BOs in database: {[bo[0] for bo in sample_bos if bo[0]]}")

                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(
                            channel=channel,
                            text=config.markdown_to_slack(f"❌ **Booking Order Not Found**\n\nBO Number: `{bo_number}` does not exist in the database.")
                        )
                        return

                    # Extract backend bo_ref for internal use and sanitize bo_number for filename
                    bo_ref = bo_data.get("bo_ref")
                    safe_bo_number = sanitize_filename(bo_number)

                    # Check if schema/syntax is outdated and regenerate if needed
                    # For now, we'll just fetch and send - regeneration logic can be added later

                    # Get the combined PDF path
                    combined_pdf_path = bo_data.get("original_file_path") or bo_data.get("parsed_excel_path")

                    if combined_pdf_path and os.path.exists(combined_pdf_path):
                        # Delete status message before uploading file
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)

                        # Send BO details with file (show user-facing bo_number)
                        details = f"📋 **Booking Order Found**\n\n"
                        details += f"**BO Number:** {bo_number}\n"
                        details += f"**Client:** {bo_data.get('client', 'N/A')}\n"
                        details += f"**Campaign:** {bo_data.get('brand_campaign', 'N/A')}\n"
                        details += f"**Gross Total:** AED {bo_data.get('gross_amount', 0):,.2f}\n"
                        details += f"**Created:** {bo_data.get('created_at', 'N/A')}\n"

                        await config.slack_client.files_upload_v2(
                            channel=channel,
                            file=combined_pdf_path,
                            filename=f"{safe_bo_number}.pdf",
                            initial_comment=config.markdown_to_slack(details)
                        )
                    else:
                        # File not found, regenerate from data
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)

                        parser = BookingOrderParser(company=bo_data.get("company", "backlite"))

                        # Generate Excel from stored data (use bo_ref for internal reference)
                        excel_path = await parser.generate_excel(bo_data, bo_ref)

                        details = f"📋 **Booking Order Found (Regenerated)**\n\n"
                        details += f"**BO Number:** {bo_number}\n"
                        details += f"**Client:** {bo_data.get('client', 'N/A')}\n"
                        details += f"**Campaign:** {bo_data.get('brand_campaign', 'N/A')}\n"
                        details += f"**Gross Total:** AED {bo_data.get('gross_amount', 0):,.2f}\n"
                        details += f"**Created:** {bo_data.get('created_at', 'N/A')}\n"
                        details += f"\n⚠️ _Original file not found - regenerated from database_"

                        await config.slack_client.files_upload_v2(
                            channel=channel,
                            file=str(excel_path),
                            filename=f"{safe_bo_number}.xlsx",
                            initial_comment=config.markdown_to_slack(details)
                        )

                        # Clean up temp file
                        try:
                            excel_path.unlink()
                        except:
                            pass

                except Exception as e:
                    logger.error(f"[BO_FETCH] Error: {e}", exc_info=True)
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack(f"❌ **Error:** Failed to fetch booking order `{bo_ref}`. Error: {str(e)}")
                    )

            elif msg.name == "get_proposals_stats":
                logger.info("[STATS] User requested proposals statistics")
                try:
                    stats = db.get_proposals_summary()

                    # Format the statistics message
                    message = "📊 **Proposals Database Summary**\n\n"
                    message += f"**Total Proposals:** {stats['total_proposals']}\n\n"

                    if stats['by_package_type']:
                        message += "**By Package Type:**\n"
                        for pkg_type, count in stats['by_package_type'].items():
                            message += f"• {pkg_type.title()}: {count}\n"
                        message += "\n"

                    if stats['recent_proposals']:
                        message += "**Recent Proposals:**\n"
                        for proposal in stats['recent_proposals']:
                            date_str = datetime.fromisoformat(proposal['date']).strftime('%Y-%m-%d %H:%M')
                            message += f"• {proposal['client']} - {proposal['locations']} ({date_str})\n"
                    else:
                        message += "_No proposals generated yet._"

                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack(message)
                    )

                except Exception as e:
                    logger.error(f"[STATS] Error: {e}", exc_info=True)

            elif msg.name == "parse_booking_order":
                # Available to all users (admin check removed per new workflow)
                args = json.loads(msg.arguments)
                company = args.get("company")
                user_notes = args.get("user_notes", "")
                await config.slack_client.chat_update(channel=channel, ts=status_ts, text="⏳ _Parsing booking order..._")
                await _handle_booking_order_parse(
                    company=company,
                    slack_event=slack_event,
                    channel=channel,
                    status_ts=status_ts,
                    user_notes=user_notes,
                    user_id=user_id,
                    user_message=user_input
                )
                return

            elif msg.name == "generate_mockup":
                # Handle mockup generation with AI or user upload
                logger.info("[MOCKUP] User requested mockup generation")

                # Parse the location from arguments
                args = json.loads(msg.arguments)
                location_name = args.get("location", "").strip()
                time_of_day = args.get("time_of_day", "").strip().lower() or "all"
                finish = args.get("finish", "").strip().lower() or "all"
                ai_prompt = args.get("ai_prompt", "").strip()
                num_ai_frames = args.get("num_ai_frames", 1) or 1  # Default to 1 if not specified

                # Convert display name to location key
                location_key = config.get_location_key_from_display_name(location_name)
                if not location_key:
                    location_key = location_name.lower().replace(" ", "_")

                # Validate location exists
                if location_key not in config.LOCATION_METADATA:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack(f"❌ **Error:** Location '{location_name}' not found. Please choose from available locations.")
                    )
                    return

                # Handle time_of_day and finish selection
                import mockup_generator
                import db

                variation_note = ""

                # Check if location has any mockup photos configured
                variations = db.list_mockup_variations(location_key)
                if not variations:
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    mockup_url = os.getenv("RENDER_EXTERNAL_URL", "http://localhost:3000") + "/mockup"
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack(
                            f"❌ **Error:** No billboard photos configured for *{location_name}* (location key: `{location_key}`).\n\n"
                            f"Ask an admin to set up mockup frames at {mockup_url}"
                        )
                    )
                    return

                # Get frame count for this location (needed for validation and storage)
                new_location_frame_count = get_location_frame_count(location_key, time_of_day, finish)

                # Get user's mockup history if exists (will validate later after checking for uploads)
                mockup_user_hist = get_mockup_history(user_id)

                # Check if user uploaded image(s) with the request
                has_images = False
                uploaded_creatives = []

                if slack_event and ("files" in slack_event or slack_event.get("subtype") == "file_share"):
                    files = slack_event.get("files", [])
                    if not files and slack_event.get("subtype") == "file_share" and "file" in slack_event:
                        files = [slack_event["file"]]

                    # Look for image files
                    for f in files:
                        filetype = f.get("filetype", "")
                        mimetype = f.get("mimetype", "")
                        filename = f.get("name", "").lower()

                        if (filetype in ["jpg", "jpeg", "png", "gif", "bmp"] or
                            mimetype.startswith("image/") or
                            any(filename.endswith(ext) for ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp"])):
                            try:
                                creative_file = await _download_slack_file(f)
                                uploaded_creatives.append(creative_file)
                                has_images = True
                                logger.info(f"[MOCKUP] Found uploaded image: {f.get('name')}")
                            except Exception as e:
                                logger.error(f"[MOCKUP] Failed to download image: {e}")

                # Determine mode based on what user provided
                # Priority: 1) New upload 2) AI prompt 3) History reuse 4) Error

                # NO EARLY FRAME VALIDATION - Allow 1 creative to be tiled across multiple frames
                # Validation happens later at line ~1916 where we check creative count, not stored frame count

                # FOLLOW-UP MODE: Check if this is a follow-up request (no upload, no AI, has history)
                if not has_images and not ai_prompt and mockup_user_hist:
                    # This is a follow-up request to apply previous creatives to a different location
                    stored_frames = mockup_user_hist.get("metadata", {}).get("num_frames", 1)
                    stored_creative_paths = mockup_user_hist.get("creative_paths", [])
                    stored_location = mockup_user_hist.get("metadata", {}).get("location_name", "unknown")

                    # Verify all creative files still exist on disk
                    missing_files = []
                    for creative_path in stored_creative_paths:
                        if not creative_path.exists():
                            missing_files.append(str(creative_path))

                    if missing_files:
                        logger.error(f"[MOCKUP] Creative files missing from history: {missing_files}")
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(
                            channel=channel,
                            text=config.markdown_to_slack(
                                f"❌ **Error:** Your previous creative files are no longer available.\n\n"
                                f"Please upload new images or use AI generation."
                            )
                        )
                        # Clean up corrupted history
                        del mockup_history[user_id]
                        return

                    # Validate creative count: Allow if 1 creative (tile across frames) OR matches frame count
                    num_stored_creatives = len(stored_creative_paths)
                    is_valid_count = (num_stored_creatives == 1) or (num_stored_creatives == new_location_frame_count)

                    if is_valid_count:
                        logger.info(f"[MOCKUP] Follow-up request detected - reusing {len(stored_creative_paths)} creative(s) from history")

                        await config.slack_client.chat_update(
                            channel=channel,
                            ts=status_ts,
                            text=f"⏳ _Applying your previous creative(s) to {location_name}..._"
                        )

                        try:
                            # Generate mockup using stored creatives (queued)
                            result_path, _ = await _generate_mockup_queued(
                                location_key,
                                stored_creative_paths,
                                time_of_day=time_of_day,
                                finish=finish
                            )

                            if not result_path:
                                raise Exception("Failed to generate mockup")

                            # Upload mockup
                            await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                            variation_info = ""
                            if time_of_day != "all" or finish != "all":
                                variation_info = f" ({time_of_day}/{finish})"

                            frames_info = f" ({stored_frames} frame(s))" if stored_frames > 1 else ""

                            await config.slack_client.files_upload_v2(
                                channel=channel,
                                file=str(result_path),
                                filename=f"mockup_{location_key}_{time_of_day}_{finish}.jpg",
                                initial_comment=config.markdown_to_slack(
                                    f"🎨 **Billboard Mockup Generated** (Follow-up)\n\n"
                                    f"📍 New Location: {location_name}{variation_info}\n"
                                    f"🔄 Using creative(s) from: {stored_location}{frames_info}\n"
                                    f"✨ Your creative has been applied to this location."
                                )
                            )

                            # Update history with new location (but keep same creatives)
                            mockup_user_hist["metadata"]["location_key"] = location_key
                            mockup_user_hist["metadata"]["location_name"] = location_name
                            mockup_user_hist["metadata"]["time_of_day"] = time_of_day
                            mockup_user_hist["metadata"]["finish"] = finish

                            logger.info(f"[MOCKUP] Follow-up mockup generated successfully for user {user_id}")

                            # Cleanup final mockup
                            try:
                                os.unlink(result_path)
                            except:
                                pass

                            # Force garbage collection to free memory from numpy arrays
                            import gc
                            gc.collect()
                            logger.debug(f"[MOCKUP] Follow-up mode: Forced garbage collection")

                            return  # Done with follow-up

                        except Exception as e:
                            logger.error(f"[MOCKUP] Error generating follow-up mockup: {e}", exc_info=True)
                            await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                            await config.slack_client.chat_postMessage(
                                channel=channel,
                                text=config.markdown_to_slack(f"❌ **Error:** Failed to generate follow-up mockup. {str(e)}")
                            )

                            # Cleanup result file if it was created before error
                            try:
                                if 'result_path' in locals() and result_path and result_path.exists():
                                    os.unlink(result_path)
                                    logger.info(f"[MOCKUP] Cleaned up partial result file after error")
                            except Exception as cleanup_error:
                                logger.error(f"[MOCKUP] Failed to cleanup result file: {cleanup_error}")

                            # Force garbage collection
                            import gc
                            gc.collect()

                            return
                    else:
                        # Invalid creative count (e.g., 2 creatives for 3 frames)
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(
                            channel=channel,
                            text=config.markdown_to_slack(
                                f"⚠️ **Creative Count Mismatch**\n\n"
                                f"I have **{num_stored_creatives} creative(s)** from your previous mockup (**{stored_location}**), "
                                f"but **{location_name}** requires **{new_location_frame_count} frame(s)**.\n\n"
                                f"**Valid options:**\n"
                                f"• Upload **1 image** (will be tiled across all frames)\n"
                                f"• Upload **{new_location_frame_count} images** (one per frame)\n"
                                f"• Use AI generation with a creative description"
                            )
                        )
                        return

                # Now proceed with normal modes (Priority: Upload > AI > Error)
                if has_images:
                    # UPLOAD MODE: User uploaded image(s) - this takes priority over AI
                    logger.info(f"[MOCKUP] Processing {len(uploaded_creatives)} uploaded image(s)")

                    # Validate creative count: Allow 1 (tile) OR match frame count
                    num_uploaded = len(uploaded_creatives)
                    is_valid_upload_count = (num_uploaded == 1) or (num_uploaded == new_location_frame_count)

                    if not is_valid_upload_count:
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(
                            channel=channel,
                            text=config.markdown_to_slack(
                                f"⚠️ **Creative Count Mismatch**\n\n"
                                f"You uploaded **{num_uploaded} image(s)**, but **{location_name}** requires **{new_location_frame_count} frame(s)**.\n\n"
                                f"**Valid options:**\n"
                                f"• Upload **1 image** (will be tiled across all frames)\n"
                                f"• Upload **{new_location_frame_count} images** (one per frame)"
                            )
                        )
                        return

                    await config.slack_client.chat_update(
                        channel=channel,
                        ts=status_ts,
                        text="⏳ _Generating mockup from uploaded image(s)..._"
                    )

                    try:
                        # Generate mockup using uploaded creatives (queued)
                        result_path, _ = await _generate_mockup_queued(
                            location_key,
                            uploaded_creatives,
                            time_of_day=time_of_day,
                            finish=finish
                        )

                        if not result_path:
                            raise Exception("Failed to generate mockup")

                        # Delete status and upload mockup
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        variation_info = ""
                        if time_of_day != "all" or finish != "all":
                            variation_info = f" ({time_of_day}/{finish})"
                        await config.slack_client.files_upload_v2(
                            channel=channel,
                            file=str(result_path),
                            filename=f"mockup_{location_key}_{time_of_day}_{finish}.jpg",
                            initial_comment=config.markdown_to_slack(
                                f"🎨 **Billboard Mockup Generated**\n\n"
                                f"📍 Location: {location_name}{variation_info}\n"
                                f"🖼️ Creative(s): {len(uploaded_creatives)} image(s)\n"
                                f"✨ Your creative has been applied to a billboard photo.{variation_note}"
                            )
                        )

                        # Get frame count for validation in follow-ups
                        location_frame_count = get_location_frame_count(location_key, time_of_day, finish)

                        # Store creative files in 30-minute history for follow-ups on other locations
                        store_mockup_history(user_id, uploaded_creatives, {
                            "location_key": location_key,
                            "location_name": location_name,
                            "time_of_day": time_of_day,
                            "finish": finish,
                            "mode": "uploaded",
                            "num_frames": location_frame_count or 1
                        })
                        logger.info(f"[MOCKUP] Stored {len(uploaded_creatives)} uploaded creative(s) in history for user {user_id} ({location_frame_count} frames)")

                        # Cleanup final mockup (we keep creatives in history, not the result)
                        try:
                            os.unlink(result_path)
                        except:
                            pass

                        # Force garbage collection to free memory from numpy arrays
                        import gc
                        gc.collect()
                        logger.debug(f"[MOCKUP] Upload mode: Forced garbage collection")

                    except Exception as e:
                        logger.error(f"[MOCKUP] Error generating mockup from upload: {e}", exc_info=True)
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(
                            channel=channel,
                            text=config.markdown_to_slack(f"❌ **Error:** Failed to generate mockup. {str(e)}")
                        )

                        # Cleanup uploaded creative files on error
                        try:
                            for creative_file in uploaded_creatives:
                                os.unlink(creative_file)
                        except:
                            pass

                        # Force garbage collection
                        import gc
                        gc.collect()

                elif ai_prompt:
                    # AI MODE: User provided a description for AI generation

                    # Validate num_ai_frames: Allow 1 (tile) OR match frame count
                    is_valid_ai_count = (num_ai_frames == 1) or (num_ai_frames == new_location_frame_count)

                    if not is_valid_ai_count:
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(
                            channel=channel,
                            text=config.markdown_to_slack(
                                f"⚠️ **Creative Count Mismatch**\n\n"
                                f"You requested **{num_ai_frames} AI creatives**, but **{location_name}** requires **{new_location_frame_count} frame(s)**.\n\n"
                                f"**Valid options:**\n"
                                f"• Generate **1 creative** (will be tiled across all frames) - omit `num_ai_frames` parameter\n"
                                f"• Generate **{new_location_frame_count} creatives** (one per frame) - set `num_ai_frames={new_location_frame_count}`"
                            )
                        )
                        return

                    await config.slack_client.chat_update(
                        channel=channel,
                        ts=status_ts,
                        text="⏳ _Generating AI creative and mockup..._"
                    )

                    try:
                        # Detect orientation for this location
                        import mockup_generator
                        is_portrait = mockup_generator.is_portrait_location(location_key)

                        if is_portrait:
                            orientation_text = """📐 FORMAT & DIMENSIONS:
- Aspect ratio: Tall portrait (roughly 2:3 ratio)
- Orientation: Vertical/portrait ONLY
- Canvas: Perfectly flat, rectangular, no warping or perspective
- Fill entire frame edge-to-edge with design
- No white borders, frames, or margins around the design"""
                        else:
                            orientation_text = """📐 FORMAT & DIMENSIONS:
- Aspect ratio: Wide landscape (roughly 3:2 ratio)
- Orientation: Horizontal/landscape ONLY
- Canvas: Perfectly flat, rectangular, no warping or perspective
- Fill entire frame edge-to-edge with design
- No white borders, frames, or margins around the design"""

                        # Extensive system prompt for billboard artwork generation
                        enhanced_prompt = f"""Create a professional outdoor advertising billboard creative - IMPORTANT: This is the FLAT 2D ARTWORK FILE that will be printed and placed ON a billboard, NOT a photograph of an existing billboard.

═══════════════════════════════════════════════════════════
CRITICAL DISTINCTIONS:
═══════════════════════════════════════════════════════════

✅ CORRECT OUTPUT (what we want):
- A flat, rectangular advertisement design (like a Photoshop/Illustrator file)
- The actual graphic design artwork that goes ON the billboard surface
- Think: magazine ad, poster design, digital banner creative
- Perfectly rectangular, no perspective, no angle, no depth
- Edge-to-edge design filling the entire rectangular canvas
- Like looking at a computer screen showing the ad design

❌ INCORRECT OUTPUT (what we DON'T want):
- A photograph of a physical billboard in a street scene
- 3D rendering showing billboard from an angle/perspective
- Image with billboard frame, poles, or support structure visible
- Photo showing buildings, sky, roads, or environment around billboard
- Any mockup showing how the billboard looks in real life
- Perspective view, vanishing points, or dimensional representation

═══════════════════════════════════════════════════════════
DETAILED DESIGN REQUIREMENTS:
═══════════════════════════════════════════════════════════

{orientation_text}

🎨 VISUAL DESIGN PRINCIPLES:
- Bold, high-impact composition that catches attention immediately
- Large hero image or visual focal point (50-70% of design)
- Vibrant, saturated colors that pop in daylight
- High contrast between elements for maximum visibility
- Simple, uncluttered layout (viewer has 5-7 seconds max)
- Professional photo quality or clean vector graphics
- Modern, contemporary advertising aesthetic

✍️ TYPOGRAPHY (if text is needed):
- LARGE, bold, highly readable fonts
- Sans-serif typefaces work best for outdoor viewing
- Maximum 7-10 words total (fewer is better)
- High contrast text-to-background ratio
- Text size: headlines should occupy 15-25% of vertical height
- Clear hierarchy: one main message, optional supporting text
- Avoid script fonts, thin fonts, or decorative typefaces
- Letter spacing optimized for distance reading

🎯 COMPOSITION STRATEGY:
- Rule of thirds or strong visual hierarchy
- One clear focal point (don't scatter attention)
- Negative space used strategically
- Visual flow guides eye to key message/CTA
- Brand logo prominent but not dominating (10-15% of space)
- Clean, professional layout with breathing room

💡 COLOR THEORY FOR OUTDOOR:
- Vibrant, saturated colors (avoid pastels or muted tones)
- High contrast pairings: dark on light or light on dark
- Colors that work in bright sunlight and shadows
- Use colors appropriate to the brand and message
- Background should enhance, not compete with message
- Avoid repetitive color schemes - vary your palette based on the creative brief

🔍 QUALITY STANDARDS:
- Sharp, crisp graphics (no blur, pixelation, or artifacts)
- Professional commercial photography or illustration
- Consistent lighting across all design elements
- No watermarks, stock photo markers, or placeholder text
- Print-ready quality at large scale
- Polished, agency-level execution

═══════════════════════════════════════════════════════════
CREATIVE BRIEF:
═══════════════════════════════════════════════════════════

{ai_prompt}

═══════════════════════════════════════════════════════════
⚠️ CRITICAL - FINAL REMINDER - READ CAREFULLY:
═══════════════════════════════════════════════════════════

🚫 ABSOLUTELY DO NOT INCLUDE:
- NO billboards, signs, or advertising structures
- NO street scenes, highways, or roads
- NO people holding/viewing the ad
- NO frames, borders, or physical contexts
- NO 3D perspective or mockup views
- NO environmental surroundings whatsoever

✅ YOU MUST CREATE:
- The FLAT ARTWORK FILE ONLY - the pure advertisement design
- A rectangular graphic that will be PLACED onto a billboard LATER
- Think: graphic designer working in Photoshop/Illustrator
- The final output is the CREATIVE CONTENT, not a mockup

📐 DELIVERABLE:
Imagine you're delivering a print file to a billboard company.
They will take YOUR flat design and apply it to their billboard.
Your job: create the artwork. Their job: put it on the billboard.

Example: If asked for a "Nike shoe ad," create the advertisement graphic (shoe + slogan + logo),
NOT a photo of a billboard displaying that ad on the street.

DELIVER ONLY THE FLAT, RECTANGULAR ADVERTISEMENT ARTWORK - NOTHING ELSE."""

                        # Update status to show we're generating
                        frames_text = f"{num_ai_frames} artworks and mockup" if num_ai_frames > 1 else "AI artwork and mockup"
                        await config.slack_client.chat_update(
                            channel=channel,
                            ts=status_ts,
                            text=f"⏳ _Generating {frames_text}..._"
                        )

                        # Generate AI creative(s) + mockup through queue (prevents memory spikes)
                        result_path, ai_creative_paths = await _generate_ai_mockup_queued(
                            ai_prompt=ai_prompt,
                            enhanced_prompt=enhanced_prompt,
                            num_ai_frames=num_ai_frames,
                            location_key=location_key,
                            time_of_day=time_of_day,
                            finish=finish
                        )

                        if not result_path:
                            raise Exception("Failed to generate mockup")

                        # Delete status and upload mockup
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        variation_info = ""
                        if time_of_day != "all" or finish != "all":
                            variation_info = f" ({time_of_day}/{finish})"

                        frames_info = f" ({num_ai_frames} frames)" if num_ai_frames > 1 else ""

                        await config.slack_client.files_upload_v2(
                            channel=channel,
                            file=str(result_path),
                            filename=f"ai_mockup_{location_key}_{time_of_day}_{finish}.jpg",
                            initial_comment=config.markdown_to_slack(
                                f"🎨 **AI-Generated Billboard Mockup**\n\n"
                                f"📍 Location: {location_name}{variation_info}{frames_info}\n"
                            )
                        )

                        # Store creative files in 30-minute history for follow-ups on other locations
                        store_mockup_history(user_id, ai_creative_paths, {
                            "location_key": location_key,
                            "location_name": location_name,
                            "time_of_day": time_of_day,
                            "finish": finish,
                            "mode": "ai_generated",
                            "num_frames": num_ai_frames
                        })
                        logger.info(f"[MOCKUP] Stored {len(ai_creative_paths)} AI creative(s) in history for user {user_id}")

                        # Cleanup final mockup (we keep creatives in history, not the result)
                        try:
                            os.unlink(result_path)
                        except:
                            pass

                        # Force garbage collection to free memory from numpy arrays
                        import gc
                        gc.collect()
                        logger.debug(f"[MOCKUP] AI mode: Forced garbage collection")

                    except Exception as e:
                        logger.error(f"[MOCKUP] Error generating AI mockup: {e}", exc_info=True)
                        await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                        await config.slack_client.chat_postMessage(
                            channel=channel,
                            text=config.markdown_to_slack(f"❌ **Error:** Failed to generate AI mockup. {str(e)}")
                        )

                        # Cleanup any AI creative files that were generated before the error
                        try:
                            for creative_path in ai_creative_paths:
                                if creative_path and creative_path.exists():
                                    os.unlink(creative_path)
                            logger.info(f"[MOCKUP] Cleaned up {len(ai_creative_paths)} AI creative file(s) after error")
                        except Exception as cleanup_error:
                            logger.error(f"[MOCKUP] Failed to cleanup AI creatives: {cleanup_error}")

                        # Force garbage collection
                        import gc
                        gc.collect()

                else:
                    # NO AI PROMPT, NO IMAGE UPLOADED, NO HISTORY: Error - user needs to provide creative
                    await config.slack_client.chat_delete(channel=channel, ts=status_ts)
                    await config.slack_client.chat_postMessage(
                        channel=channel,
                        text=config.markdown_to_slack(
                            f"❌ **Sorry!** You need to provide a creative for the mockup.\n\n"
                            f"**Three ways to generate mockups:**\n\n"
                            f"1️⃣ **Upload Your Image:** Attach your creative when you send the request\n"
                            f"   Example: [Upload creative.jpg] + \"mockup for {location_name}\"\n\n"
                            f"2️⃣ **AI Generation (No upload needed):** Describe what you want\n"
                            f"   Example: \"mockup for {location_name} with luxury watch ad, gold and elegant typography\"\n"
                            f"   The AI will generate the creative for you!\n\n"
                            f"3️⃣ **Follow-up Request:** If you recently generated a mockup (within 30 min), just ask!\n"
                            f"   Example: \"show me this on {location_name}\" or \"apply to {location_name}\"\n"
                            f"   I'll reuse your previous creative(s) automatically.\n\n"
                            f"Please try again with an image attachment, creative description, or generate a mockup first!"
                        )
                    )

        else:
            reply = msg.content[-1].text if hasattr(msg, 'content') and msg.content else "How can I help you today?"
            # Format any markdown-style text from the LLM
            formatted_reply = reply
            # Ensure bullet points are properly formatted
            formatted_reply = formatted_reply.replace('\n- ', '\n• ')
            formatted_reply = formatted_reply.replace('\n* ', '\n• ')
            # Ensure headers are bolded
            import re
            formatted_reply = re.sub(r'^(For .+:)$', r'**\1**', formatted_reply, flags=re.MULTILINE)
            formatted_reply = re.sub(r'^([A-Z][A-Z\s]+:)$', r'**\1**', formatted_reply, flags=re.MULTILINE)
            # Delete status message before sending reply
            await config.slack_client.chat_delete(channel=channel, ts=status_ts)
            await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack(formatted_reply))

        user_history[user_id] = history[-10:]

    except Exception as e:
        config.logger.error(f"LLM loop error: {e}", exc_info=True)
        # Try to delete status message if it exists
        try:
            if 'status_ts' in locals() and status_ts:
                await config.slack_client.chat_delete(channel=channel, ts=status_ts)
        except:
            pass
        await config.slack_client.chat_postMessage(channel=channel, text=config.markdown_to_slack("❌ **Error:** Something went wrong. Please try again.")) 