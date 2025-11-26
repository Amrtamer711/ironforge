import os
import tempfile
import subprocess
import platform
import shutil
from pathlib import Path
import asyncio

from pypdf import PdfWriter, PdfReader
from pptx import Presentation

import config

# Limit concurrent conversions to avoid CPU/app contention
# With 2 CPUs, we can handle more concurrent conversions
_CONVERT_SEMAPHORE = asyncio.Semaphore(int(os.getenv("PDF_CONVERT_CONCURRENCY", "4")))


async def convert_pptx_to_pdf_async(pptx_path: str) -> str:
    """Async wrapper for PDF conversion with semaphore protection"""
    async with _CONVERT_SEMAPHORE:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, convert_pptx_to_pdf, pptx_path)


def convert_pptx_to_pdf(pptx_path: str) -> str:
    logger = config.logger
    logger.info(f"[PDF_CONVERT] Starting conversion of '{pptx_path}'")
    
    pdf_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    pdf_file.close()
    logger.info(f"[PDF_CONVERT] Target PDF path: '{pdf_file.name}'")

    system = platform.system()
    logger.info(f"[PDF_CONVERT] Operating system: {system}")

    libreoffice_paths = [
        '/usr/bin/libreoffice',  # Docker/Linux standard location
        '/usr/bin/soffice',      # Alternative name
        '/opt/libreoffice/program/soffice',  # Some installations
        '/usr/local/bin/libreoffice',
        '/opt/homebrew/bin/soffice',  # macOS homebrew
        'libreoffice',  # PATH lookup
        'soffice',      # PATH lookup
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS
    ]

    for lo_path in libreoffice_paths:
        if shutil.which(lo_path) or os.path.exists(lo_path):
            try:
                logger.info(f"[PDF_CONVERT] Trying LibreOffice at '{lo_path}'")
                cmd = [lo_path, '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(pdf_file.name), pptx_path]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
                if result.returncode == 0:
                    converted_pdf = os.path.join(
                        os.path.dirname(pdf_file.name),
                        os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
                    )
                    if os.path.exists(converted_pdf):
                        shutil.move(converted_pdf, pdf_file.name)
                        logger.info(f"[PDF_CONVERT] Successfully converted using LibreOffice at '{lo_path}'")
                        return pdf_file.name
                    else:
                        logger.warning(f"[PDF_CONVERT] Converted file not found at expected location: {converted_pdf}")
                else:
                    logger.warning(f"[PDF_CONVERT] LibreOffice at '{lo_path}' failed with code {result.returncode}")
                    logger.warning(f"[PDF_CONVERT] stdout: {result.stdout}")
                    logger.warning(f"[PDF_CONVERT] stderr: {result.stderr}")
            except Exception as e:
                logger.debug(f"[PDF_CONVERT] LibreOffice conversion failed: {e}")
                continue

    if shutil.which('unoconv'):
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', pdf_file.name, pptx_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            if result.returncode == 0 and os.path.exists(pdf_file.name):
                return pdf_file.name
        except Exception as e:
            config.logger.debug(f"unoconv conversion failed: {e}")

    if system == "Darwin":
        try:
            powerpoint_script = f'''
            tell application "Microsoft PowerPoint"
                open POSIX file "{pptx_path}"
                save active presentation in POSIX file "{pdf_file.name}" as save as PDF
                close active presentation
            end tell
            '''
            result = subprocess.run(['osascript', '-e', powerpoint_script], capture_output=True, text=True, timeout=60)
            if result.returncode == 0 and os.path.exists(pdf_file.name):
                return pdf_file.name
        except Exception as e:
            config.logger.debug(f"PowerPoint conversion failed: {e}")

        try:
            keynote_script = f'''
            tell application "Keynote"
                open POSIX file "{pptx_path}"
                export front document to POSIX file "{pdf_file.name}" as PDF
                close front document
            end tell
            '''
            result = subprocess.run(['osascript', '-e', keynote_script], capture_output=True, text=True, timeout=60)
            if result.returncode == 0 and os.path.exists(pdf_file.name):
                return pdf_file.name
        except Exception as e:
            config.logger.debug(f"Keynote conversion failed: {e}")

    # Fallback: text-only extraction
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.pdfgen import canvas
    from reportlab.lib import colors

    try:
        pres = Presentation(pptx_path)
        page_width, page_height = landscape(letter)
        c = canvas.Canvas(pdf_file.name, pagesize=landscape(letter))
        for slide_idx, slide in enumerate(pres.slides):
            if slide.background and hasattr(slide.background, 'fill'):
                try:
                    if slide.background.fill.type == 1:
                        bg_color = slide.background.fill.fore_color.rgb
                        if bg_color:
                            c.setFillColorRGB(bg_color[0]/255.0, bg_color[1]/255.0, bg_color[2]/255.0)
                            c.rect(0, 0, page_width, page_height, fill=1, stroke=0)
                except:
                    pass
            c.setFillColor(colors.black)
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text') and shape.text.strip():
                        left = float(shape.left) / 914400 * 72
                        top = float(shape.top) / 914400 * 72
                        y_pos = page_height - top - 50
                        text = shape.text.strip()
                        font_size = 12
                        if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                            for para in shape.text_frame.paragraphs:
                                if para.runs:
                                    run = para.runs[0]
                                    if run.font.size:
                                        font_size = run.font.size.pt
                        c.setFont("Helvetica", min(font_size, 24))
                        lines = text.split('\n')
                        for line in lines:
                            if line.strip():
                                c.drawString(left, y_pos, line.strip())
                                y_pos -= font_size + 5
                except Exception as e:
                    config.logger.debug(f"Error processing shape: {e}")
            c.setFont("Helvetica", 10)
            c.drawString(page_width - 100, 30, f"Slide {slide_idx + 1}")
            if slide_idx < len(pres.slides) - 1:
                c.showPage()
        c.save()
        config.logger.warning("PDF created using fallback text extraction. Install LibreOffice for fidelity.")
        return pdf_file.name
    except Exception as e:
        config.logger.error(f"PDF conversion failed: {e}")
        raise


def merge_pdfs(pdf_files: list) -> str:
    logger = config.logger
    logger.info(f"[PDF_MERGE] Merging {len(pdf_files)} PDF files")
    for idx, pdf in enumerate(pdf_files):
        logger.info(f"[PDF_MERGE]   File {idx + 1}: '{pdf}'")
    
    output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    output_file.close()
    logger.info(f"[PDF_MERGE] Output file: '{output_file.name}'")
    
    pdf_writer = PdfWriter()
    for pdf_path in pdf_files:
        pdf_reader = PdfReader(pdf_path)
        page_count = len(pdf_reader.pages)
        logger.info(f"[PDF_MERGE] Adding {page_count} pages from '{pdf_path}'")
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
    
    with open(output_file.name, 'wb') as output:
        pdf_writer.write(output)
    
    logger.info(f"[PDF_MERGE] Successfully merged PDFs to '{output_file.name}'")
    return output_file.name


async def remove_slides_and_convert_to_pdf(pptx_path: str, remove_first: bool = False, remove_last: bool = False) -> str:
    import shutil as _sh
    import tempfile as _tf
    
    logger = config.logger
    logger.info(f"[REMOVE_SLIDES] Processing '{pptx_path}'")
    logger.info(f"[REMOVE_SLIDES] Remove first: {remove_first}, Remove last: {remove_last}")

    async with _CONVERT_SEMAPHORE:
        temp_pptx = _tf.NamedTemporaryFile(delete=False, suffix=".pptx")
        temp_pptx.close()
        _sh.copy2(pptx_path, temp_pptx.name)
        logger.info(f"[REMOVE_SLIDES] Created temp file: '{temp_pptx.name}'")

        pres = Presentation(temp_pptx.name)
        xml_slides = pres.slides._sldIdLst
        slides_to_remove = []

        if remove_first and len(pres.slides) > 0:
            slides_to_remove.append(list(xml_slides)[0])
        if remove_last and len(pres.slides) > 1:
            slides_to_remove.append(list(xml_slides)[-1])

        for slide_id in slides_to_remove:
            if slide_id in xml_slides:
                xml_slides.remove(slide_id)

        pres.save(temp_pptx.name)
        pdf_path = convert_pptx_to_pdf(temp_pptx.name)
        try:
            os.unlink(temp_pptx.name)
        except:
            pass
        return pdf_path 