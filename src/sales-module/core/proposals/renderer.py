"""
Proposal Renderer Module.

Handles PPTX/PDF generation for proposals.
"""

import tempfile
from typing import Any

from pptx import Presentation

from generators.pptx import create_combined_financial_proposal_slide, create_financial_proposal_slide


class ProposalRenderer:
    """
    Renders proposals to PPTX/PDF format.

    Responsibilities:
    - Create financial slides in PPTX templates
    - Generate separate proposal slides
    - Generate combined package slides
    """

    def create_proposal_with_template(
        self,
        source_path: str,
        financial_data: dict,
        currency: str = None
    ) -> tuple[str, list[str], list[str]]:
        """
        Create single proposal with financial slide.

        Args:
            source_path: Path to template PPTX file
            financial_data: Financial data dict with location, dates, rates, etc.
            currency: Optional currency code (e.g., 'USD', 'EUR')

        Returns:
            Tuple of (pptx_path, vat_amounts, total_amounts)
            - pptx_path: Path to generated PPTX file
            - vat_amounts: List of VAT amounts as strings
            - total_amounts: List of total amounts as strings

        Example:
            >>> renderer = ProposalRenderer()
            >>> pptx_path, vat, totals = renderer.create_proposal_with_template(
            ...     "/templates/location.pptx",
            ...     {"location": "dubai_gateway", "durations": [4, 8], ...},
            ...     currency="AED"
            ... )
        """
        pres = Presentation(source_path)
        insert_position = max(len(pres.slides) - 1, 0)
        slide_width = pres.slide_width
        slide_height = pres.slide_height

        # Add blank slide for financial data
        blank_layout = pres.slide_layouts[6] if len(pres.slide_layouts) > 6 else pres.slide_layouts[0]
        financial_slide = pres.slides.add_slide(blank_layout)

        # Create financial slide content
        vat_amounts, total_amounts = create_financial_proposal_slide(
            financial_slide,
            financial_data,
            slide_width,
            slide_height,
            currency
        )

        # Move financial slide to second-to-last position
        if len(pres.slides) > 1 and insert_position < len(pres.slides) - 1:
            xml_slides = pres.slides._sldIdLst
            new_slide_element = xml_slides[-1]
            xml_slides.remove(new_slide_element)
            xml_slides.insert(insert_position, new_slide_element)

        # Save to temporary file
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        pres.save(tmp.name)

        return tmp.name, vat_amounts, total_amounts

    def create_combined_proposal_with_template(
        self,
        source_path: str,
        proposals_data: list,
        combined_net_rate: str,
        client_name: str,
        payment_terms: str = "100% upfront",
        currency: str = None
    ) -> tuple[str, str]:
        """
        Create combined package proposal with financial slide.

        Args:
            source_path: Path to template PPTX file
            proposals_data: List of proposal dicts
            combined_net_rate: Combined net rate for package
            client_name: Client name
            payment_terms: Payment terms text
            currency: Optional currency code

        Returns:
            Tuple of (pptx_path, total_combined)
            - pptx_path: Path to generated PPTX file
            - total_combined: Total amount as string

        Example:
            >>> renderer = ProposalRenderer()
            >>> pptx_path, total = renderer.create_combined_proposal_with_template(
            ...     "/templates/location.pptx",
            ...     proposals_data,
            ...     "50000",
            ...     "ABC Corp",
            ...     currency="AED"
            ... )
        """
        pres = Presentation(source_path)
        insert_position = max(len(pres.slides) - 1, 0)
        slide_width = pres.slide_width
        slide_height = pres.slide_height

        # Add slide for combined financial data
        layout = pres.slide_layouts[0]
        financial_slide = pres.slides.add_slide(layout)

        # Clear existing content
        for shape in list(financial_slide.shapes):
            if hasattr(shape, "text_frame"):
                shape.text_frame.clear()

        # Create combined financial slide content
        total_combined = create_combined_financial_proposal_slide(
            financial_slide,
            proposals_data,
            combined_net_rate,
            slide_width,
            slide_height,
            client_name,
            payment_terms,
            currency,
        )

        # Move financial slide to second-to-last position
        xml_slides = pres.slides._sldIdLst
        slides_list = list(xml_slides)
        new_slide_element = slides_list[-1]
        xml_slides.remove(new_slide_element)
        xml_slides.insert(insert_position, new_slide_element)

        # Save to temporary file
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        pres.save(tmp.name)

        return tmp.name, total_combined
