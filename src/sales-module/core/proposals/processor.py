"""
Proposal Processing Orchestrator.

Coordinates validation, rendering, and PDF generation for proposals.
"""

import asyncio
import os
import shutil
import tempfile
import time
from datetime import datetime, timedelta, timezone
from typing import Any

from pptx import Presentation
from pypdf import PdfReader

import config

# Logger shorthand for timing logs
logger = config.logger
from core.services.template_service import TemplateService
from db.database import db
from generators.pdf import convert_pptx_to_pdf, merge_pdfs, remove_slides_and_convert_to_pdf

from .intro_outro import IntroOutroHandler
from .renderer import ProposalRenderer
from .validator import ProposalValidator


class ProposalProcessor:
    """
    Orchestrates proposal generation workflow.

    Responsibilities:
    - Coordinate validation, rendering, PDF generation
    - Handle intro/outro slide extraction
    - Merge PDFs for multi-proposal packages
    - Log proposals to database
    - Clean up temporary files
    """

    def __init__(
        self,
        validator: ProposalValidator,
        renderer: ProposalRenderer,
        intro_outro_handler: IntroOutroHandler,
        template_service: TemplateService,
        user_companies: list[str] | None = None,
    ):
        """
        Initialize processor with dependencies.

        Args:
            validator: ProposalValidator instance
            renderer: ProposalRenderer instance
            intro_outro_handler: IntroOutroHandler instance
            template_service: TemplateService instance (required)
            user_companies: List of company schemas for package expansion
        """
        self.validator = validator
        self.renderer = renderer
        self.intro_outro_handler = intro_outro_handler
        self.template_service = template_service
        self.logger = config.logger

        # Request-scoped cache for extracted intro/outro slides
        # Avoids re-extracting pages from the same PDF within a single request
        # Maps: pdf_name -> (intro_pdf_path, outro_pdf_path)
        self._extracted_slides_cache: dict[str, tuple[str, str]] = {}

        # User companies for package expansion
        self._user_companies = user_companies

    async def _check_if_package(self, location_key: str) -> tuple[bool, list[str]]:
        """
        Check if location_key is a package and return its network keys.

        Args:
            location_key: Location key to check

        Returns:
            Tuple of (is_package, network_keys)
            e.g., (True, ["dna01", "dna02", "dna03", "dna04"])
        """
        if not self._user_companies:
            return False, []

        from core.services.mockup_service.package_expander import PackageExpander

        expander = PackageExpander(self._user_companies)
        package_data = await expander._find_package(location_key.lower().strip())

        if package_data:
            targets = await expander._expand_package(package_data)
            network_keys = [t.network_key for t in targets]
            self.logger.info(f"[PROCESSOR] Package '{location_key}' expands to: {network_keys}")
            return True, network_keys

        return False, []

    async def _process_package_networks_pdf(
        self,
        network_keys: list[str],
        proposal: dict,
        company_hint: str | None = None
    ) -> tuple[str | None, list[str]]:
        """
        Process all networks in a package using PDF-first strategy.

        Downloads PDF for each network, extracts content (strips intro/outro),
        and merges them into a single PDF. Continues with available networks
        if some are missing.

        Args:
            network_keys: List of network keys in the package
            proposal: The original proposal data
            company_hint: Company schema hint for template lookup

        Returns:
            Tuple of (merged_pdf_path, missing_networks):
            - merged_pdf_path: Path to merged PDF, or None if ALL networks unavailable
            - missing_networks: List of network keys that had no templates
        """
        loop = asyncio.get_event_loop()
        network_pdfs = []
        missing_networks = []

        for network_key in network_keys:
            pdf_path = await self.template_service.download_to_temp(
                network_key, company_hint=company_hint, format="pdf"
            )

            if not pdf_path:
                # Track missing network but continue with others
                self.logger.warning(f"[PROCESSOR] PDF template not found for network: {network_key}")
                missing_networks.append(network_key)
                continue

            # Extract content pages (strip first and last for intro/outro)
            reader = PdfReader(pdf_path)
            total_pages = len(reader.pages)

            if total_pages > 2:
                # Keep middle pages (content only)
                content_pdf = self._extract_pages_from_pdf(pdf_path, list(range(1, total_pages - 1)))
            elif total_pages == 2:
                # Just 2 pages, keep the second (likely content)
                content_pdf = self._extract_pages_from_pdf(pdf_path, [1])
            else:
                # Single page, keep it
                content_pdf = pdf_path
                pdf_path = None  # Don't delete below

            if pdf_path:
                try:
                    os.unlink(pdf_path)
                except OSError:
                    pass

            network_pdfs.append(content_pdf)

        # If no networks had PDFs, return None to trigger PPTX fallback
        if not network_pdfs:
            return None, missing_networks

        # Merge all network PDFs into one
        if len(network_pdfs) == 1:
            return network_pdfs[0], missing_networks

        from generators.pdf import merge_pdfs
        merged_pdf = await loop.run_in_executor(None, merge_pdfs, network_pdfs)

        # Clean up individual network PDFs
        for pdf in network_pdfs:
            try:
                os.unlink(pdf)
            except OSError:
                pass

        return merged_pdf, missing_networks

    async def _process_package_networks_pptx(
        self,
        network_keys: list[str],
        proposal: dict,
        company_hint: str | None = None,
        already_missing: list[str] | None = None
    ) -> tuple[str | None, list[str]]:
        """
        Process all networks in a package using PPTX fallback strategy.

        Downloads PPTX for each network, converts to PDF (strips intro/outro),
        and merges them into a single PDF. Continues with available networks
        if some are missing.

        Args:
            network_keys: List of network keys in the package
            proposal: The original proposal data
            company_hint: Company schema hint for template lookup
            already_missing: Networks already known to be missing (from PDF attempt)

        Returns:
            Tuple of (merged_pdf_path, missing_networks):
            - merged_pdf_path: Path to merged PDF, or None if ALL networks unavailable
            - missing_networks: List of network keys that had no templates
        """
        from generators.pdf import remove_slides_and_convert_to_pdf, merge_pdfs

        network_pdfs = []
        missing_networks = list(already_missing) if already_missing else []

        for network_key in network_keys:
            # Skip networks already known to be missing
            if network_key in missing_networks:
                continue

            pptx_path = await self.template_service.download_to_temp(
                network_key, company_hint=company_hint, format="pptx"
            )

            if not pptx_path:
                # Track missing network but continue with others
                self.logger.warning(f"[PROCESSOR] PPTX template not found for network: {network_key}")
                missing_networks.append(network_key)
                continue

            # Convert to PDF, removing intro (first) and outro (last) slides
            pdf_path = await remove_slides_and_convert_to_pdf(str(pptx_path), remove_first=True, remove_last=True)

            try:
                os.unlink(pptx_path)
            except OSError:
                pass

            network_pdfs.append(pdf_path)

        # If no networks had templates, return None
        if not network_pdfs:
            return None, missing_networks

        # Merge all network PDFs into one
        if len(network_pdfs) == 1:
            return network_pdfs[0], missing_networks

        loop = asyncio.get_event_loop()
        merged_pdf = await loop.run_in_executor(None, merge_pdfs, network_pdfs)

        # Clean up individual network PDFs
        for pdf in network_pdfs:
            try:
                os.unlink(pdf)
            except OSError:
                pass

        return merged_pdf, missing_networks

    @staticmethod
    def _generate_timestamp_code() -> str:
        """
        Generate compact timestamp code for filenames (HHMMDMMYY).

        Format: HH (hour) MM (minute) D (day) M (month) YY (year)
        Example: 0729072511 = 07:29 AM on July 25, 2011
        Uses UAE timezone (UTC+4)

        Returns:
            Timestamp code string
        """
        uae_tz = timezone(timedelta(hours=4))
        now = datetime.now(uae_tz)
        return f"{now.strftime('%H%M')}{now.day}{now.month}{now.strftime('%y')}"

    @staticmethod
    def _extract_pages_from_pdf(pdf_path: str, pages: list[int]) -> str:
        """
        Extract specific pages from PDF.

        Args:
            pdf_path: Source PDF path
            pages: List of page numbers (0-indexed)

        Returns:
            Path to new PDF with extracted pages
        """
        reader = PdfReader(pdf_path)
        from pypdf import PdfWriter
        writer = PdfWriter()

        for page_num in pages:
            if page_num < len(reader.pages):
                writer.add_page(reader.pages[page_num])

        output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
        output_file.close()

        with open(output_file.name, 'wb') as f:
            writer.write(f)

        return output_file.name

    async def _create_intro_outro_slides(
        self,
        intro_outro_info: dict[str, Any]
    ) -> tuple[str, str]:
        """
        Create intro and outro slide PDFs.

        Strategy:
        1. Try to use pre-made PDF from Supabase Storage
        2. Fall back to PowerPoint extraction if not found

        Args:
            intro_outro_info: Location info dict from IntroOutroHandler

        Returns:
            Tuple of (intro_pdf_path, outro_pdf_path)
        """
        intro_outro_start = time.time()
        series = intro_outro_info.get('series', '')
        location_key = intro_outro_info.get('key', '')
        display_name = intro_outro_info.get('metadata', {}).get('display_name', location_key)

        self.logger.info(f"[TIMING] Intro/outro creation START - {display_name}")
        self.logger.info(f"[PROCESSOR] Series: '{series}'")

        # Determine which pre-made PDF to use
        pdf_name = None
        is_landmark = intro_outro_info.get('is_landmark', False)
        is_non_landmark = intro_outro_info.get('is_non_landmark', False)

        if is_landmark or series == 'The Landmark Series':
            pdf_name = "landmark_series"
            self.logger.info("[PROCESSOR] LANDMARK SERIES - looking for pre-made PDF...")
        elif is_non_landmark:
            pdf_name = "rest"
            self.logger.info("[PROCESSOR] NON-LANDMARK - using rest.pdf...")
        elif 'Digital Icons' in series:
            pdf_name = "digital_icons"
            self.logger.info("[PROCESSOR] DIGITAL ICONS - looking for pre-made PDF...")
        else:
            self.logger.info(f"[PROCESSOR] No pre-made PDF mapping for series '{series}'")

        # Try to download pre-made PDF from storage
        # Extract company_hint from metadata for O(1) lookup
        metadata = intro_outro_info.get('metadata', {})
        company_hint = metadata.get('company_schema') or metadata.get('company')

        if pdf_name:
            # Check cache first for extracted slides
            if pdf_name in self._extracted_slides_cache:
                cached_intro, cached_outro = self._extracted_slides_cache[pdf_name]
                if os.path.exists(cached_intro) and os.path.exists(cached_outro):
                    self.logger.info(f"[PROCESSOR] Cache hit for extracted slides: {pdf_name}")
                    return cached_intro, cached_outro
                else:
                    # Cached files were deleted, remove from cache
                    del self._extracted_slides_cache[pdf_name]

            t0 = time.time()
            pdf_temp_path = await self.template_service.download_intro_outro_to_temp(
                pdf_name, company_hint=company_hint
            )
            self.logger.info(f"[TIMING] Intro/outro PDF download: {(time.time() - t0)*1000:.0f}ms")
            if pdf_temp_path:
                t0 = time.time()
                intro_pdf = self._extract_pages_from_pdf(pdf_temp_path, [0])
                reader = PdfReader(pdf_temp_path)
                last_page = len(reader.pages) - 1
                outro_pdf = self._extract_pages_from_pdf(pdf_temp_path, [last_page])
                self.logger.info(f"[TIMING] Intro/outro extraction: {(time.time() - t0)*1000:.0f}ms")
                # Clean up the temp PDF (but keep extracted pages)
                try:
                    os.unlink(pdf_temp_path)
                except OSError:
                    pass
                # Cache the extracted slides
                self._extracted_slides_cache[pdf_name] = (intro_pdf, outro_pdf)
                self.logger.info(f"[TIMING] Intro/outro creation DONE: {(time.time() - intro_outro_start)*1000:.0f}ms")
                return intro_pdf, outro_pdf
            else:
                self.logger.info(f"[PROCESSOR] PRE-MADE PDF NOT FOUND: {pdf_name}")

        # Fall back to PowerPoint extraction
        self.logger.info("[PROCESSOR] FALLING BACK to PowerPoint extraction")

        # Download template from storage (company_hint already extracted above)
        template_path = await self.template_service.download_to_temp(
            location_key, company_hint=company_hint
        )
        if not template_path:
            # Last resort: try the path from intro_outro_info
            template_path = intro_outro_info.get('template_path')
            if not template_path or not os.path.exists(template_path):
                raise FileNotFoundError(f"Template not found for {location_key}")

        self.logger.info(f"[PROCESSOR] Using PowerPoint template: {template_path}")

        loop = asyncio.get_event_loop()

        # Create intro (first slide only)
        intro_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        intro_pptx.close()
        shutil.copy2(template_path, intro_pptx.name)

        pres = Presentation(intro_pptx.name)
        xml_slides = pres.slides._sldIdLst
        slides_to_remove = list(xml_slides)[1:]
        for slide_id in slides_to_remove:
            xml_slides.remove(slide_id)
        pres.save(intro_pptx.name)

        intro_pdf = await loop.run_in_executor(None, convert_pptx_to_pdf, intro_pptx.name)

        # Create outro (last slide only)
        outro_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        outro_pptx.close()
        shutil.copy2(template_path, outro_pptx.name)

        pres = Presentation(outro_pptx.name)
        xml_slides = pres.slides._sldIdLst
        slides_to_remove = list(xml_slides)[:-1]
        for slide_id in slides_to_remove:
            xml_slides.remove(slide_id)
        pres.save(outro_pptx.name)

        outro_pdf = await loop.run_in_executor(None, convert_pptx_to_pdf, outro_pptx.name)

        # Clean up temp files
        try:
            os.unlink(intro_pptx.name)
            os.unlink(outro_pptx.name)
            # Clean up downloaded template if it was from storage
            if template_path != intro_outro_info.get('template_path'):
                os.unlink(template_path)
        except OSError as e:
            self.logger.warning(f"[PROCESSOR] Failed to cleanup temp files: {e}")

        return intro_pdf, outro_pdf

    async def process_separate(
        self,
        proposals_data: list[dict[str, Any]],
        submitted_by: str = "",
        client_name: str = "",
        payment_terms: str = "100% upfront",
        currency: str = None
    ) -> dict[str, Any]:
        """
        Process separate proposals (one PDF per location or merged).

        Uses PDF-first strategy with PPTX fallback for faster processing.

        Args:
            proposals_data: List of proposal dicts
            submitted_by: User who submitted
            client_name: Client name
            payment_terms: Payment terms text
            currency: Currency code (e.g., 'USD', 'EUR')

        Returns:
            Dict with success status and file paths:
            - For single proposal: {success, is_single, pptx_path, pdf_path, location, ...}
            - For multiple: {success, is_single: False, individual_files, merged_pdf_path, ...}

        Example:
            >>> processor = ProposalProcessor(validator, renderer, intro_outro)
            >>> result = await processor.process_separate(
            ...     proposals_data,
            ...     submitted_by="user@example.com",
            ...     client_name="ABC Corp",
            ...     payment_terms="100% upfront",
            ...     currency="AED"
            ... )
        """
        process_start = time.time()
        self.logger.info("[TIMING] ========== SEPARATE PROPOSAL START ==========")
        self.logger.info(f"[TIMING] Proposals: {len(proposals_data)}, Client: {client_name}")

        # Validate proposals (async)
        t0 = time.time()
        validated_proposals, errors = await self.validator.validate_proposals(proposals_data)
        self.logger.info(f"[TIMING] Validation: {(time.time() - t0)*1000:.0f}ms")
        if errors:
            return {"success": False, "errors": errors}

        is_single = len(validated_proposals) == 1
        loop = asyncio.get_event_loop()
        total_proposals = len(validated_proposals)

        # Get intro/outro info for multiple proposals
        intro_outro_info = None
        if len(validated_proposals) > 1:
            intro_outro_info = self.intro_outro_handler.get_intro_outro_location(validated_proposals)

        # Build financial data helper
        def build_financial_data(proposal: dict) -> dict:
            # Get location metadata for display name
            location_meta = proposal.get("location_metadata", {})
            financial_data = {
                "location": proposal["location"],
                "durations": proposal["durations"],
                "net_rates": proposal.get("net_rates", []),
                "spots": proposal.get("spots", 1),
                "client_name": client_name,
                "payment_terms": payment_terms,
                "location_metadata": location_meta,  # Pass metadata for display name
            }
            if "start_dates" in proposal and proposal["start_dates"]:
                financial_data["start_dates"] = proposal["start_dates"]
            else:
                financial_data["start_date"] = proposal.get("start_date", "1st December 2025")
            if "end_date" in proposal:
                financial_data["end_date"] = proposal["end_date"]
            if "production_fee" in proposal:
                financial_data["production_fee"] = proposal["production_fee"]
            return financial_data

        # PDF-first strategy for processing proposals
        async def process_separate_pdf_first(idx: int, proposal: dict) -> dict | None:
            """PDF-First: Download PDF template, create standalone financial slide."""
            location_key = proposal.get("location", "").lower().strip()
            company_hint = proposal.get("company_schema")
            single_start = time.time()

            # Try to download PDF template
            t0 = time.time()
            pdf_template_path = await self.template_service.download_to_temp(
                location_key, company_hint=company_hint, format="pdf"
            )

            if not pdf_template_path:
                self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PDF not available, will use fallback")
                return None  # Signal to use fallback

            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PDF download: {(time.time() - t0)*1000:.0f}ms (PDF-FIRST)")

            # Build financial data and create standalone financial slide
            financial_data = build_financial_data(proposal)
            t0 = time.time()
            financial_pptx_path, vat_amounts, total_amounts = await loop.run_in_executor(
                None,
                self.renderer.create_standalone_financial_slide,
                financial_data,
                currency,
                None,  # slide_width (use default)
                None,  # slide_height (use default)
            )
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - Financial slide created: {(time.time() - t0)*1000:.0f}ms")

            # Convert financial slide to PDF
            t0 = time.time()
            financial_pdf = await loop.run_in_executor(None, convert_pptx_to_pdf, financial_pptx_path)
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - Financial PDF conversion: {(time.time() - t0)*1000:.0f}ms")

            # Clean up financial PPTX
            try:
                os.unlink(financial_pptx_path)
            except OSError:
                pass

            result = {
                "location": proposal["location"].title(),
                "filename": f"{proposal['location'].title()}_Proposal.pptx",
                "totals": total_amounts,
                "idx": idx,
                "pdf_template_path": pdf_template_path,
                "financial_pdf": financial_pdf,
                "is_pdf_first": True,
            }

            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - TOTAL: {(time.time() - single_start)*1000:.0f}ms (PDF-FIRST)")
            return result

        async def process_separate_pptx_fallback(idx: int, proposal: dict) -> dict:
            """PPTX Fallback: Original flow using PPTX templates."""
            location_key = proposal.get("location", "").lower().strip()
            company_hint = proposal.get("company_schema")
            single_start = time.time()
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - START (PPTX FALLBACK)")

            # Download PPTX template
            t0 = time.time()
            template_path = await self.template_service.download_to_temp(
                location_key, company_hint=company_hint, format="pptx"
            )
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PPTX download: {(time.time() - t0)*1000:.0f}ms")

            if not template_path:
                raise FileNotFoundError(f"Template not found for {location_key}")

            # Build financial data
            financial_data = build_financial_data(proposal)

            # Render PPTX with financial slide
            t0 = time.time()
            pptx_path, vat_amounts, total_amounts = await loop.run_in_executor(
                None,
                self.renderer.create_proposal_with_template,
                str(template_path),
                financial_data,
                currency
            )
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PPTX rendering: {(time.time() - t0)*1000:.0f}ms")

            # Clean up downloaded template
            try:
                os.unlink(template_path)
            except OSError:
                pass

            result = {
                "path": pptx_path,
                "location": proposal["location"].title(),
                "filename": f"{proposal['location'].title()}_Proposal.pptx",
                "totals": total_amounts,
                "idx": idx,
                "is_pdf_first": False,
            }

            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - TOTAL: {(time.time() - single_start)*1000:.0f}ms (PPTX FALLBACK)")
            return result

        async def process_proposal(idx: int, proposal: dict) -> dict:
            """Process a proposal: try PDF-first, fall back to PPTX if needed."""
            location_key = proposal.get("location", "").lower().strip()
            company_hint = proposal.get("company_schema")

            # Check if this is a package (bundle of multiple networks)
            is_package, network_keys = await self._check_if_package(location_key)

            if is_package:
                self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PACKAGE with {len(network_keys)} networks: {network_keys}")

                # Guardrail: Package must have at least one network
                if not network_keys:
                    self.logger.error(f"[PROCESSOR] Package '{location_key}' has no networks - cannot generate proposal")
                    return {
                        "idx": idx,
                        "error": f"Package '{location_key}' has no networks configured. Contact admin.",
                        "location": location_key,
                    }

                single_start = time.time()
                missing_networks = []

                # Try PDF-first for package
                merged_content_pdf, pdf_missing = await self._process_package_networks_pdf(
                    network_keys, proposal, company_hint
                )

                if merged_content_pdf is None:
                    # Fallback to PPTX, passing networks already known to be missing
                    self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - Package PDF not available, using PPTX fallback")
                    merged_content_pdf, pptx_missing = await self._process_package_networks_pptx(
                        network_keys, proposal, company_hint, already_missing=pdf_missing
                    )
                    missing_networks = pptx_missing
                else:
                    missing_networks = pdf_missing

                # If ALL networks are missing, return error
                if merged_content_pdf is None:
                    missing_str = ", ".join(missing_networks)
                    self.logger.error(f"[PROCESSOR] Package '{location_key}' has no available templates - all networks missing: {missing_str}")
                    return {
                        "idx": idx,
                        "error": f"No templates available for package '{location_key}'. Missing networks: {missing_str}. Contact administrator.",
                        "location": location_key,
                    }

                # Log if some networks were missing
                if missing_networks:
                    self.logger.warning(f"[PROCESSOR] Package '{location_key}' generated with partial content. Missing networks: {missing_networks}")

                # Build financial data and create standalone financial slide
                financial_data = build_financial_data(proposal)
                financial_pptx_path, vat_amounts, total_amounts = await loop.run_in_executor(
                    None,
                    self.renderer.create_standalone_financial_slide,
                    financial_data,
                    currency,
                    None,
                    None,
                )

                # Convert financial slide to PDF
                financial_pdf = await loop.run_in_executor(None, convert_pptx_to_pdf, financial_pptx_path)

                # Clean up financial PPTX
                try:
                    os.unlink(financial_pptx_path)
                except OSError:
                    pass

                self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PACKAGE TOTAL: {(time.time() - single_start)*1000:.0f}ms")

                return {
                    "location": proposal["location"].title(),
                    "filename": f"{proposal['location'].title()}_Proposal.pptx",
                    "totals": total_amounts,
                    "idx": idx,
                    "pdf_template_path": merged_content_pdf,
                    "financial_pdf": financial_pdf,
                    "is_pdf_first": True,
                    "is_package": True,
                    "network_keys": network_keys,
                    "missing_networks": missing_networks,  # Track for warning message
                }

            # Normal flow for non-packages
            result = await process_separate_pdf_first(idx, proposal)
            if result is None:
                result = await process_separate_pptx_fallback(idx, proposal)
            return result

        # Process all in parallel
        self.logger.info(f"[TIMING] Starting parallel processing of {total_proposals} proposals...")
        t0 = time.time()
        tasks = [process_proposal(idx, p) for idx, p in enumerate(validated_proposals)]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        self.logger.info(f"[TIMING] Parallel processing complete: {(time.time() - t0)*1000:.0f}ms")

        # Check for errors
        for idx, result in enumerate(results):
            if isinstance(result, Exception):
                return {"success": False, "error": f"Error processing proposal {idx + 1}: {str(result)}"}
            # Check for error dicts (e.g., from package with no networks)
            if isinstance(result, dict) and result.get("error"):
                return {"success": False, "error": result["error"]}

        # Sort by original index
        sorted_results = sorted(results, key=lambda x: x["idx"])

        # Handle single proposal
        if is_single:
            result = sorted_results[0]
            total_str = result["totals"][0] if result["totals"] else "AED 0"
            timestamp_code = self._generate_timestamp_code()
            client_prefix = client_name.replace(" ", "_") if client_name else "Client"

            if result.get("is_pdf_first"):
                # PDF-first flow: insert financial slide BEFORE outro (last slide)
                t0 = time.time()

                # Check if this is a package - packages need special handling
                # because their content PDF has NO intro/outro (stripped from each network)
                if result.get("is_package"):
                    # Get intro/outro from first network in the package
                    network_keys = result.get("network_keys", [])
                    if network_keys:
                        first_network = network_keys[0]
                        synthetic_proposal = {"location": first_network}
                        package_intro_outro_info = self.intro_outro_handler.get_intro_outro_location([synthetic_proposal])

                        if package_intro_outro_info:
                            self.logger.info(f"[INTRO_OUTRO] Package using intro/outro from network: {first_network}")
                            intro_pdf, outro_pdf = await self._create_intro_outro_slides(package_intro_outro_info)
                            # Merge: intro + content + financial + outro
                            pdf_path = await loop.run_in_executor(
                                None, merge_pdfs, [intro_pdf, result["pdf_template_path"], result["financial_pdf"], outro_pdf]
                            )
                            # Clean up intro/outro PDFs
                            try:
                                os.unlink(intro_pdf)
                                os.unlink(outro_pdf)
                            except OSError:
                                pass
                        else:
                            # No intro/outro found - just merge content + financial
                            self.logger.warning(f"[INTRO_OUTRO] No intro/outro found for package, merging without")
                            pdf_path = await loop.run_in_executor(
                                None, merge_pdfs, [result["pdf_template_path"], result["financial_pdf"]]
                            )
                    else:
                        # No network keys - just merge content + financial
                        pdf_path = await loop.run_in_executor(
                            None, merge_pdfs, [result["pdf_template_path"], result["financial_pdf"]]
                        )
                else:
                    # Regular location - last page is outro
                    reader = PdfReader(result["pdf_template_path"])
                    total_pages = len(reader.pages)

                    if total_pages > 1:
                        # Extract content (all pages except last/outro)
                        content_pdf = self._extract_pages_from_pdf(
                            result["pdf_template_path"], list(range(0, total_pages - 1))
                        )
                        # Extract outro (last page)
                        outro_pdf = self._extract_pages_from_pdf(
                            result["pdf_template_path"], [total_pages - 1]
                        )
                        # Merge: content + financial + outro
                        pdf_path = await loop.run_in_executor(
                            None, merge_pdfs, [content_pdf, result["financial_pdf"], outro_pdf]
                        )
                        # Clean up extracted PDFs
                        try:
                            os.unlink(content_pdf)
                            os.unlink(outro_pdf)
                        except OSError:
                            pass
                    else:
                        # Single page template - just append financial
                        pdf_path = await loop.run_in_executor(
                            None, merge_pdfs, [result["pdf_template_path"], result["financial_pdf"]]
                        )

                self.logger.info(f"[TIMING] Single proposal PDF merge: {(time.time() - t0)*1000:.0f}ms")

                # Clean up temp files
                try:
                    os.unlink(result["pdf_template_path"])
                    os.unlink(result["financial_pdf"])
                except OSError:
                    pass

                db.log_proposal(
                    submitted_by=submitted_by,
                    client_name=client_name,
                    package_type="separate",
                    locations=result["location"],
                    total_amount=total_str,
                )

                total_time = (time.time() - process_start) * 1000
                self.logger.info(f"[TIMING] ========== SEPARATE PROPOSAL COMPLETE ==========")
                self.logger.info(f"[TIMING] Total time: {total_time:.0f}ms ({total_time/1000:.1f}s)")

                # Build warning message if some networks were missing
                warning = None
                missing = result.get("missing_networks", [])
                if missing:
                    missing_str = ", ".join(missing)
                    warning = f"⚠️ Note: Some networks in this package are missing templates ({missing_str}). The proposal was generated with available content only. Please contact administrator to add missing templates."

                return {
                    "success": True,
                    "is_single": True,
                    "pptx_path": None,  # No PPTX in PDF-first flow
                    "pdf_path": pdf_path,
                    "location": result["location"],
                    "pptx_filename": None,
                    "pdf_filename": f"{client_prefix}_{timestamp_code}.pdf",
                    "warning": warning,
                }
            else:
                # PPTX fallback flow: convert PPTX to PDF
                t0 = time.time()
                pdf_path = await loop.run_in_executor(None, convert_pptx_to_pdf, result["path"])
                self.logger.info(f"[TIMING] Single proposal PDF conversion: {(time.time() - t0)*1000:.0f}ms")

                db.log_proposal(
                    submitted_by=submitted_by,
                    client_name=client_name,
                    package_type="separate",
                    locations=result["location"],
                    total_amount=total_str,
                )

                total_time = (time.time() - process_start) * 1000
                self.logger.info(f"[TIMING] ========== SEPARATE PROPOSAL COMPLETE ==========")
                self.logger.info(f"[TIMING] Total time: {total_time:.0f}ms ({total_time/1000:.1f}s)")

                return {
                    "success": True,
                    "is_single": True,
                    "pptx_path": result["path"],
                    "pdf_path": pdf_path,
                    "location": result["location"],
                    "pptx_filename": result["filename"],
                    "pdf_filename": f"{client_prefix}_{timestamp_code}.pdf",
                }

        # Handle multiple proposals
        individual_files = []
        pdf_files = []
        locations = []

        for result in sorted_results:
            if result.get("is_pdf_first"):
                # PDF-first flow: extract middle pages, merge with financial slide
                t0 = time.time()
                reader = PdfReader(result["pdf_template_path"])
                total_pages = len(reader.pages)

                # For packages, content is already stripped (no intro/outro) - keep all pages
                # For regular locations, skip first/last for intro/outro
                if result.get("is_package"):
                    # Package content is already clean - keep all pages
                    pages_to_keep = list(range(total_pages))
                elif intro_outro_info:
                    pages_to_keep = list(range(1, total_pages - 1)) if total_pages > 2 else list(range(total_pages))
                else:
                    idx = result["idx"]
                    if idx == 0:
                        pages_to_keep = list(range(0, total_pages - 1))
                    elif idx < total_proposals - 1:
                        pages_to_keep = list(range(1, total_pages - 1))
                    else:
                        pages_to_keep = list(range(1, total_pages))

                if pages_to_keep:
                    template_content_pdf = self._extract_pages_from_pdf(result["pdf_template_path"], pages_to_keep)
                else:
                    template_content_pdf = result["pdf_template_path"]

                # Merge template content with financial slide
                location_pdf = await loop.run_in_executor(
                    None, merge_pdfs, [template_content_pdf, result["financial_pdf"]]
                )
                self.logger.info(f"[TIMING] Location PDF merge ({result['location']}): {(time.time() - t0)*1000:.0f}ms")

                # Clean up temp files
                try:
                    if template_content_pdf != result["pdf_template_path"]:
                        os.unlink(template_content_pdf)
                    os.unlink(result["pdf_template_path"])
                    os.unlink(result["financial_pdf"])
                except OSError:
                    pass

                individual_files.append({
                    "path": None,  # No PPTX in PDF-first flow
                    "location": result["location"],
                    "filename": result["filename"],
                    "totals": result["totals"],
                })
                pdf_files.append(location_pdf)
            else:
                # PPTX fallback flow: remove slides and convert
                t0 = time.time()
                if intro_outro_info:
                    remove_first = True
                    remove_last = True
                else:
                    idx = result["idx"]
                    remove_first = False
                    remove_last = False
                    if idx == 0:
                        remove_last = True
                    elif idx < total_proposals - 1:
                        remove_first = True
                        remove_last = True
                    else:
                        remove_first = True

                pdf_path = await remove_slides_and_convert_to_pdf(result["path"], remove_first, remove_last)
                self.logger.info(f"[TIMING] Location PDF conversion ({result['location']}): {(time.time() - t0)*1000:.0f}ms")

                individual_files.append({
                    "path": result["path"],
                    "location": result["location"],
                    "filename": result["filename"],
                    "totals": result["totals"],
                })
                pdf_files.append(pdf_path)

            locations.append(result["location"])

        # Add intro/outro slides
        if intro_outro_info:
            t0 = time.time()
            intro_pdf, outro_pdf = await self._create_intro_outro_slides(intro_outro_info)
            self.logger.info(f"[TIMING] Intro/outro slides: {(time.time() - t0)*1000:.0f}ms")
            pdf_files.insert(0, intro_pdf)
            pdf_files.append(outro_pdf)

        # Merge PDFs
        t0 = time.time()
        merged_pdf = await loop.run_in_executor(None, merge_pdfs, pdf_files)
        self.logger.info(f"[TIMING] Final PDF merge ({len(pdf_files)} files): {(time.time() - t0)*1000:.0f}ms")

        # Clean up temp PDFs
        for pdf_file in pdf_files:
            try:
                os.unlink(pdf_file)
            except Exception as e:
                self.logger.warning(f"Failed to clean up PDF: {pdf_file} - {e}")

        # Log to database
        first_totals = [f.get("totals", ["AED 0"])[0] for f in individual_files]
        summary_total = ", ".join(first_totals)

        db.log_proposal(
            submitted_by=submitted_by,
            client_name=client_name,
            package_type="separate",
            locations=", ".join(locations),
            total_amount=summary_total,
        )

        timestamp_code = self._generate_timestamp_code()
        client_prefix = client_name.replace(" ", "_") if client_name else "Client"

        total_time = (time.time() - process_start) * 1000
        self.logger.info(f"[TIMING] ========== SEPARATE PROPOSAL COMPLETE ==========")
        self.logger.info(f"[TIMING] Total time: {total_time:.0f}ms ({total_time/1000:.1f}s)")
        self.logger.info(f"[TIMING] Locations: {', '.join(locations)}")

        return {
            "success": True,
            "is_single": False,
            "individual_files": individual_files,
            "merged_pdf_path": merged_pdf,
            "locations": ", ".join(locations),
            "merged_pdf_filename": f"{client_prefix}_{timestamp_code}.pdf",
        }

    async def process_combined(
        self,
        proposals_data: list[dict[str, Any]],
        combined_net_rate: str,
        submitted_by: str = "",
        client_name: str = "",
        payment_terms: str = "100% upfront",
        currency: str = None
    ) -> dict[str, Any]:
        """
        Process combined package (one PDF with all locations).

        Args:
            proposals_data: List of proposal dicts
            combined_net_rate: Combined rate for package
            submitted_by: User who submitted
            client_name: Client name
            payment_terms: Payment terms text
            currency: Currency code

        Returns:
            Dict with success status and file paths:
            {success, is_combined, pdf_path, locations, pdf_filename}

        Example:
            >>> result = await processor.process_combined(
            ...     proposals_data,
            ...     "50000",
            ...     submitted_by="user@example.com",
            ...     client_name="ABC Corp",
            ...     payment_terms="100% upfront",
            ...     currency="AED"
            ... )
        """
        process_start = time.time()
        self.logger.info("[TIMING] ========== COMBINED PROPOSAL START ==========")
        self.logger.info(f"[TIMING] Proposals: {len(proposals_data)}, Client: {client_name}")

        # Validate proposals (async)
        t0 = time.time()
        validated_proposals, errors = await self.validator.validate_combined_package(
            proposals_data,
            combined_net_rate
        )
        self.logger.info(f"[TIMING] Validation: {(time.time() - t0)*1000:.0f}ms")
        if errors:
            return {"success": False, "errors": errors}

        loop = asyncio.get_event_loop()

        # Get intro/outro info
        t0 = time.time()
        intro_outro_info = self.intro_outro_handler.get_intro_outro_location(validated_proposals)
        self.logger.info(f"[TIMING] Intro/outro info: {(time.time() - t0)*1000:.0f}ms")

        # If intro_outro_info is None, check if any proposal is a package
        # If so, expand it and use the first network for intro/outro
        if intro_outro_info is None:
            for proposal in validated_proposals:
                location_key = proposal.get("location", "").lower().strip()
                is_package, network_keys = await self._check_if_package(location_key)
                if is_package and network_keys:
                    self.logger.info(f"[INTRO_OUTRO] Package '{location_key}' detected, using first network '{network_keys[0]}' for intro/outro")
                    # Create a synthetic proposal with the first network for intro/outro lookup
                    synthetic_proposal = {"location": network_keys[0]}
                    intro_outro_info = self.intro_outro_handler.get_intro_outro_location([synthetic_proposal])
                    if intro_outro_info:
                        self.logger.info(f"[INTRO_OUTRO] Found intro/outro from package network: {intro_outro_info.get('key')}")
                    break

        # Deduplicate deck slides (same location can have multiple durations)
        seen_locations = set()
        unique_proposals_for_deck = []

        for proposal in validated_proposals:
            location_key = proposal["location"]
            if location_key not in seen_locations:
                seen_locations.add(location_key)
                unique_proposals_for_deck.append(proposal)

        self.logger.info(f"[TIMING] Unique locations: {len(unique_proposals_for_deck)} (from {len(validated_proposals)} proposals)")

        # Process each unique location (in parallel)
        total_proposals = len(unique_proposals_for_deck)

        async def process_combined_single_pdf_first(idx: int, proposal: dict) -> dict:
            """
            PDF-First Strategy: Download PDF templates, build standalone financial slide.

            This is the fast path (~1-3s per location vs 10-30s for PPTX).
            Returns None if PDF not available (triggers fallback).
            """
            location_key = proposal.get("location", "").lower().strip()
            company_hint = proposal.get("company_schema")
            single_start = time.time()

            # Try to download PDF template
            t0 = time.time()
            pdf_template_path = await self.template_service.download_to_temp(
                location_key, company_hint=company_hint, format="pdf"
            )

            if not pdf_template_path:
                self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PDF not available, will use fallback")
                return None  # Signal to use fallback

            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PDF download: {(time.time() - t0)*1000:.0f}ms (PDF-FIRST)")

            # Extract middle pages (skip first and last for intro/outro)
            t0 = time.time()
            reader = PdfReader(pdf_template_path)
            total_pages = len(reader.pages)

            if intro_outro_info:
                pages_to_keep = list(range(1, total_pages - 1)) if total_pages > 2 else list(range(total_pages))
            else:
                if idx == 0:
                    pages_to_keep = list(range(0, total_pages - 1))
                elif idx < total_proposals - 1:
                    pages_to_keep = list(range(1, total_pages - 1))
                else:
                    pages_to_keep = list(range(1, total_pages))

            if pages_to_keep:
                pdf_path = self._extract_pages_from_pdf(pdf_template_path, pages_to_keep)
                try:
                    os.unlink(pdf_template_path)
                except OSError:
                    pass
            else:
                pdf_path = pdf_template_path

            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PDF extraction: {(time.time() - t0)*1000:.0f}ms")
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - TOTAL: {(time.time() - single_start)*1000:.0f}ms (PDF-FIRST)")

            return {"pdf_path": pdf_path, "idx": idx}

        async def process_combined_single_pptx_fallback(idx: int, proposal: dict) -> dict:
            """
            PPTX Fallback: Original flow using PPTX templates.

            Used when PDF template is not available.
            """
            location_key = proposal.get("location", "").lower().strip()
            company_hint = proposal.get("company_schema")
            single_start = time.time()
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - START (PPTX FALLBACK)")

            t0 = time.time()
            template_path = await self.template_service.download_to_temp(
                location_key, company_hint=company_hint, format="pptx"
            )
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PPTX download: {(time.time() - t0)*1000:.0f}ms")

            if not template_path:
                raise FileNotFoundError(f"Template not found for {location_key}")

            pptx_path = str(template_path)

            # Determine which slides to remove
            if intro_outro_info:
                remove_first = True
                remove_last = True
            else:
                remove_first = False
                remove_last = False
                if idx == 0:
                    remove_last = True
                elif idx < total_proposals - 1:
                    remove_first = True
                    remove_last = True
                else:
                    remove_first = True

            t0 = time.time()
            pdf_path = await remove_slides_and_convert_to_pdf(pptx_path, remove_first, remove_last)
            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PDF conversion: {(time.time() - t0)*1000:.0f}ms")

            try:
                os.unlink(pptx_path)
            except OSError:
                pass

            self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - TOTAL: {(time.time() - single_start)*1000:.0f}ms (PPTX FALLBACK)")
            return {"pdf_path": pdf_path, "idx": idx}

        async def process_location(idx: int, proposal: dict) -> dict:
            """Process a location: try PDF-first, fall back to PPTX if needed."""
            location_key = proposal.get("location", "").lower().strip()
            company_hint = proposal.get("company_schema")

            # Check if this is a package (bundle of multiple networks)
            is_package, network_keys = await self._check_if_package(location_key)

            if is_package:
                self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PACKAGE with {len(network_keys)} networks: {network_keys}")
                single_start = time.time()
                missing_networks = []

                # Try PDF-first for package (already strips intro/outro from each network)
                merged_content_pdf, pdf_missing = await self._process_package_networks_pdf(
                    network_keys, proposal, company_hint
                )

                if merged_content_pdf is None:
                    # Fallback to PPTX, passing networks already known to be missing
                    self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - Package PDF not available, using PPTX fallback")
                    merged_content_pdf, pptx_missing = await self._process_package_networks_pptx(
                        network_keys, proposal, company_hint, already_missing=pdf_missing
                    )
                    missing_networks = pptx_missing
                else:
                    missing_networks = pdf_missing

                # If ALL networks are missing, return error
                if merged_content_pdf is None:
                    missing_str = ", ".join(missing_networks)
                    self.logger.error(f"[PROCESSOR] Package '{location_key}' has no available templates - all networks missing: {missing_str}")
                    return {
                        "idx": idx,
                        "error": f"No templates available for package '{location_key}'. Missing networks: {missing_str}. Contact administrator.",
                    }

                # Log if some networks were missing
                if missing_networks:
                    self.logger.warning(f"[PROCESSOR] Package '{location_key}' generated with partial content. Missing networks: {missing_networks}")

                self.logger.info(f"[TIMING] [{idx+1}/{total_proposals}] {location_key} - PACKAGE TOTAL: {(time.time() - single_start)*1000:.0f}ms")

                return {"pdf_path": merged_content_pdf, "idx": idx, "is_package": True, "network_keys": network_keys, "missing_networks": missing_networks}

            # Normal flow for non-packages
            result = await process_combined_single_pdf_first(idx, proposal)
            if result is None:
                result = await process_combined_single_pptx_fallback(idx, proposal)
            return result

        # Process all locations in parallel (PDF-first with PPTX fallback)
        self.logger.info(f"[TIMING] Starting parallel processing of {total_proposals} locations...")
        t0 = time.time()
        tasks = [process_location(idx, p) for idx, p in enumerate(unique_proposals_for_deck)]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        self.logger.info(f"[TIMING] Parallel processing complete: {(time.time() - t0)*1000:.0f}ms")

        # Check for errors
        for idx, result in enumerate(results):
            if isinstance(result, Exception):
                return {"success": False, "error": f"Error processing proposal {idx + 1}: {str(result)}"}
            # Check for error dicts (e.g., from package with no networks)
            if isinstance(result, dict) and result.get("error"):
                return {"success": False, "error": result["error"]}

        # Sort by original index and extract PDFs
        sorted_results = sorted(results, key=lambda x: x["idx"])
        pdf_files = [r["pdf_path"] for r in sorted_results]

        # Aggregate missing networks from all results
        all_missing_networks = []
        for r in sorted_results:
            missing = r.get("missing_networks", [])
            if missing:
                all_missing_networks.extend(missing)

        # Create standalone combined financial slide
        t0 = time.time()
        financial_pptx_path, total_combined = await loop.run_in_executor(
            None,
            self.renderer.create_standalone_combined_financial_slide,
            validated_proposals,
            combined_net_rate,
            client_name,
            payment_terms,
            currency,
            None,  # slide_width (use default)
            None,  # slide_height (use default)
        )
        self.logger.info(f"[TIMING] Standalone financial slide created: {(time.time() - t0)*1000:.0f}ms")

        # Convert financial slide to PDF
        t0 = time.time()
        financial_pdf_path = await loop.run_in_executor(None, convert_pptx_to_pdf, financial_pptx_path)
        self.logger.info(f"[TIMING] Financial slide PDF conversion: {(time.time() - t0)*1000:.0f}ms")

        # Clean up financial PPTX
        try:
            os.unlink(financial_pptx_path)
        except OSError:
            pass

        # Add financial slide PDF to the list (before outro)
        pdf_files.append(financial_pdf_path)

        # Add intro/outro slides
        if intro_outro_info:
            t0 = time.time()
            intro_pdf, outro_pdf = await self._create_intro_outro_slides(intro_outro_info)
            self.logger.info(f"[TIMING] Intro/outro slides: {(time.time() - t0)*1000:.0f}ms")
            pdf_files.insert(0, intro_pdf)
            pdf_files.append(outro_pdf)

        # Merge PDFs
        t0 = time.time()
        merged_pdf = await loop.run_in_executor(None, merge_pdfs, pdf_files)
        self.logger.info(f"[TIMING] PDF merge ({len(pdf_files)} files): {(time.time() - t0)*1000:.0f}ms")

        # Clean up temp PDFs
        for pdf_file in pdf_files:
            try:
                os.unlink(pdf_file)
            except Exception as e:
                self.logger.warning(f"Failed to clean up PDF: {pdf_file} - {e}")

        # Calculate total if not provided
        if total_combined is None:
            municipality_fee = 520
            total_upload_fees = sum(
                config.UPLOAD_FEES_MAPPING.get(p["location"].lower(), 3000)
                for p in validated_proposals
            )
            net_rate_numeric = float(combined_net_rate.replace("AED", "").replace(",", "").strip())
            subtotal = net_rate_numeric + total_upload_fees + municipality_fee
            vat = subtotal * 0.05
            total_combined = f"AED {subtotal + vat:,.0f}"

        # Log to database
        locations_str = ", ".join([p["location"].title() for p in validated_proposals])

        db.log_proposal(
            submitted_by=submitted_by,
            client_name=client_name,
            package_type="combined",
            locations=locations_str,
            total_amount=total_combined,
        )

        timestamp_code = self._generate_timestamp_code()
        client_prefix = client_name.replace(" ", "_") if client_name else "Client"

        total_time = (time.time() - process_start) * 1000
        self.logger.info(f"[TIMING] ========== COMBINED PROPOSAL COMPLETE ==========")
        self.logger.info(f"[TIMING] Total time: {total_time:.0f}ms ({total_time/1000:.1f}s)")
        self.logger.info(f"[TIMING] Locations: {locations_str}")

        # Build warning message if some networks were missing
        warning = None
        if all_missing_networks:
            missing_str = ", ".join(all_missing_networks)
            warning = f"⚠️ Note: Some networks in this package are missing templates ({missing_str}). The proposal was generated with available content only. Please contact administrator to add missing templates."

        return {
            "success": True,
            "is_combined": True,
            "pptx_path": None,
            "pdf_path": merged_pdf,
            "locations": locations_str,
            "pdf_filename": f"{client_prefix}_{timestamp_code}.pdf",
            "warning": warning,
        }
