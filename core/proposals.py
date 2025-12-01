import os
import asyncio
import tempfile
import shutil
from pathlib import Path
from typing import Dict, Any, List, Tuple, Optional
from datetime import datetime, timezone, timedelta

from pptx import Presentation
from pypdf import PdfReader, PdfWriter

import config
from db.database import db
from generators.pptx import create_financial_proposal_slide, create_combined_financial_proposal_slide
from generators.pdf import convert_pptx_to_pdf, merge_pdfs, remove_slides_and_convert_to_pdf


def _generate_timestamp_code() -> str:
    """Generate a compact timestamp code for filenames (HHMMDMMYY)
    Format: HH (hour) MM (minute) D (day) M (month) YY (year)
    Example: 0729072511 = 07:29 AM on July 25, 2011
    Uses UAE timezone (UTC+4)
    """
    # UAE timezone is UTC+4
    uae_tz = timezone(timedelta(hours=4))
    now = datetime.now(uae_tz)
    return f"{now.strftime('%H%M')}{now.day}{now.month}{now.strftime('%y')}"


def _template_path_for_key(key: str) -> Path:
    mapping = config.get_location_mapping()
    filename = mapping.get(key)
    if not filename:
        raise FileNotFoundError(f"Unknown location '{key}'. Available: {', '.join(config.available_location_names())}")
    return config.TEMPLATES_DIR / filename


def _extract_pages_from_pdf(pdf_path: str, pages: List[int]) -> str:
    """Extract specific pages from a PDF and save to a new PDF file.
    
    Args:
        pdf_path: Path to the source PDF
        pages: List of page numbers to extract (0-indexed)
    
    Returns:
        Path to the new PDF file
    """
    logger = config.logger
    logger.info(f"[EXTRACT_PDF] Extracting pages {pages} from {pdf_path}")
    
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    
    for page_num in pages:
        if page_num < len(reader.pages):
            writer.add_page(reader.pages[page_num])
    
    output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    output_file.close()
    
    with open(output_file.name, 'wb') as f:
        writer.write(f)
    
    logger.info(f"[EXTRACT_PDF] Saved extracted pages to {output_file.name}")
    return output_file.name




def _get_location_info_for_intro_outro(proposals_data: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    """Find a suitable location for intro/outro slides based on series."""
    logger = config.logger

    logger.info(f"[INTRO_OUTRO] 🔍 Searching for suitable location from {len(proposals_data)} proposals")

    # First, look for locations with "The Landmark Series"
    mapping = config.get_location_mapping()
    for idx, proposal in enumerate(proposals_data):
        location = proposal.get("location", "").lower().strip()
        logger.info(f"[INTRO_OUTRO] Checking proposal {idx+1}: location='{location}'")

        # Get the actual key from display name or direct match
        matched_key = config.get_location_key_from_display_name(location)
        if not matched_key:
            # Try old matching logic
            for key in mapping.keys():
                if key in location or location in key:
                    matched_key = key
                    break

        if matched_key:
            location_meta = config.LOCATION_METADATA.get(matched_key, {})
            display_type = location_meta.get('display_type', 'Unknown')
            series = location_meta.get('series', '')
            display_name = location_meta.get('display_name', matched_key)

            logger.info(f"[INTRO_OUTRO] Found location: '{display_name}' (key: {matched_key})")
            logger.info(f"[INTRO_OUTRO]   - Display Type: {display_type}")
            logger.info(f"[INTRO_OUTRO]   - Series: '{series}'")

            if series == 'The Landmark Series':
                logger.info(f"[INTRO_OUTRO] ✅ LANDMARK SERIES FOUND! Using '{display_name}' for intro/outro")
                return {
                    'key': matched_key,
                    'series': series,
                    'template_path': str(config.TEMPLATES_DIR / mapping[matched_key]),
                    'metadata': location_meta,
                    'is_landmark': True
                }

    # If no Landmark Series found, use the first location from proposals
    logger.info(f"[INTRO_OUTRO] ❌ No Landmark Series location found in proposals")
    
    if proposals_data:
        first_location = proposals_data[0].get("location", "").lower().strip()
        logger.info(f"[INTRO_OUTRO] 📍 Falling back to first location: '{first_location}'")
        
        # Get the actual key from display name or direct match
        matched_key = config.get_location_key_from_display_name(first_location)
        if not matched_key:
            # Try old matching logic
            for key in mapping.keys():
                if key in first_location or first_location in key:
                    matched_key = key
                    break
        
        if matched_key:
            location_meta = config.LOCATION_METADATA.get(matched_key, {})
            series = location_meta.get('series', '')
            display_name = location_meta.get('display_name', matched_key)
            display_type = location_meta.get('display_type', 'Unknown')

            logger.info(f"[INTRO_OUTRO] 🎯 Using first location: '{display_name}' (key: {matched_key})")
            logger.info(f"[INTRO_OUTRO]   - Display Type: {display_type}")
            logger.info(f"[INTRO_OUTRO]   - Series: {series}")

            # Mark as non-landmark for rest.pdf usage
            return {
                'key': matched_key,
                'series': series,
                'template_path': str(config.TEMPLATES_DIR / mapping[matched_key]),
                'metadata': location_meta,
                'is_non_landmark': True  # Flag to use rest.pdf
            }
    
    logger.info(f"[INTRO_OUTRO] ⚠️ No suitable location found for intro/outro")
    return None


def create_proposal_with_template(source_path: str, financial_data: dict, currency: str = None) -> Tuple[str, List[str], List[str]]:
    import tempfile

    pres = Presentation(source_path)
    insert_position = max(len(pres.slides) - 1, 0)
    slide_width = pres.slide_width
    slide_height = pres.slide_height

    blank_layout = pres.slide_layouts[6] if len(pres.slide_layouts) > 6 else pres.slide_layouts[0]
    financial_slide = pres.slides.add_slide(blank_layout)

    vat_amounts, total_amounts = create_financial_proposal_slide(financial_slide, financial_data, slide_width, slide_height, currency)

    if len(pres.slides) > 1 and insert_position < len(pres.slides) - 1:
        xml_slides = pres.slides._sldIdLst
        new_slide_element = xml_slides[-1]
        xml_slides.remove(new_slide_element)
        xml_slides.insert(insert_position, new_slide_element)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    pres.save(tmp.name)
    return tmp.name, vat_amounts, total_amounts


def create_combined_proposal_with_template(source_path: str, proposals_data: list, combined_net_rate: str, client_name: str, payment_terms: str = "100% upfront", currency: str = None) -> Tuple[str, str]:
    import tempfile

    pres = Presentation(source_path)
    insert_position = max(len(pres.slides) - 1, 0)
    slide_width = pres.slide_width
    slide_height = pres.slide_height

    layout = pres.slide_layouts[0]
    financial_slide = pres.slides.add_slide(layout)

    for shape in list(financial_slide.shapes):
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()

    total_combined = create_combined_financial_proposal_slide(
        financial_slide,
        proposals_data,
        combined_net_rate,
        slide_width,
        slide_height,
        client_name,
        payment_terms,
        currency,
    )

    xml_slides = pres.slides._sldIdLst
    slides_list = list(xml_slides)
    new_slide_element = slides_list[-1]
    xml_slides.remove(new_slide_element)
    xml_slides.insert(insert_position, new_slide_element)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    pres.save(tmp.name)
    return tmp.name, total_combined


async def process_combined_package(proposals_data: list, combined_net_rate: str, submitted_by: str, client_name: str, payment_terms: str = "100% upfront", currency: str = None) -> Dict[str, Any]:
    logger = config.logger
    logger.info(f"[COMBINED] Starting process_combined_package")
    logger.info(f"[COMBINED] Proposals: {proposals_data}")
    logger.info(f"[COMBINED] Combined rate: {combined_net_rate}")
    logger.info(f"[COMBINED] Client: {client_name}, Submitted by: {submitted_by}")
    logger.info(f"[COMBINED] Payment terms received: {payment_terms}")
    logger.info(f"[COMBINED] Currency: {currency or 'AED'}")
    
    validated_proposals = []
    for idx, proposal in enumerate(proposals_data):
        location = proposal.get("location", "").lower().strip()
        start_date = proposal.get("start_date", "1st December 2025")
        durations = proposal.get("durations", [])
        spots = int(proposal.get("spots", 1))
        
        logger.info(f"[COMBINED] Validating proposal {idx + 1}:")
        logger.info(f"[COMBINED]   Location: '{location}'")
        logger.info(f"[COMBINED]   Start date: {start_date}")
        logger.info(f"[COMBINED]   Durations: {durations}")
        logger.info(f"[COMBINED]   Spots: {spots}")

        # Get the mapping first (we'll need it later)
        mapping = config.get_location_mapping()
        
        # First try to get key from display name
        matched_key = config.get_location_key_from_display_name(location)
        
        # If that didn't work, try the old matching logic
        if not matched_key:
            logger.info(f"[COMBINED] Available mappings: {list(mapping.keys())}")
            
            for key in mapping.keys():
                if key in location or location in key:
                    matched_key = key
                    logger.info(f"[COMBINED] Matched '{location}' to '{key}'")
                    break
        else:
            logger.info(f"[COMBINED] Matched display name '{location}' to key '{matched_key}'")
                
        if not matched_key:
            logger.error(f"[COMBINED] No match found for location '{location}'")
            return {"success": False, "error": f"Unknown location '{location}' in proposal {idx + 1}"}
        if not durations:
            return {"success": False, "error": f"No duration specified for {matched_key}"}

        validated_proposal = {
            "location": matched_key,
            "start_date": start_date,
            "durations": durations,
            "spots": spots,
            "filename": mapping[matched_key],
        }

        # Add end_date if provided
        end_date = proposal.get("end_date")
        if end_date:
            validated_proposal["end_date"] = end_date

        # Add production fee if provided
        production_fee = proposal.get("production_fee")
        if production_fee:
            validated_proposal["production_fee"] = production_fee

        validated_proposals.append(validated_proposal)

    loop = asyncio.get_event_loop()
    pdf_files: List[str] = []

    # Check if we'll have intro/outro slides
    intro_outro_info = _get_location_info_for_intro_outro(validated_proposals)

    # Deduplicate deck slides - same location can appear multiple times (multiple durations)
    # but we only want to show the deck slides once
    seen_locations = set()
    unique_proposals_for_deck = []

    for proposal in validated_proposals:
        location_key = proposal["location"]
        if location_key not in seen_locations:
            seen_locations.add(location_key)
            unique_proposals_for_deck.append(proposal)
            logger.info(f"[COMBINED] Including deck slides for '{location_key}' (first occurrence)")
        else:
            logger.info(f"[COMBINED] Skipping duplicate deck slides for '{location_key}'")

    for idx, proposal in enumerate(unique_proposals_for_deck):
        src = config.TEMPLATES_DIR / proposal["filename"]
        if not src.exists():
            return {"success": False, "error": f"{proposal['filename']} not found"}

        if idx == len(unique_proposals_for_deck) - 1:
            # Investment sheet uses ALL validated_proposals (including duplicates for multiple durations)
            pptx_file, total_combined = await loop.run_in_executor(
                None,
                create_combined_proposal_with_template,
                str(src),
                validated_proposals,  # Use original list with all proposals
                combined_net_rate,
                client_name,
                payment_terms,
                currency,
            )
        else:
            pptx_file = str(src)
            total_combined = None

        # When we have intro/outro slides, remove both first and last from all PPTs
        if intro_outro_info:
            remove_first = True
            remove_last = True
        else:
            # Legacy behavior when no intro/outro template
            remove_first = False
            remove_last = False
            if idx == 0:
                remove_last = True
            elif idx < len(validated_proposals) - 1:
                remove_first = True
                remove_last = True
            else:
                remove_first = True

        pdf_file = await remove_slides_and_convert_to_pdf(pptx_file, remove_first, remove_last)
        pdf_files.append(pdf_file)

        if idx == len(validated_proposals) - 1:
            try:
                os.unlink(pptx_file)
            except:
                pass
    
    # For combined proposals, create intro and outro slides
    if intro_outro_info:
        series = intro_outro_info.get('series', '')
        location_key = intro_outro_info.get('key', '')
        display_name = intro_outro_info.get('metadata', {}).get('display_name', location_key)
        
        logger.info(f"[COMBINED] 🎬 Creating intro/outro slides")
        logger.info(f"[COMBINED] 📍 Selected location: '{display_name}' (key: {location_key})")
        logger.info(f"[COMBINED] 📂 Series: '{series}'")
        
        # Check for pre-made PDFs in intro_outro directory
        intro_outro_dir = config.TEMPLATES_DIR / "intro_outro"
        pdf_path = None
        
        # Map series to PDF filenames
        is_landmark = intro_outro_info.get('is_landmark', False)
        is_non_landmark = intro_outro_info.get('is_non_landmark', False)

        if is_landmark or series == 'The Landmark Series':
            # Use landmark_series.pdf for The Landmark Series
            pdf_path = intro_outro_dir / "landmark_series.pdf"
            logger.info(f"[COMBINED] 🏆 LANDMARK SERIES DETECTED! Looking for pre-made PDF...")
        elif is_non_landmark:
            # Use rest.pdf for non-Landmark Series locations
            pdf_path = intro_outro_dir / "rest.pdf"
            logger.info(f"[COMBINED] 🏢 NON-LANDMARK LOCATIONS DETECTED! Using rest.pdf...")
        elif 'Digital Icons' in series:
            pdf_path = intro_outro_dir / "digital_icons.pdf"
            logger.info(f"[COMBINED] 💎 DIGITAL ICONS SERIES DETECTED! Looking for pre-made PDF...")
        else:
            logger.info(f"[COMBINED] ❓ No pre-made PDF mapping for series '{series}'")
        
        if pdf_path and pdf_path.exists():
            logger.info(f"[COMBINED] ✅ PRE-MADE PDF FOUND! Using: {pdf_path}")
            # Extract first page for intro
            intro_pdf = _extract_pages_from_pdf(str(pdf_path), [0])
            # Extract last page for outro (assuming 2-page PDF)
            reader = PdfReader(str(pdf_path))
            last_page = len(reader.pages) - 1
            outro_pdf = _extract_pages_from_pdf(str(pdf_path), [last_page])
        else:
            # Fall back to PowerPoint extraction
            if pdf_path:
                logger.info(f"[COMBINED] ❌ PRE-MADE PDF NOT FOUND at: {pdf_path}")
            logger.info(f"[COMBINED] 🔄 FALLING BACK to PowerPoint extraction")
            template_path = intro_outro_info['template_path']
            logger.info(f"[COMBINED] 📄 Using PowerPoint template: {template_path}")
            
            # Create intro by keeping only the first slide
            intro_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            intro_pptx.close()
            shutil.copy2(template_path, intro_pptx.name)
            
            # Remove all slides except the first
            pres = Presentation(intro_pptx.name)
            xml_slides = pres.slides._sldIdLst
            slides_to_remove = list(xml_slides)[1:]  # All slides except first
            for slide_id in slides_to_remove:
                xml_slides.remove(slide_id)
            pres.save(intro_pptx.name)
            
            intro_pdf = await loop.run_in_executor(None, convert_pptx_to_pdf, intro_pptx.name)
            
            # Create outro by keeping only the last slide
            outro_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            outro_pptx.close()
            shutil.copy2(template_path, outro_pptx.name)
            
            # Remove all slides except the last
            pres = Presentation(outro_pptx.name)
            xml_slides = pres.slides._sldIdLst
            slides_to_remove = list(xml_slides)[:-1]  # All slides except last
            for slide_id in slides_to_remove:
                xml_slides.remove(slide_id)
            pres.save(outro_pptx.name)
            
            outro_pdf = await loop.run_in_executor(None, convert_pptx_to_pdf, outro_pptx.name)
            
            # Clean up temp files
            try:
                os.unlink(intro_pptx.name)
                os.unlink(outro_pptx.name)
            except Exception as e:
                logger.warning(f"Failed to clean up intro/outro files: {e}")
        
        # Insert intro at beginning and outro at end
        pdf_files.insert(0, intro_pdf)
        pdf_files.append(outro_pdf)

    merged_pdf = await loop.run_in_executor(None, merge_pdfs, pdf_files)
    for pdf_file in pdf_files:
        try:
            os.unlink(pdf_file)
        except Exception as e:
            logger.warning(f"Failed to clean up PDF file {pdf_file}: {e}")

    locations_str = ", ".join([p["location"].title() for p in validated_proposals])

    if total_combined is None:
        municipality_fee = 520
        total_upload_fees = sum(config.UPLOAD_FEES_MAPPING.get(p["location"].lower(), 3000) for p in validated_proposals)
        net_rate_numeric = float(combined_net_rate.replace("AED", "").replace(",", "").strip())
        subtotal = net_rate_numeric + total_upload_fees + municipality_fee
        vat = subtotal * 0.05
        total_combined = f"AED {subtotal + vat:,.0f}"

    db.log_proposal(
        submitted_by=submitted_by,
        client_name=client_name,
        package_type="combined",
        locations=locations_str,
        total_amount=total_combined,
    )

    timestamp_code = _generate_timestamp_code()
    client_prefix = client_name.replace(" ", "_") if client_name else "Client"

    return {
        "success": True,
        "is_combined": True,
        "pptx_path": None,
        "pdf_path": merged_pdf,
        "locations": locations_str,
        "pdf_filename": f"{client_prefix}_{timestamp_code}.pdf",
    }


async def process_proposals(
    proposals_data: list,
    package_type: str = "separate",
    combined_net_rate: str = None,
    submitted_by: str = "",
    client_name: str = "",
    payment_terms: str = "100% upfront",
    currency: str = None,
) -> Dict[str, Any]:
    """Process proposal generation with optional currency conversion.

    Args:
        proposals_data: List of proposal dicts
        package_type: "separate" or "combined"
        combined_net_rate: Net rate for combined packages
        submitted_by: User who submitted
        client_name: Client name
        payment_terms: Payment terms text
        currency: Target currency code (e.g., 'USD', 'EUR'). If None or 'AED', uses AED.

    Returns:
        Dict with success status and file paths
    """
    logger = config.logger
    logger.info(f"[PROCESS] Starting process_proposals")
    logger.info(f"[PROCESS] Package type: {package_type}")
    logger.info(f"[PROCESS] Proposals data: {proposals_data}")
    logger.info(f"[PROCESS] Combined rate: {combined_net_rate}")
    logger.info(f"[PROCESS] Submitted by: {submitted_by}")
    logger.info(f"[PROCESS] Client: {client_name}")
    logger.info(f"[PROCESS] Currency: {currency or 'AED'}")

    if not proposals_data:
        return {"success": False, "error": "No proposals provided"}

    is_single = len(proposals_data) == 1 and package_type != "combined"
    logger.info(f"[PROCESS] Is single: {is_single}")

    if package_type == "combined" and len(proposals_data) > 1:
        logger.info("[PROCESS] Routing to process_combined_package")
        return await process_combined_package(proposals_data, combined_net_rate, submitted_by, client_name, payment_terms, currency)

    individual_files = []
    pdf_files = []
    locations = []

    loop = asyncio.get_event_loop()
    
    # Check if we'll have intro/outro slides for multiple proposals
    intro_outro_info = None
    if len(proposals_data) > 1:
        intro_outro_info = _get_location_info_for_intro_outro(proposals_data)

    # Process all proposals in parallel for better performance
    async def process_single_proposal(idx: int, proposal: dict):
        location = proposal.get("location", "").lower().strip()
        start_date = proposal.get("start_date", "1st December 2025")
        durations = proposal.get("durations", [])
        net_rates = proposal.get("net_rates", [])
        spots = int(proposal.get("spots", 1))
        
        logger.info(f"[PROCESS] Processing proposal {idx + 1}:")
        logger.info(f"[PROCESS]   Location: '{location}'")
        logger.info(f"[PROCESS]   Start date: {start_date}")
        logger.info(f"[PROCESS]   Durations: {durations}")
        logger.info(f"[PROCESS]   Net rates: {net_rates}")
        logger.info(f"[PROCESS]   Spots: {spots}")

        # Get the mapping first (we'll need it later)
        mapping = config.get_location_mapping()
        
        # First try to get key from display name
        matched_key = config.get_location_key_from_display_name(location)
        
        # If that didn't work, try the old matching logic
        if not matched_key:
            logger.info(f"[PROCESS] Available location mappings: {list(mapping.keys())}")
            
            for key in mapping.keys():
                if key in location or location in key:
                    matched_key = key
                    logger.info(f"[PROCESS] Matched '{location}' to '{key}'")
                    break
        else:
            logger.info(f"[PROCESS] Matched display name '{location}' to key '{matched_key}'")
        
        if not matched_key:
            logger.error(f"[PROCESS] No match found for location '{location}'")
            return {"success": False, "error": f"Unknown location '{location}' in proposal {idx + 1}"}

        if len(durations) != len(net_rates):
            return {"success": False, "error": f"Mismatched durations and rates for {matched_key} - {len(durations)} durations but {len(net_rates)} rates"}
        if not durations:
            return {"success": False, "error": f"No duration specified for {matched_key}"}

        src = config.TEMPLATES_DIR / mapping[matched_key]
        if not src.exists():
            return {"success": False, "error": f"{mapping[matched_key]} not found"}

        financial_data = {
            "location": matched_key,
            "start_date": start_date,
            "durations": durations,
            "net_rates": net_rates,
            "spots": spots,
            "client_name": client_name,
        }

        # Add end_date if provided
        end_date = proposal.get("end_date")
        if end_date:
            financial_data["end_date"] = end_date

        # Add production fee if provided
        production_fee = proposal.get("production_fee")
        if production_fee:
            financial_data["production_fee"] = production_fee

        # Add payment_terms if provided
        payment_terms = proposal.get("payment_terms")
        if payment_terms:
            financial_data["payment_terms"] = payment_terms

        pptx_file, vat_amounts, total_amounts = await loop.run_in_executor(None, create_proposal_with_template, str(src), financial_data, currency)

        result = {
            "path": pptx_file,
            "location": matched_key.title(),
            "filename": f"{matched_key.title()}_Proposal.pptx",
            "totals": total_amounts,
            "matched_key": matched_key,
            "idx": idx
        }

        if is_single:
            pdf_file = await loop.run_in_executor(None, convert_pptx_to_pdf, pptx_file)
            result["pdf_path"] = pdf_file
            timestamp_code = _generate_timestamp_code()
            client_prefix = client_name.replace(" ", "_") if client_name else "Client"
            result["pdf_filename"] = f"{client_prefix}_{timestamp_code}.pdf"
        else:
            # When we have intro/outro slides, remove both first and last from all PPTs
            if intro_outro_info:
                remove_first = True
                remove_last = True
            else:
                # Legacy behavior when no intro/outro template
                remove_first = False
                remove_last = False
                if idx == 0:
                    remove_last = True
                elif idx < len(proposals_data) - 1:
                    remove_first = True
                    remove_last = True
                else:
                    remove_first = True
            pdf_file = await remove_slides_and_convert_to_pdf(pptx_file, remove_first, remove_last)
            result["pdf_file"] = pdf_file
            
        return {"success": True, "result": result}

    # Process all proposals in parallel
    tasks = [process_single_proposal(idx, proposal) for idx, proposal in enumerate(proposals_data)]
    results = await asyncio.gather(*tasks, return_exceptions=True)
    
    # Check for errors and organize results
    for idx, result in enumerate(results):
        if isinstance(result, Exception):
            return {"success": False, "error": f"Error processing proposal {idx + 1}: {str(result)}"}
        if isinstance(result, dict) and not result.get("success"):
            return result  # Return the error
    
    # Sort results by original index to maintain order
    sorted_results = sorted(
        [r for r in results if r.get("success")],
        key=lambda x: x["result"]["idx"]
    )
    
    # Extract successful results in order
    for result in sorted_results:
        proposal_result = result["result"]
        individual_files.append({
            "path": proposal_result["path"],
            "location": proposal_result["location"],
            "filename": proposal_result["filename"],
            "totals": proposal_result["totals"],
        })
        if "pdf_path" in proposal_result:
            individual_files[-1]["pdf_path"] = proposal_result["pdf_path"]
            individual_files[-1]["pdf_filename"] = proposal_result["pdf_filename"]
        if "pdf_file" in proposal_result:
            pdf_files.append(proposal_result["pdf_file"])
        locations.append(proposal_result["location"])
    
    # For multiple proposals, create intro and outro slides
    if len(pdf_files) > 1 and intro_outro_info:
            series = intro_outro_info.get('series', '')
            location_key = intro_outro_info.get('key', '')
            display_name = intro_outro_info.get('metadata', {}).get('display_name', location_key)
            
            logger.info(f"[PROCESS] 🎬 Creating intro/outro slides")
            logger.info(f"[PROCESS] 📍 Selected location: '{display_name}' (key: {location_key})")
            logger.info(f"[PROCESS] 📂 Series: '{series}'")
            
            # Check for pre-made PDFs in intro_outro directory
            intro_outro_dir = config.TEMPLATES_DIR / "intro_outro"
            pdf_path = None
            
            # Map series to PDF filenames
            if 'Landmark' in series:
                pdf_path = intro_outro_dir / "landmark_series.pdf"
                logger.info(f"[PROCESS] 🏆 LANDMARK SERIES DETECTED! Looking for pre-made PDF...")
            elif 'Digital Icons' in series:
                pdf_path = intro_outro_dir / "digital_icons.pdf"
                logger.info(f"[PROCESS] 💎 DIGITAL ICONS SERIES DETECTED! Looking for pre-made PDF...")
            else:
                logger.info(f"[PROCESS] ❓ No pre-made PDF mapping for series '{series}'")
            
            if pdf_path and pdf_path.exists():
                logger.info(f"[PROCESS] ✅ PRE-MADE PDF FOUND! Using: {pdf_path}")
                # Extract first page for intro
                intro_pdf = _extract_pages_from_pdf(str(pdf_path), [0])
                # Extract last page for outro (assuming 2-page PDF)
                reader = PdfReader(str(pdf_path))
                last_page = len(reader.pages) - 1
                outro_pdf = _extract_pages_from_pdf(str(pdf_path), [last_page])
            else:
                # Fall back to PowerPoint extraction
                if pdf_path:
                    logger.info(f"[PROCESS] ❌ PRE-MADE PDF NOT FOUND at: {pdf_path}")
                logger.info(f"[PROCESS] 🔄 FALLING BACK to PowerPoint extraction")
                template_path = intro_outro_info['template_path']
                logger.info(f"[PROCESS] 📄 Using PowerPoint template: {template_path}")
                
                # Create intro by keeping only the first slide
                intro_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
                intro_pptx.close()
                shutil.copy2(template_path, intro_pptx.name)
                
                # Remove all slides except the first
                pres = Presentation(intro_pptx.name)
                xml_slides = pres.slides._sldIdLst
                slides_to_remove = list(xml_slides)[1:]  # All slides except first
                for slide_id in slides_to_remove:
                    xml_slides.remove(slide_id)
                pres.save(intro_pptx.name)
                
                intro_pdf = await loop.run_in_executor(None, convert_pptx_to_pdf, intro_pptx.name)
                
                # Create outro by keeping only the last slide
                outro_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
                outro_pptx.close()
                shutil.copy2(template_path, outro_pptx.name)
                
                # Remove all slides except the last
                pres = Presentation(outro_pptx.name)
                xml_slides = pres.slides._sldIdLst
                slides_to_remove = list(xml_slides)[:-1]  # All slides except last
                for slide_id in slides_to_remove:
                    xml_slides.remove(slide_id)
                pres.save(outro_pptx.name)
                
                outro_pdf = await loop.run_in_executor(None, convert_pptx_to_pdf, outro_pptx.name)
                
                # Clean up temp files
                try:
                    os.unlink(intro_pptx.name)
                    os.unlink(outro_pptx.name)
                except:
                    pass
            
            # Insert intro at beginning and outro at end
            pdf_files.insert(0, intro_pdf)
            pdf_files.append(outro_pdf)

    if is_single:
        totals = individual_files[0].get("totals", [])
        total_str = totals[0] if totals else "AED 0"
        db.log_proposal(
            submitted_by=submitted_by,
            client_name=client_name,
            package_type="single",
            locations=individual_files[0]["location"],
            total_amount=total_str,
        )
        return {
            "success": True,
            "is_single": True,
            "pptx_path": individual_files[0]["path"],
            "pdf_path": individual_files[0]["pdf_path"],
            "location": individual_files[0]["location"],
            "pptx_filename": individual_files[0]["filename"],
            "pdf_filename": individual_files[0]["pdf_filename"],
        }

    merged_pdf = await loop.run_in_executor(None, merge_pdfs, pdf_files)
    for pdf_file in pdf_files:
        try:
            os.unlink(pdf_file)
        except Exception as e:
            logger.warning(f"Failed to clean up PDF file {pdf_file}: {e}")

    first_totals = [files.get("totals", ["AED 0"])[0] for files in individual_files]
    summary_total = ", ".join(first_totals)
    db.log_proposal(
        submitted_by=submitted_by,
        client_name=client_name,
        package_type="separate",
        locations=", ".join(locations),
        total_amount=summary_total,
    )

    timestamp_code = _generate_timestamp_code()
    client_prefix = client_name.replace(" ", "_") if client_name else "Client"

    return {
        "success": True,
        "is_single": False,
        "individual_files": individual_files,
        "merged_pdf_path": merged_pdf,
        "locations": ", ".join(locations),
        "merged_pdf_filename": f"{client_prefix}_{timestamp_code}.pdf",
    } 