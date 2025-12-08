"""
File Utilities - Download, validate, and convert files from messaging channels.

Uses the unified channel abstraction layer for file downloads.
"""

from pathlib import Path
from typing import Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

import config
from utils.memory import cleanup_memory

logger = config.logger

def _validate_pdf_file(file_path: Path) -> bool:
    """Validate that uploaded file is actually a PDF."""
    try:
        # Check if file exists first
        if not file_path.exists():
            config.logger.error(f"[VALIDATION] File does not exist: {file_path}")
            return False

        file_size = file_path.stat().st_size
        config.logger.info(f"[VALIDATION] Validating PDF file: {file_path} (size: {file_size} bytes)")

        # Quick sanity checks before invoking PyMuPDF
        if file_size <= 0:
            config.logger.warning(f"[VALIDATION] PDF file is empty: {file_path}")
            return False

        # Try to open with PyMuPDF
        import fitz
        doc = fitz.open(file_path)
        page_count = len(doc)
        doc.close()

        if page_count <= 0:
            config.logger.warning(f"[VALIDATION] PDF has no pages: {file_path}")
            return False

        config.logger.info(f"[VALIDATION] ✓ Valid PDF with {page_count} pages")
        return True
    except Exception as e:
        config.logger.error(f"[VALIDATION] Invalid PDF file: {e}")
        return False


async def _convert_pdf_to_pptx(pdf_path: Path) -> Optional[Path]:
    """Convert PDF to PowerPoint with maximum quality (300 DPI).

    Args:
        pdf_path: Path to input PDF file

    Returns:
        Path to converted PPTX file, or None if conversion failed
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import tempfile
        import os

        config.logger.info(f"[PDF_CONVERT] Starting PDF to PPTX conversion: {pdf_path}")

        # Open PDF
        doc = fitz.open(pdf_path)
        page_count = len(doc)
        config.logger.info(f"[PDF_CONVERT] PDF has {page_count} pages")

        # Create temporary directory for images
        temp_dir = tempfile.mkdtemp(prefix="pdf_convert_")

        # Convert pages to high-resolution images (4x zoom = ~300 DPI)
        zoom = 4.0
        matrix = fitz.Matrix(zoom, zoom)
        image_paths = []

        for page_num in range(page_count):
            page = doc[page_num]
            pix = page.get_pixmap(matrix=matrix, alpha=False)

            # Save as PNG
            img_path = os.path.join(temp_dir, f"page_{page_num + 1}.png")
            pix.save(img_path)
            image_paths.append(img_path)
            config.logger.info(f"[PDF_CONVERT] Extracted page {page_num + 1}/{page_count} ({pix.width}x{pix.height}px)")

        doc.close()

        # Create PowerPoint presentation
        config.logger.info("[PDF_CONVERT] Creating PowerPoint presentation...")
        prs = Presentation()

        # Set slide size to 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Add pages as slides
        for i, img_path in enumerate(image_paths, 1):
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Get image dimensions
            img = Image.open(img_path)
            img_width, img_height = img.size
            img_aspect = img_width / img_height
            slide_aspect = float(prs.slide_width) / float(prs.slide_height)

            # Calculate dimensions to fill slide while maintaining aspect ratio
            if img_aspect > slide_aspect:
                # Image is wider - fit to width
                width = prs.slide_width
                height = int(float(prs.slide_width) / img_aspect)
                left = 0
                top = int((prs.slide_height - height) / 2)
            else:
                # Image is taller - fit to height
                height = prs.slide_height
                width = int(float(prs.slide_height) * img_aspect)
                left = int((prs.slide_width - width) / 2)
                top = 0

            # Add image to slide
            slide.shapes.add_picture(img_path, left, top, width, height)
            config.logger.info(f"[PDF_CONVERT] Added slide {i}/{page_count}")

        # Save PPTX to temporary file
        pptx_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(pptx_temp.name)
        pptx_temp.close()

        # Cleanup temporary images
        for img_path in image_paths:
            try:
                os.remove(img_path)
            except OSError as cleanup_err:
                logger.debug(f"[PDF_CONVERT] Failed to cleanup temp image {img_path}: {cleanup_err}")
        try:
            os.rmdir(temp_dir)
        except OSError as cleanup_err:
            logger.debug(f"[PDF_CONVERT] Failed to cleanup temp directory {temp_dir}: {cleanup_err}")

        file_size_mb = os.path.getsize(pptx_temp.name) / (1024 * 1024)
        config.logger.info(f"[PDF_CONVERT] ✓ Conversion complete: {pptx_temp.name} ({file_size_mb:.1f} MB, {page_count} slides, 300 DPI)")

        return Path(pptx_temp.name)

    except Exception as e:
        config.logger.error(f"[PDF_CONVERT] Conversion failed: {e}", exc_info=True)
        return None


def _validate_powerpoint_file(file_path: Path) -> bool:
    """Validate that uploaded file is actually a PowerPoint presentation."""
    try:
        # Check if file exists first
        if not file_path.exists():
            config.logger.error(f"[VALIDATION] File does not exist: {file_path}")
            return False

        file_size = file_path.stat().st_size
        config.logger.info(f"[VALIDATION] Validating PowerPoint file: {file_path} (size: {file_size} bytes)")

        # Quick sanity checks before invoking python-pptx
        if file_size <= 0:
            config.logger.warning(f"[VALIDATION] PowerPoint file is empty: {file_path}")
            return False

        # PPTX files are ZIP packages starting with 'PK\x03\x04'
        try:
            with open(file_path, 'rb') as fp:
                magic = fp.read(4)
            if magic != b'PK\x03\x04':
                config.logger.warning(
                    f"[VALIDATION] File does not look like a PPTX (ZIP signature missing). "
                    f"Likely an HTML error/permission issue from Slack. Path: {file_path}"
                )
                return False
        except Exception as e:
            config.logger.warning(f"[VALIDATION] Failed to read file header for {file_path}: {e}")
            return False

        # Try to open as PowerPoint - this will fail if not a valid PPTX
        pres = Presentation(str(file_path))
        # Basic validation: must have at least 1 slide
        slide_count = len(pres.slides)

        # CRITICAL: Delete presentation object to free memory
        # python-pptx loads entire PPT into RAM (50-100MB+)
        del pres
        cleanup_memory(context="pptx_validation", aggressive=False, log_stats=False)

        if slide_count < 1:
            config.logger.warning(f"[VALIDATION] PowerPoint file has no slides: {file_path}")
            return False

        config.logger.info(f"[VALIDATION] PowerPoint validation successful: {slide_count} slides")
        return True
    except Exception as e:
        config.logger.warning(f"[VALIDATION] PowerPoint validation failed: {e}")
        # Cleanup on error path too
        try:
            del pres
        except NameError:
            pass  # pres was never assigned
        cleanup_memory(context="pptx_validation_error", aggressive=False, log_stats=False)
        return False


async def download_file(file_info: Dict[str, Any]) -> Path:
    """
    Download a file from the messaging channel.

    Uses the unified channel abstraction layer for file downloads.
    Platform-agnostic: works with any channel adapter (Slack, Web, etc.)
    """
    channel = config.get_channel_adapter()
    if not channel:
        raise RuntimeError("No channel adapter available")

    config.logger.info(f"[DOWNLOAD] Downloading file: {file_info.get('name', 'unknown')}")

    # Use channel adapter to download file
    file_path = await channel.download_file(file_info)

    # Verify file was written
    if file_path.exists():
        config.logger.info(f"[DOWNLOAD] File successfully written: {file_path} (size: {file_path.stat().st_size} bytes)")
    else:
        config.logger.error(f"[DOWNLOAD] File not found after write: {file_path}")

    return file_path
